#!/usr/bin/env python3
"""Generate IT246 Unit 2 PowerPoint deck (with native-shape diagrams).
Run: python3 build_unit2_pptx.py  ->  IT246_Unit2.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

NAVY=RGBColor(0x0C,0x2B,0x4A); BLUE=RGBColor(0x18,0x5F,0xA5); TEAL=RGBColor(0x0F,0x6E,0x56)
AMBER=RGBColor(0x85,0x4F,0x0B); CORAL=RGBColor(0x99,0x3C,0x1D); LIGHT=RGBColor(0xF2,0xF5,0xF8)
PALET=RGBColor(0xE1,0xF5,0xEE); WHITE=RGBColor(0xFF,0xFF,0xFF); DARK=RGBColor(0x1A,0x1A,0x1A)
GREY=RGBColor(0x55,0x60,0x6B); LGREY=RGBColor(0xC7,0xCF,0xD6)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

def _bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def _box(s,l,t,w,h): return s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)).text_frame
def _bar(s,c,l,t,w,h):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=c; sp.line.fill.background(); sp.shadow.inherit=False; return sp
def _notes(s,t): s.notes_slide.notes_text_frame.text=t

def _settext(sp,text,font=WHITE,size=14,bold=True,align=PP_ALIGN.CENTER):
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Inches(0.06); tf.margin_right=Inches(0.06); tf.margin_top=Inches(0.03); tf.margin_bottom=Inches(0.03)
    lines=text.split("\n"); p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=lines[0]; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=font
    for extra in lines[1:]:
        p2=tf.add_paragraph(); p2.alignment=align; r=p2.add_run(); r.text=extra
        r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=font

def shp(s,kind,l,t,w,h,text="",fill=BLUE,font=WHITE,size=14,bold=True,line=None,align=PP_ALIGN.CENTER):
    sp=s.shapes.add_shape(kind,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1.25)
    sp.shadow.inherit=False
    if text: _settext(sp,text,font,size,bold,align)
    return sp

def conn(s,x1,y1,x2,y2,color=GREY,w=1.75):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb=color; c.line.width=Pt(w); c.shadow.inherit=False; return c

def add_title(title,subtitle,footer):
    s=prs.slides.add_slide(BLANK); _bg(s,NAVY); _bar(s,AMBER,0.9,3.05,1.6,0.10)
    tf=_box(s,0.9,2.0,11.5,1.0); tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(38); r.font.bold=True; r.font.color.rgb=WHITE
    tf2=_box(s,0.9,3.25,11.5,1.2); tf2.word_wrap=True
    r=tf2.paragraphs[0].add_run(); r.text=subtitle; r.font.size=Pt(20); r.font.color.rgb=RGBColor(0xCF,0xDD,0xEA)
    tf3=_box(s,0.9,6.5,11.5,0.6); r=tf3.paragraphs[0].add_run(); r.text=footer
    r.font.size=Pt(13); r.font.color.rgb=RGBColor(0x9A,0xB2,0xC6); return s

def add_divider(kicker,title,notes=""):
    s=prs.slides.add_slide(BLANK); _bg(s,BLUE)
    tfk=_box(s,0.9,2.4,11.5,0.6); r=tfk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=RGBColor(0xBF,0xDC,0xF2)
    tft=_box(s,0.9,3.0,11.5,1.6); tft.word_wrap=True; r=tft.paragraphs[0].add_run(); r.text=title
    r.font.size=Pt(33); r.font.bold=True; r.font.color.rgb=WHITE
    if notes: _notes(s,notes)
    return s

def _header(s,kicker,header,kc=BLUE):
    _bar(s,kc,0,0,SW.inches,0.18)
    tfk=_box(s,0.7,0.32,12.0,0.45); r=tfk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=kc
    tfh=_box(s,0.7,0.7,12.0,1.0); tfh.word_wrap=True; r=tfh.paragraphs[0].add_run(); r.text=header
    r.font.size=Pt(27); r.font.bold=True; r.font.color.rgb=NAVY

def add_content(header,kicker,blocks,notes="",kc=BLUE):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,header,kc)
    tf=_box(s,0.7,1.8,11.9,5.3); tf.word_wrap=True; first=True
    for text,level in blocks:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=text
        if level==0: r.font.size=Pt(18); r.font.color.rgb=DARK; p.space_after=Pt(7); p.level=0
        elif level==1: r.font.size=Pt(15.5); r.font.color.rgb=GREY; p.space_after=Pt(3); p.level=1
        else: r.font.size=Pt(15.5); r.font.color.rgb=TEAL; r.font.italic=True; p.space_after=Pt(6); p.level=1
    if notes: _notes(s,notes)
    return s

def add_quiz(title,items,notes=""):
    s=prs.slides.add_slide(BLANK); _bg(s,LIGHT); _bar(s,AMBER,0,0,SW.inches,0.18)
    tfk=_box(s,0.7,0.32,12,0.45); r=tfk.paragraphs[0].add_run(); r.text="CHECK FOR UNDERSTANDING"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=AMBER
    tfh=_box(s,0.7,0.7,12,0.9); r=tfh.paragraphs[0].add_run(); r.text=title
    r.font.size=Pt(26); r.font.bold=True; r.font.color.rgb=NAVY
    tf=_box(s,0.7,1.8,11.9,5.2); tf.word_wrap=True; first=True
    for text,kind in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=text
        if kind=="q": r.font.size=Pt(18); r.font.bold=True; r.font.color.rgb=NAVY; p.space_before=Pt(8)
        elif kind=="a": r.font.size=Pt(16); r.font.color.rgb=TEAL; r.font.bold=True; p.level=1
        else: r.font.size=Pt(16); r.font.color.rgb=DARK; p.level=1
    if notes: _notes(s,notes)
    return s

def add_activity(title,blocks,notes=""):
    s=prs.slides.add_slide(BLANK); _bg(s,PALET); _bar(s,TEAL,0,0,SW.inches,0.18)
    tfk=_box(s,0.7,0.32,12,0.45); r=tfk.paragraphs[0].add_run(); r.text="🛠  IN-CLASS ACTIVITY"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=TEAL
    tfh=_box(s,0.7,0.7,12,0.9); r=tfh.paragraphs[0].add_run(); r.text=title
    r.font.size=Pt(26); r.font.bold=True; r.font.color.rgb=NAVY
    tf=_box(s,0.7,1.85,11.9,5.2); tf.word_wrap=True; first=True
    for text,level in blocks:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=text
        if level==0: r.font.size=Pt(18); r.font.color.rgb=DARK; p.space_after=Pt(8)
        else: r.font.size=Pt(15.5); r.font.color.rgb=GREY; p.space_after=Pt(4); p.level=1
    if notes: _notes(s,notes)
    return s

def diagram_slide(kicker,header,kc=BLUE):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,header,kc); return s

# =============== DIAGRAMS ===============
def diag_relationship_web():
    s=diagram_slide("S6 · Concept 1  ·  Diagram","The IT Worker's Six Relationships")
    cx,cy,cw,ch=5.45,3.15,2.3,1.0; ctr=(cx+cw/2,cy+ch/2)
    boxes=[(5.45,1.85,"EMPLOYER\nloyalty · IP",BLUE),
           (1.2,2.7,"SUPPLIER\nno bribery",TEAL),
           (9.6,2.7,"CLIENT\nhonest advice",TEAL),
           (1.2,4.85,"PEERS\nrespect",AMBER),
           (9.6,4.85,"SOCIETY\nsafe systems",AMBER),
           (5.45,5.65,"IT USERS\nsupport, don't exploit",CORAL)]
    bw,bh=2.3,0.85
    for x,y,_,_c in boxes:  # lines behind
        conn(s,ctr[0],ctr[1],x+bw/2,y+bh/2,color=LGREY,w=2)
    for x,y,txt,col in boxes:
        shp(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,y,bw,bh,txt,fill=col,size=12.5)
    shp(s,MSO_SHAPE.OVAL,cx,cy,cw,ch,"IT\nWORKER",fill=NAVY,size=16)
    _notes(s,"Access = power = responsibility. Each spoke carries a duty. The 'free iPhone' lives on the SUPPLIER spoke; moonlighting for a competitor lives on EMPLOYER.")

def diag_cert_vs_license():
    s=diagram_slide("S7 · Concepts 3–4  ·  Comparison","Certification vs Licensing")
    rows=[("","Certification","Licensing"),
          ("Granted by","Vendor / industry body","Government"),
          ("Legal force","None (voluntary)","Yes — legally required"),
          ("Proves","A specific skill","Right to practice"),
          ("Example","CCNA, CISSP, AWS","NMC (doctors), NEC (engineers)")]
    L,T,W,H=1.6,2.1,10.1,4.0
    tbl=s.shapes.add_table(len(rows),3,Inches(L),Inches(T),Inches(W),Inches(H)).table
    tbl.columns[0].width=Inches(2.7); tbl.columns[1].width=Inches(3.7); tbl.columns[2].width=Inches(3.7)
    for ri,row in enumerate(rows):
        for ci,val in enumerate(row):
            cell=tbl.cell(ri,ci); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.12); cell.margin_top=Inches(0.03); cell.margin_bottom=Inches(0.03)
            p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
            r=p.add_run(); r.text=val
            if ri==0:
                cell.fill.solid(); cell.fill.fore_color.rgb=(WHITE if ci==0 else (BLUE if ci==1 else TEAL))
                r.font.color.rgb=(NAVY if ci==0 else WHITE); r.font.bold=True; r.font.size=Pt(16)
            elif ci==0:
                cell.fill.solid(); cell.fill.fore_color.rgb=NAVY; r.font.color.rgb=WHITE; r.font.bold=True; r.font.size=Pt(14)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb=(RGBColor(0xE6,0xF1,0xFB) if ci==1 else PALET)
                r.font.color.rgb=DARK; r.font.size=Pt(13.5)
    _notes(s,"Key line: certification proves SKILL (voluntary); licensing grants the RIGHT to practice (legal). In Nepal doctors/engineers are licensed; IT workers usually aren't — that's the S7 debate.")

def diag_aup_card():
    s=diagram_slide("S8 · Concept 2  ·  Diagram","Acceptable Use Policy — Allowed vs Prohibited")
    shp(s,MSO_SHAPE.ROUNDED_RECTANGLE,0.9,2.15,5.6,0.7,"✔  ALLOWED",fill=TEAL,size=18)
    shp(s,MSO_SHAPE.ROUNDED_RECTANGLE,6.85,2.15,5.6,0.7,"✘  PROHIBITED",fill=CORAL,size=18)
    allowed=["Work tasks & approved software","Reasonable, brief personal use","Reporting security issues","Using your own credentials only"]
    prohib=["Pirated / cracked software","Torrenting, crypto-mining, gaming","Sharing passwords / USB drives","Personal side-business on company email"]
    tfa=_box(s,1.05,3.05,5.3,3.6); tfa.word_wrap=True
    for i,t in enumerate(allowed):
        p=tfa.paragraphs[0] if i==0 else tfa.add_paragraph(); r=p.add_run(); r.text="•  "+t
        r.font.size=Pt(15); r.font.color.rgb=DARK; p.space_after=Pt(10)
    tfp=_box(s,7.0,3.05,5.3,3.6); tfp.word_wrap=True
    for i,t in enumerate(prohib):
        p=tfp.paragraphs[0] if i==0 else tfp.add_paragraph(); r=p.add_run(); r.text="•  "+t
        r.font.size=Pt(15); r.font.color.rgb=DARK; p.space_after=Pt(10)
    _notes(s,"Every job hands you an AUP on day one. Breaking it (piracy, misuse) is a common reason people are fired or sued. Use the shared-Wi-Fi hook to populate both columns.")

def diag_data_trail():
    s=diagram_slide("S9 · Concept 1  ·  Diagram","Your Data Trail — where consent disappears")
    steps=[("YOU\nyour personal data",GREY),("Loan form · SIM ·\nQR · KYC",BLUE),
           ("Company\ncollects & stores",TEAL),("Third parties\nspam · profiling · resale",CORAL)]
    w,h,y=2.85,1.6,3.0; gap=0.12; total=4*w+3*gap; x0=(13.333-total)/2
    for i,(txt,col) in enumerate(steps):
        shp(s,MSO_SHAPE.PENTAGON,x0+i*(w+gap),y,w,h,txt,fill=col,size=13.5)
    cap=_box(s,0.7,5.0,11.9,1.0); cap.word_wrap=True
    r=cap.paragraphs[0].add_run()
    r.text="At each hop you lose control. Data minimization + consent are the brakes. (Nepal: Individual Privacy Act 2075.)"
    r.font.size=Pt(14); r.font.italic=True; r.font.color.rgb=GREY
    _notes(s,"The spam-call hook IS this trail. 'I have nothing to hide' misses the point — privacy is about control at each hop, not hiding.")

def diag_defamation_spectrum():
    s=diagram_slide("S10 · Concept 4  ·  Diagram","Opinion → Criticism → Harassment → Defamation")
    cells=[("OPINION","'I didn't enjoy it'\nlow risk",TEAL),
           ("CRITICISM","'Service is slow\n& overpriced'",BLUE),
           ("HARASSMENT","repeated targeting\nof a person",AMBER),
           ("DEFAMATION","false claim of fact\nthat harms — ETA risk",CORAL)]
    w,h,y=2.9,1.7,2.9; gap=0.12; total=4*w+3*gap; x0=(13.333-total)/2
    for i,(head,sub,col) in enumerate(cells):
        sp=shp(s,MSO_SHAPE.ROUNDED_RECTANGLE,x0+i*(w+gap),y,w,h,"",fill=col)
        tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=head
        r.font.size=Pt(17); r.font.bold=True; r.font.color.rgb=WHITE
        p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER; r=p2.add_run(); r.text=sub
        r.font.size=Pt(12.5); r.font.color.rgb=WHITE
        if i<3:
            shp(s,MSO_SHAPE.RIGHT_ARROW,x0+i*(w+gap)+w-0.04,y+h/2-0.14,gap+0.02,0.28,fill=LGREY)
    cap=_box(s,0.7,5.0,11.9,0.6); r=cap.paragraphs[0].add_run()
    r.text="Legal & ethical risk rises left → right. Opinion/criticism are protected; defamation is a false statement of FACT."
    r.font.size=Pt(14); r.font.italic=True; r.font.color.rgb=GREY
    _notes(s,"The 'shop X overcharges' vs 'shop X is a fraud' case maps onto criticism vs defamation. In Nepal defamation can trigger arrest under the ETA 2063 (Unit 9).")

# =============== BUILD ===============
add_title("Unit 2 — Ethics for IT Workers and IT Users",
          "IT 246: IT Ethics and Cybersecurity  ·  BIM 6th Semester",
          "5 lecture hours (S6–S10)  ·  50 min each  ·  Nepal / South Asia context")
add_content("Unit 2 — Learning Outcomes","Overview",
    [("By the end of this unit, you will be able to:",0),
     ("Identify the 6 IT-worker relationships and the ethical duty in each",1),
     ("Explain codes, certification, licensing, and professional organizations",1),
     ("Describe Acceptable Use Policies and user ethics issues (piracy, misuse)",1),
     ("Explain privacy, surveillance, anonymity, and identity issues",1),
     ("Reason about social-networking issues: bullying, fake news, defamation",1)],
    notes="Five sessions: workers (S6–S7), users (S8), privacy (S9), social media (S10).")

# S6 (50 min: 5 hook + 35 content + 5 CFU + 3 app + 2 summary)
add_divider("Session 6 · Lecture hour 1","Managing the IT Worker Relationship",
            "HOOK [~5 min]: 2 scenarios — a vendor offers the IT manager a 'free' iPhone the week before a Rs 2-crore tender; a dev moonlights for a client's competitor. Ask 'has anyone broken a LAW?' (no) then 'does it feel WRONG?' (yes). That gap is the lesson. HOLD the gift/bribe question — resolve in Concept 4. Agenda on board.")
add_content("Access = Power = Responsibility","S6 · Concept 1  [THEORY]  ·  ~8 min",
    [("IT worker = designs/builds/maintains IT systems AND the data inside them",0),
     ("What makes them different: PRIVILEGED ACCESS — admin/root, personal data, source code",0),
     ("Logic chain: access is power, and power is responsibility — 'can' ≠ 'should'",2),
     ("Nepal: a bank IT staffer COULD view any balance — only trust & ethics stop them",1),
     ("Misconception: 'IT is back-office, ethics is for managers.' Access IS power.",2)],
    notes="Write the definition on the board. Land the logic chain slowly — it's the session's spine. Global: a sysadmin reading executives' emails with root — the system LET them, only ethics should have stopped it. Transition: 'we have six arrows — follow each and ask what you owe the person at the other end.'")
diag_relationship_web()
add_content("Employer & Client","S6 · Concepts 2–3  [THEORY] [EXAMPLE]  ·  ~9 + ~8 min",
    [("EMPLOYER — duties of loyalty (protect assets/IP) and honesty (no resume fraud)",0),
     ("Strain points: conflicts of interest, resume fraud, license compliance, whistle-blowing tension",1),
     ("Nepal: copying client source code for a weekend gig = IP theft, betrays employer AND client",1),
     ("CLIENT — honest advice, deliver what's promised, disclose risks",0),
     ("Information asymmetry: the client can't judge technical quality — they rely on your honesty",2),
     ("Misconception: 'if the contract doesn't ban it, it's fine.' Duty goes beyond the contract.",2)],
    notes="Employer (~9 min): explain loyalty + honesty WITH why; name the 4 strain points (conflict of interest, resume fraud, license compliance, whistle-blowing). Mini case: boss says install pirated MS/AutoCAD — 'what do you do?' AND 'who is liable?' Tie to Unit 1 blind-obedience trap. Client (~8 min): the e-governance over-promise to a municipality — they lack expertise to catch under-delivery until it fails.")
add_content("Suppliers · Peers · Users · Society","S6 · Concept 4  [THEORY]  ·  ~5 min",
    [("Suppliers: no bribery/kickbacks — the 'free iPhone' IS a bribe (resolve the hook here)",0),
     ("Peers: respect, no resume inflation, no poaching trade secrets",0),
     ("IT users: support & train them; don't exploit their lack of skill",0),
     ("Society: build safe systems; don't enable harm (e.g. a system that leaks citizen data)",0),
     ("Analogy: an IT worker is a pharmacist — trusted with powerful 'medicine' (access/data)",2)],
    notes="Resolve the hook now. Gift-vs-bribe test: would it survive going public? timed around a decision? could it sway judgment? The iPhone fails all three.")
add_activity("S6 — 'Map your future job'  ·  ~5 min",
    [("Think–Pair–Share (2 min pairs, 3 min sharing)",0),
     ("Each pair picks an IT role (developer, sysadmin, support, data analyst)",1),
     ("Name one realistic ethical risk on TWO of the six relationship arrows for that role",1),
     ("Share aloud; place each on the relationship web",1)],
    notes="Makes the six-relationship web concrete and personal. Good seeds: support tech reading a user's private files 'to fix it'; analyst selling insights to a supplier.")
add_quiz("S6 — Quick Check  ·  ~5 min",
    [("Q1. A supplier's phone gift right before a tender is…","q"),
     ("✅ b) A conflict of interest / bribery","a"),
     ("Q2. Honest technical advice is owed most directly to the…","q"),
     ("✅ c) Client","a"),
     ("Discussion: which relationship is hardest for a junior IT worker in Nepal to manage honestly?","o")],
    notes="APPLICATION [~3 min]: your first job hands you access + vendor contacts on day one; how you handle gifts/access/honesty in week one sets your reputation for years. SUMMARY [~2 min]: 6 relationships; privileged access creates duties; loyalty/honesty/no-conflict run through all. Next: how the profession encourages good behavior.")

# S7
add_divider("Session 7 · Lecture hour 2","Encouraging Professionalism of IT Workers",
            "HOOK [~5 min]: a doctor swears an oath and can lose a license; an engineer signs off on a bridge. But anyone can call themselves an 'IT professional' tomorrow — there's usually no license to lose. Is IT really a profession? Should it be regulated like the others?")
add_content("Is IT a Profession?","S7 · Concept 1  [THEORY]  ·  ~7 min",
    [("Profession = specialized education + body of knowledge + code of ethics + self-regulation",0),
     ("IT partly qualifies (knowledge + codes) but lacks UNIVERSAL licensing — debated status",2),
     ("Nepal: doctors (NMC) & engineers (NEC) are licensed; IT workers generally are not",1),
     ("Misconception: 'anyone who fixes computers is a professional.' Standards & accountability > skill.",2)],
    notes="Set up the whole session's debate: does IT meet the bar of a 'profession'? Keep it open — resolved through codes/cert/licensing.")
add_content("Professional Codes of Ethics","S7 · Concept 2  [THEORY] [EXAMPLE]  ·  ~8 min",
    [("Code of ethics = published principles members commit to (ACM Code, IEEE-CS)",0),
     ("Codes set expectations, build public trust, guide decisions, and enable discipline",0),
     ("Nepal: a software firm adopting an ACM-style internal code; CAN community norms",1),
     ("Mini case: ACM 'avoid harm' — a dev asked to build a HIDDEN tracking feature; what does the code say?",1)],
    notes="Work the 'avoid harm' case live — the code converts a vague feeling into a clear obligation to refuse or disclose.")
diag_cert_vs_license()
add_content("Certification · Licensing · Professional Bodies","S7 · Concepts 3–5  [THEORY]  ·  ~9 min",
    [("Certification = vendor/industry proof of a SKILL (CCNA, CISSP, AWS); voluntary; not a license",0),
     ("Misconception: 'certification = ethical behavior.' It proves skill, NOT integrity.",2),
     ("Licensing = government permission with LEGAL force; protects the public (US 'PE-Software' debate)",0),
     ("Professional bodies set standards, run certs, publish codes: ACM, IEEE-CS, ISACA, CAN-Nepal",0),
     ("Analogy: a code of ethics is the 'terms of service' you accept to be trusted as a pro",2)],
    notes="Lean on the comparison table: certification = voluntary proof of skill; licensing = legal right to practice. Nepali grads chase CCNA/CEH for outsourcing work — a cert signals competence, not character.")
add_activity("S7 — 'License IT? — quick debate'  ·  ~5 min",
    [("Split the room: half argue FOR licensing IT workers in Nepal, half AGAINST (2 min prep)",0),
     ("Take 2 points from each side (3 min)",1),
     ("Surface the trade-off: public safety vs innovation/cost/who-decides-the-standard",1)],
    notes="Primes the S7 discussion question. Good FOR points: safety-critical systems, accountability. AGAINST: IT changes too fast, stifles innovation, who sets the syllabus?")
add_quiz("S7 — Quick Check  ·  ~5 min",
    [("Q1. A government-granted, legally enforced permission to practice is…","q"),
     ("✅ b) A license","a"),
     ("Q2. ACM and IEEE-CS are examples of…","q"),
     ("✅ c) Professional organizations (with codes of ethics)","a"),
     ("Discussion: should IT workers in Nepal be licensed like engineers and doctors?","o")],
    notes="APPLICATION [~3 min]: employers/outsourcing clients want certs + adherence to a code — shapes hireability and pay. SUMMARY [~2 min]: IT = partial profession; codes guide, certs prove skill, licensing has legal force; bodies set standards. Next: keeping USERS ethical.")

# S8
add_divider("Session 8 · Lecture hour 3","Ethical Use of IT Resources among Users",
            "HOOK [~5 min]: one shared Wi-Fi password — someone torrents movies, another installs cracked Photoshop, a third runs a crypto miner on the server overnight. No outsider hacked anything. Whose problem is it? Users need rules, not just tools.")
add_content("IT Users & the Insider Risk","S8 · Concept 1  [THEORY]  ·  ~8 min",
    [("IT user = anyone using the org's IT resources (staff, students, customers)",0),
     ("Most security & ethics incidents start with ordinary USERS, not hackers — the insider risk",2),
     ("Nepal: a college lab where students install games/cracked software on shared PCs",1),
     ("Misconception: 'only the IT dept is responsible for misuse.' Every user shares responsibility.",2)],
    notes="Reframe security as everyone's job, not just IT's. The insider risk is the thread for the whole session.")
add_content("Acceptable Use Policy (AUP)","S8 · Concept 2  [THEORY] [EXAMPLE]  ·  ~8 min",
    [("AUP = a document of what users MAY and MAY NOT do with the org's IT resources",0),
     ("Defines: permitted use, prohibited use, monitoring, and consequences — users sign/acknowledge",0),
     ("Nepal: NTC/Ncell staff AUP bans personal USBs & unlicensed software; a university IT policy",1)],
    notes="See the allowed/prohibited card next slide. Tell students they'll sign one of these on day one of any job.")
diag_aup_card()
add_content("Software Piracy · Inappropriate Use","S8 · Concepts 3–4  [THEORY] [EXAMPLE]  ·  ~9 + ~5 min",
    [("Piracy = copying/installing/distributing software without a valid license (→ Unit 3 IP)",0),
     ("Nepal: high piracy rates; offices on cracked Windows/Office; 'everyone does it' normalization",1),
     ("Misconception: 'it's only piracy if I sell it.' Using an unlicensed copy is ALREADY a violation.",2),
     ("Analogy: cracked software = sneaking into a movie hall — cheaper today, never entitled",2),
     ("Inappropriate use = org resources for unauthorized/personal/harmful purposes",0)],
    notes="Piracy (~9 min): forms — counterfeiting, unlicensed installs, license violation; legal AND ethical harm. Inappropriate use (~5 min): the line-drawing case — is checking eSewa + a family message at work 'inappropriate'? Discuss reasonableness vs the AUP.")
add_activity("S8 — 'Draft the lab AUP'  ·  ~5 min",
    [("In pairs (2 min), then collect aloud (3 min)",0),
     ("Write 3 rules you'd put in YOUR college lab's Acceptable Use Policy",1),
     ("For each: state the rule, why it exists, and the consequence for breaking it",1)],
    notes="Concrete and fun. Watch for rules with no enforcement (a 'poster') — tie back to Unit 1's enforceability point.")
add_quiz("S8 — Quick Check  ·  ~5 min",
    [("Q1. The document defining permitted/prohibited IT use is the…","q"),
     ("✅ b) Acceptable Use Policy (AUP)","a"),
     ("Q2. Installing a cracked Office copy for work (not selling) is…","q"),
     ("✅ b) Software piracy","a"),
     ("Discussion: is piracy in Nepal mainly an ethics, affordability, or enforcement problem?","o")],
    notes="APPLICATION [~3 min]: you'll sign an AUP on day one; breaking it (piracy/misuse) is a common reason people are fired or sued. SUMMARY [~2 min]: users cause most incidents; AUPs set the rules; piracy + inappropriate use are the big two. Next: the user's flip side — privacy.")

# S9
add_divider("Session 9 · Lecture hour 4","Key Privacy and Anonymity Issues",
            "HOOK [~5 min]: you fill a loan form, hand your number to a SIM dealer, scan a QR to enter a mall — a week later you're flooded with spam loan calls. Your data traveled somewhere. When did you consent?")
add_content("Data Privacy","S9 · Concept 1  [THEORY] [EXAMPLE]  ·  ~9 min",
    [("Data privacy = your right to CONTROL how personal data is collected, used, shared",0),
     ("Principles: consent, purpose limitation, data minimization",0),
     ("Nepal: KYC data at banks/telecoms; Individual Privacy Act 2075 (2018); citizen-data leaks",1),
     ("Misconception: 'nothing to hide, so privacy doesn't matter.' Privacy = control & power, not hiding.",2)],
    notes="Deliver the 'locked door' rebuttal to 'nothing to hide': you lock your door not because you're doing wrong inside, but because it's yours to control.")
add_content("Surveillance","S9 · Concept 2  [THEORY]  ·  ~8 min",
    [("Surveillance = systematic monitoring of activities, communications, or location",0),
     ("Forms: workplace monitoring, CCTV, government surveillance, app/location tracking",0),
     ("Core tension: security vs privacy",2),
     ("Nepal: rising CCTV in Kathmandu; employer email/screen monitoring; SIM/biometric registration",1)],
    notes="Mini case: employer logs every keystroke + screenshot — ethical? Does telling staff in advance change your answer? See the data-trail diagram next.")
diag_data_trail()
add_content("Anonymity · Identity & Identity Theft","S9 · Concepts 3–4  [THEORY] [EXAMPLE]  ·  ~8 + ~5 min",
    [("Anonymity = acting without revealing identity — a DOUBLE-EDGED tool",0),
     ("Protects whistle-blowers/activists; shields trolls/scammers/criminals",1),
     ("Misconception: 'I'm completely anonymous online.' IP logs/metadata/fingerprints leave a trail.",2),
     ("Identity = the data representing you (citizenship no., eSewa/bank logins, biometrics)",0),
     ("Nepal: fake FB profiles of public figures; phishing SMS posing as eSewa/bank to steal OTPs",1),
     ("Analogy: your data is a house key — easy to copy, noticed only once someone's inside",2)],
    notes="Anonymity (~8 min): same tool exposes corruption AND spreads rumors — Nepali examples both ways. Identity (~5 min): identity theft via phishing previews Unit 7.")
add_activity("S9 — 'Audit your data trail'  ·  ~5 min",
    [("In pairs (2 min), then share (3 min)",0),
     ("List 3 places this WEEK you handed over personal data (form, app, QR, SIM, KYC)",1),
     ("For each: did you give real consent? did they collect more than needed? where might it go?",1)],
    notes="Makes the data-trail diagram personal. Connects consent + data minimization to their own week.")
add_quiz("S9 — Quick Check  ·  ~5 min",
    [("Q1. 'Collect only the data you need' is the principle of…","q"),
     ("✅ b) Data minimization","a"),
     ("Q2. Online anonymity is best described as…","q"),
     ("✅ c) A double-edged tool","a"),
     ("Discussion: should the government require real-name registration to post on social media?","o")],
    notes="APPLICATION [~3 min]: you'll handle customer data; mishandling breaks Nepal's privacy law and destroys trust. SUMMARY [~2 min]: privacy = control; surveillance trades security for privacy; anonymity cuts both ways; identity is the new target. Next: where privacy & harm collide — social media.")

# S10
add_divider("Session 10 · Lecture hour 5 — CLOSES UNIT 2","Social Networking Ethical Issues",
            "HOOK [~5 min]: a doctored celebrity screenshot goes viral before lunch; by evening a teen is mocked in a school group chat and a 'breaking news' rumor sparks panic. Nobody hacked a server — ordinary users did all of it. Social media weaponizes everyday people.")
add_content("Why Social Networking is an Ethics Hotspot","S10 · Concept 1  [THEORY]  ·  ~8 min",
    [("Social networking = platforms to create, share, connect (FB, TikTok, Instagram, X, YouTube)",0),
     ("Scale + speed + anonymity + virality turn small acts into LARGE harm; online bystander effect",2),
     ("Nepal: very high FB/TikTok usage; viral content shaping opinion and elections",1),
     ("Misconception: 'it's just online, it isn't real.' Real legal, social & psychological harm.",2)],
    notes="Tie back to S5 of Unit 1 (scale/speed/permanence). The bystander effect sets up the cyberbullying case.")
add_content("Cyberbullying · Fake News","S10 · Concepts 2–3  [THEORY] [EXAMPLE]  ·  ~8 min",
    [("Cyberbullying/harassment: trolling, doxxing, group pile-ons, gendered harassment; mental-health harm",0),
     ("Nepal: harassment of women/public figures in comment sections; school group-chat bullying",1),
     ("Misinformation (careless) vs disinformation (deliberate); algorithms reward outrage",0),
     ("Misconception: 'sharing isn't creating, so I'm not responsible.' Forwarding spreads the harm.",2),
     ("Analogy: a lie online travels like a forwarded festival greeting — impossible to recall",2)],
    notes="Bystander mini case: you stay silent while a group chat mocks a classmate — are you complicit? Nepal: disaster rumors, fake cures, election misinformation.")
diag_defamation_spectrum()
add_content("Defamation & Reputation Harm","S10 · Concept 4  [THEORY] [EXAMPLE]  ·  ~5 min",
    [("Defamation = a false statement of FACT, presented as true, that damages reputation",0),
     ("Opinion/criticism ≠ defamation — the line is fact vs opinion",2),
     ("Nepal: arrests under the ETA 2063 for defamatory FB posts; free-speech-vs-misuse debate",1),
     ("Case: 'shop X overcharges' (opinion) vs 'shop X is a fraud that steals' (defamation)",1)],
    notes="VERIFY ETA 2063 specifics against current law before teaching (active reform debate). Full law treatment is Unit 9.")
add_activity("S10 — 'Verify before you share'  ·  ~6 min",
    [("In pairs (2 min), then share (4 min)",0),
     ("Build a 2-step check you'd run before sharing a shocking post (source? evidence? reverse-image?)",1),
     ("Apply it to a recent Nepali viral rumor — would your check have stopped it?",1),
     ("Then classify 4 sample posts: opinion / criticism / harassment / defamation",1)],
    notes="The classify-4-posts task is the assessment skill. Close the unit on personal responsibility — the share button is an ethical act.")
add_quiz("S10 — Quick Check  ·  ~5 min",
    [("Q1. Deliberately spreading false info to mislead is…","q"),
     ("✅ b) Disinformation","a"),
     ("Q2. A false factual statement that damages reputation is…","q"),
     ("✅ c) Defamation","a"),
     ("Discussion: 2 checks to do before sharing a shocking post — would they stop a recent Nepali rumor?","o")],
    notes="APPLICATION [~3 min]: public posts are part of your professional reputation; in Nepal a defamatory/harassing post can lead to arrest under the ETA. SUMMARY [~2 min]: social media scales everyday acts; bullying/fake news/defamation are the core issues; verify before you share. Next unit: Intellectual Property.")

# End
add_divider("Assessment","Unit 2 — End-of-Unit Quiz",
            "12 MCQ + 5 short answer + 2 applied cases + 1 discussion. Full questions and answer key are in Unit2_material.md.")
add_content("Exam pointers — what to master","Unit 2 · Review",
    [("The 6 IT-worker relationships + one ethical risk each; gift vs bribe",0),
     ("Certification vs licensing (granted by / legal force / proves what)",0),
     ("What an AUP contains; software piracy (using unlicensed = violation)",0),
     ("Privacy = control; data minimization; anonymity as double-edged",0),
     ("Misinformation vs disinformation vs defamation; opinion ≠ defamation",0)],
    notes="These five lines map to the most common exam questions. Full quiz + answer key in the markdown file.")
add_title("End of Unit 2",
          "Next: Unit 3 — Intellectual Property (copyright, patent, trade secret, trademark)",
          "Generated for IT 246 · diagrams are native shapes (editable) · speaker notes in each slide's Notes pane")

out="/Users/inventechg1/Desktop/2083_SEM/sixth/IT246_Unit2.pptx"
prs.save(out)
print("Saved",out,"with",len(prs.slides._sldIdLst),"slides")
