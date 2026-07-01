#!/usr/bin/env python3
"""Generate IT246 Unit 1 PowerPoint deck (with native-shape diagrams).
Run: python3 build_unit1_pptx.py  ->  IT246_Unit1.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ---- palette ----
NAVY   = RGBColor(0x0C, 0x2B, 0x4A)
BLUE   = RGBColor(0x18, 0x5F, 0xA5)
TEAL   = RGBColor(0x0F, 0x6E, 0x56)
AMBER  = RGBColor(0x85, 0x4F, 0x0B)
CORAL  = RGBColor(0x99, 0x3C, 0x1D)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF8)
PALET  = RGBColor(0xE1, 0xF5, 0xEE)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GREY   = RGBColor(0x55, 0x60, 0x6B)
LGREY  = RGBColor(0xC7, 0xCF, 0xD6)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame


def _bar(slide, color, l, t, w, h):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _settext(sp, text, font=WHITE, size=14, bold=True, align=PP_ALIGN.CENTER):
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    lines = text.split("\n")
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = lines[0]
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = font
    for extra in lines[1:]:
        p2 = tf.add_paragraph(); p2.alignment = align
        r = p2.add_run(); r.text = extra
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = font


def shp(slide, kind, l, t, w, h, text="", fill=BLUE, font=WHITE, size=14,
        bold=True, line=None, align=PP_ALIGN.CENTER):
    sp = slide.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    if text:
        _settext(sp, text, font, size, bold, align)
    return sp


def conn(slide, x1, y1, x2, y2, color=GREY, w=1.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w); c.shadow.inherit = False
    return c


# ---------- scaffolds ----------
def add_title(title, subtitle, footer):
    s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
    _bar(s, AMBER, 0.9, 3.05, 1.6, 0.10)
    tf = _box(s, 0.9, 2.0, 11.5, 1.0); tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = WHITE
    tf2 = _box(s, 0.9, 3.25, 11.5, 1.2); tf2.word_wrap = True
    r = tf2.paragraphs[0].add_run(); r.text = subtitle
    r.font.size = Pt(20); r.font.color.rgb = RGBColor(0xCF, 0xDD, 0xEA)
    tf3 = _box(s, 0.9, 6.5, 11.5, 0.6)
    r = tf3.paragraphs[0].add_run(); r.text = footer
    r.font.size = Pt(13); r.font.color.rgb = RGBColor(0x9A, 0xB2, 0xC6)
    return s


def add_divider(kicker, title, notes=""):
    s = prs.slides.add_slide(BLANK); _bg(s, BLUE)
    tfk = _box(s, 0.9, 2.4, 11.5, 0.6)
    r = tfk.paragraphs[0].add_run(); r.text = kicker.upper()
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = RGBColor(0xBF, 0xDC, 0xF2)
    tft = _box(s, 0.9, 3.0, 11.5, 1.6); tft.word_wrap = True
    r = tft.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(33); r.font.bold = True; r.font.color.rgb = WHITE
    if notes: _notes(s, notes)
    return s


def _header(s, kicker, header, kicker_color=BLUE):
    _bar(s, kicker_color, 0, 0, SW.inches, 0.18)
    tfk = _box(s, 0.7, 0.32, 12.0, 0.45)
    r = tfk.paragraphs[0].add_run(); r.text = kicker.upper()
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = kicker_color
    tfh = _box(s, 0.7, 0.7, 12.0, 1.0); tfh.word_wrap = True
    r = tfh.paragraphs[0].add_run(); r.text = header
    r.font.size = Pt(27); r.font.bold = True; r.font.color.rgb = NAVY


def add_content(header, kicker, blocks, notes="", kicker_color=BLUE):
    s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
    _header(s, kicker, header, kicker_color)
    tf = _box(s, 0.7, 1.8, 11.9, 5.3); tf.word_wrap = True
    first = True
    for text, level in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        r = p.add_run(); r.text = text
        if level == 0:
            r.font.size = Pt(18); r.font.color.rgb = DARK; p.space_after = Pt(7); p.level = 0
        elif level == 1:
            r.font.size = Pt(15.5); r.font.color.rgb = GREY; p.space_after = Pt(3); p.level = 1
        else:
            r.font.size = Pt(15.5); r.font.color.rgb = TEAL; r.font.italic = True
            p.space_after = Pt(6); p.level = 1
    if notes: _notes(s, notes)
    return s


def add_quiz(title, items, notes=""):
    s = prs.slides.add_slide(BLANK); _bg(s, LIGHT)
    _bar(s, AMBER, 0, 0, SW.inches, 0.18)
    tfk = _box(s, 0.7, 0.32, 12, 0.45)
    r = tfk.paragraphs[0].add_run(); r.text = "CHECK FOR UNDERSTANDING"
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = AMBER
    tfh = _box(s, 0.7, 0.7, 12, 0.9)
    r = tfh.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = NAVY
    tf = _box(s, 0.7, 1.8, 11.9, 5.2); tf.word_wrap = True
    first = True
    for text, kind in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        r = p.add_run(); r.text = text
        if kind == "q":
            r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = NAVY; p.space_before = Pt(8)
        elif kind == "a":
            r.font.size = Pt(16); r.font.color.rgb = TEAL; r.font.bold = True; p.level = 1
        else:
            r.font.size = Pt(16); r.font.color.rgb = DARK; p.level = 1
    if notes: _notes(s, notes)
    return s


def add_activity(title, blocks, notes=""):
    s = prs.slides.add_slide(BLANK); _bg(s, PALET)
    _bar(s, TEAL, 0, 0, SW.inches, 0.18)
    tfk = _box(s, 0.7, 0.32, 12, 0.45)
    r = tfk.paragraphs[0].add_run(); r.text = "🛠  IN-CLASS ACTIVITY"
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = TEAL
    tfh = _box(s, 0.7, 0.7, 12, 0.9)
    r = tfh.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = NAVY
    tf = _box(s, 0.7, 1.85, 11.9, 5.2); tf.word_wrap = True
    first = True
    for text, level in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        r = p.add_run(); r.text = text
        if level == 0:
            r.font.size = Pt(18); r.font.color.rgb = DARK; p.space_after = Pt(8)
        else:
            r.font.size = Pt(15.5); r.font.color.rgb = GREY; p.space_after = Pt(4); p.level = 1
    if notes: _notes(s, notes)
    return s


def diagram_slide(kicker, header, kicker_color=BLUE):
    s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
    _header(s, kicker, header, kicker_color)
    return s


# =================== DIAGRAMS ===================
def diag_legal_ethical_grid():
    s = diagram_slide("S1 · Concept 1  ·  Diagram", "Legal vs Ethical — two different axes")
    cx0, cw, cy0, ch, gap = 3.1, 4.35, 2.7, 1.85, 0.18
    # column headers
    for i, lab in enumerate(["LEGAL", "ILLEGAL"]):
        t = _box(s, cx0 + i*(cw+gap), 2.2, cw, 0.4)
        r = t.paragraphs[0].add_run(); r.text = lab
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = NAVY
        t.paragraphs[0].alignment = PP_ALIGN.CENTER
    # row headers
    for j, lab in enumerate(["ETHICAL", "UNETHICAL"]):
        t = _box(s, 1.0, cy0 + j*(ch+gap) + ch/2 - 0.3, 1.9, 0.6)
        r = t.paragraphs[0].add_run(); r.text = lab
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = NAVY
        t.paragraphs[0].alignment = PP_ALIGN.CENTER
    cells = [  # (col, row, text, color)
        (0, 0, "Paying taxes honestly\nKeeping a promise to a customer", TEAL),
        (1, 0, "Justified civil disobedience\n(rare — illegal but morally defended)", BLUE),
        (0, 1, "Exploiting a tax loophole\nQueue-jumping via 'source-force'", AMBER),
        (1, 1, "Paying a bribe ('ghus')\nLeaking someone's private data", CORAL),
    ]
    for col, row, txt, color in cells:
        shp(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx0 + col*(cw+gap), cy0 + row*(ch+gap), cw, ch,
            txt, fill=color, size=13.5)
    _notes(s, "The whole point of S1: legal and ethical are SEPARATE axes. Walk each quadrant. The dangerous quadrant is legal+unethical — that's where most business scandals live.")


def diag_csr_pyramid():
    s = diagram_slide("S2 · Concept 1  ·  Diagram", "Carroll's CSR Pyramid")
    layers = [  # bottom to top: (label, width, color)
        ("ECONOMIC — be profitable (the foundation)", 9.4, NAVY),
        ("LEGAL — obey the law", 7.4, BLUE),
        ("ETHICAL — do what is right, beyond the law", 5.4, TEAL),
        ("PHILANTHROPIC — give back to the community", 3.4, AMBER),
    ]
    h = 1.05; y = 6.3
    for label, w, color in layers:
        x = (13.333 - w) / 2
        shp(s, MSO_SHAPE.TRAPEZOID if w == 3.4 else MSO_SHAPE.RECTANGLE, x, y - h, w, h,
            label, fill=color, size=13)
        y -= (h + 0.04)
    cap = _box(s, 0.7, 1.95, 11.9, 0.5)
    r = cap.paragraphs[0].add_run()
    r.text = "Lower layers must hold before the upper ones mean anything — you can't be philanthropic if you're not first viable and legal."
    r.font.size = Pt(13.5); r.font.italic = True; r.font.color.rgb = GREY
    _notes(s, "Build bottom-up. NRB directs banks to spend on CSR — that sits in the ethical/philanthropic layers, but only works because the bank is profitable and legal first.")


def diag_ethics_cycle():
    s = diagram_slide("S3 · Concept 2  ·  Diagram", "An Ethics Program is a Cycle, not a Poster")
    steps = [("CODE\nof conduct", BLUE), ("TRAINING", TEAL), ("REPORTING\n(hotline)", AMBER),
             ("AUDIT", CORAL), ("CONSEQUENCES", NAVY)]
    w, h, y = 2.25, 1.5, 3.2; gap = 0.12
    total = 5*w + 4*gap; x0 = (13.333 - total)/2
    for i, (txt, col) in enumerate(steps):
        shp(s, MSO_SHAPE.PENTAGON, x0 + i*(w+gap), y, w, h, txt, fill=col, size=13.5)
    cap = _box(s, 0.7, 5.1, 11.9, 0.7)
    r = cap.paragraphs[0].add_run()
    r.text = "↺  …and back to refining the Code. A document on the wall changes nothing — the cycle is what changes behavior."
    r.font.size = Pt(14); r.font.italic = True; r.font.color.rgb = GREY
    _notes(s, "Stress the loop. The unequal-punishment mini case shows what happens when 'Consequences' is applied selectively — the whole cycle loses credibility.")


def diag_decision_flow():
    s = diagram_slide("S4 · Concept 1  ·  Diagram", "The 5-Step Ethical Decision Process")
    steps = ["1\nGet the\nfacts", "2\nStakeholders\n& options", "3\nEvaluate\n(lenses)",
             "4\nChoose\n& act", "5\nReflect"]
    w, h, y = 2.15, 1.5, 2.6; gap = 0.32
    total = 5*w + 4*gap; x0 = (13.333 - total)/2
    for i, txt in enumerate(steps):
        shp(s, MSO_SHAPE.ROUNDED_RECTANGLE, x0 + i*(w+gap), y, w, h, txt, fill=BLUE, size=14)
        if i < 4:
            ax = x0 + i*(w+gap) + w + 0.02
            shp(s, MSO_SHAPE.RIGHT_ARROW, ax, y + h/2 - 0.16, gap - 0.04, 0.32, fill=LGREY)
    # lenses row
    lens = [("Utilitarian", "greatest good"), ("Rights / duty", "respect rights"),
            ("Fairness", "treat alike"), ("Virtue / common good", "good character")]
    lw, ly = 2.7, 4.7; lgap = 0.22; lt = 4*lw + 3*lgap; lx0 = (13.333 - lt)/2
    cap = _box(s, lx0, 4.25, lt, 0.4)
    r = cap.paragraphs[0].add_run(); r.text = "Step 3 — test each option through these lenses:"
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = NAVY
    for i, (name, test) in enumerate(lens):
        shp(s, MSO_SHAPE.ROUNDED_RECTANGLE, lx0 + i*(lw+lgap), ly, lw, 0.95,
            name + "\n(" + test + ")", fill=TEAL, size=12.5)
    _notes(s, "This 5-step spine returns in every later unit. Run the leaked-exam-answers scenario through it live. Lenses can disagree — that's fine; the point is defensible reasoning.")


def diag_papa():
    s = diagram_slide("S5 · Concept 2  ·  Diagram", "The Four Pillars of IT Ethics — 'PAPA'")
    cells = [
        (3.0, 2.55, "PRIVACY", "Who can collect & see your data?  (Unit 2)", BLUE),
        (7.5, 2.55, "ACCURACY", "Who's liable when data / algorithms are wrong?", TEAL),
        (3.0, 4.45, "PROPERTY", "Who owns information & software?  (Unit 3)", AMBER),
        (7.5, 4.45, "ACCESS", "Who gets to use technology — and who's left out?", CORAL),
    ]
    for x, y, head, sub, color in cells:
        sp = shp(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 4.3, 1.75, "", fill=color)
        tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = head; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run(); r.text = sub; r.font.size = Pt(12.5); r.font.color.rgb = WHITE
    cap = _box(s, 0.7, 6.45, 11.9, 0.6)
    r = cap.paragraphs[0].add_run()
    r.text = "Security threads through all four (Units 5–8).  Core question of the course:  just because we CAN — should we?"
    r.font.size = Pt(13.5); r.font.italic = True; r.font.color.rgb = GREY
    _notes(s, "PAPA = Privacy, Accuracy, Property, Access. This quadrant previews the whole course. End on 'can vs should' with the 100,000-user scraping case.")


# =================== BUILD DECK ===================
add_title(
    "Unit 1 — An Overview of Ethics",
    "IT 246: IT Ethics and Cybersecurity  ·  BIM 6th Semester",
    "5 lecture hours (S1–S5)  ·  50 min each  ·  Nepal / South Asia context",
)
add_content(
    "Unit 1 — Learning Outcomes", "Overview",
    [("By the end of this unit, you will be able to:", 0),
     ("Define ethics and distinguish it from morals, law, and etiquette", 1),
     ("Explain ethics in the business world and Corporate Social Responsibility (CSR)", 1),
     ("Describe how organizations foster CSR and improve business ethics", 1),
     ("Apply a structured ethical decision-making process to a dilemma", 1),
     ("Explain why ethics matters specifically in Information Technology", 1)],
    notes="Five sessions, one outcome each. This unit sets the reasoning toolkit reused across all 9 units.",
)

# ---------------- S1 (50 min: 5 hook + 35 content + 5 CFU + 3 app + 2 summary) ----------------
add_divider("Session 1 · Lecture hour 1", "What is Ethics? · Ethics in Business",
            "HOOK [~5 min]: Read 3 scenarios — sharing one paid login among 8; taking 'chiya kharcha' to push a file; dumping e-waste in the river. Run TWO votes: 'which are illegal?' then 'which are unethical?' The counts won't match — that mismatch IS the lesson. Don't resolve yet; cold-call a student who split their vote. Write agenda on board.")
add_content("Ethics, Morals, Law & Etiquette", "S1 · Concept 1  [THEORY]  ·  ~8 min",
    [("Ethics — a code of right/wrong owned by a GROUP, profession, or society (it's shared)", 0),
     ("Morals — an individual's PERSONAL beliefs (family, religion, upbringing)", 0),
     ("Law — rules the STATE enforces, with penalties (codified, slow, can't cover everything)", 0),
     ("Etiquette — social MANNERS; breaking it is rude, not deeply wrong", 0),
     ("They usually overlap — the cases that matter are where they PULL APART", 2),
     ("Misconception: 'if it's legal, it's ethical.' Law is the FLOOR, not the ceiling.", 2)],
    notes="Write the 4 words on the board. Define each, but spend the time on how they RELATE: must-do (law) vs should-do (ethics) vs personally-believe (morals) vs polite (etiquette). Nepal: bribe = illegal AND unethical; 'source-force' queue-jump = unethical not illegal; left-hand at a feast = etiquette only. Say aloud: 'the law is the minimum society tolerates, not the maximum you should aim for.' Transition: 'if legal ≠ ethical, we need to SEE the difference — the grid.'")
diag_legal_ethical_grid()
add_content("Ethics in the Business World", "S1 · Concept 3  [THEORY]  ·  ~8 min",
    [("Business ethics = applying ethical standards to commercial conduct", 0),
     ("A business serves MANY stakeholders at once (customers, staff, society, owners) — they conflict", 0),
     ("Central tension: short-term profit vs long-term TRUST", 2),
     ("Trust is slow to build, instant to lose — unethical shortcuts spend a years-long reserve", 0),
     ("Nepal: festival shortage — Shop A triples prices (wins today); Shop B holds (longer queue every Dashain)", 1)],
    notes="Explain WHY business is special: customers can't verify most claims (does the medicine work? is my data safe?), so business runs on trust. Walk the two-shop example fully. Mini case (take 2–3 answers): a trekking agency hides a route safety risk to keep a booking — who bears the cost? Client (safety), agency (reputation/lawsuits), AND the whole tourism sector. One shortcut, an industry's loss.")
add_content("Why Good People & Businesses Still Go Wrong", "S1 · Concept 4  [THEORY]  ·  ~7 min",
    [("Most misbehavior comes from ordinary people under predictable pressure — not 'bad people'", 0),
     ("Results pressure — the target crowds out the ethics", 0),
     ("'Everyone does it' — normalization makes wrong feel normal", 0),
     ("Weak oversight + diffusion of responsibility — 'I just followed the process'", 0),
     ("If misbehavior is SITUATIONAL, good systems & culture prevent it (→ S2, S3)", 2),
     ("Analogy: ethics = the BRAKES that let you drive fast safely — not a roadblock", 2)],
    notes="Correct the 'bad apples' assumption — misbehavior is mostly situational, which is exactly why S2/S3 (systems & culture) matter. Deliver the brakes analogy with conviction: take the brakes off and you can't speed, you can only crash.")
add_activity("S1 — 'Find the grey area'  ·  ~6 min",
    [("Think–Pair–Share (2 min in pairs, 3 min sharing, 1 min close)", 0),
     ("Each pair invents ONE act that is legal-but-unethical and ONE illegal-but-arguably-ethical", 1),
     ("Draw from Nepali daily life or IT", 1),
     ("Take 4–5 answers aloud and place each on the 2×2 grid", 1),
     ("Close: notice how EASY legal-but-unethical was to find — that's the quadrant this course trains", 1)],
    notes="Seeds if pairs stall: selling legally-collected user data; a politician's legal conflict of interest; pirating unaffordable software (illegal, students argue ethical). Keep it moving — this activity is what fills the time and makes the grid stick.")
add_quiz("S1 — Quick Check  ·  ~5 min",
    [("Q1. An action that is legal but still wrong is best called…", "q"),
     ("✅ b) Unethical", "a"),
     ("Q2. A profession's written code of conduct is an example of…", "q"),
     ("✅ c) Ethics (a group code)", "a"),
     ("Discussion: name one thing legal in Nepal you consider unethical — which grid quadrant?", "o")],
    notes="APPLICATION [~3 min]: every job offer has a code of conduct; employers check your online reputation BEFORE the interview; recruiters ask 'tell me about an ethical dilemma'. SUMMARY [~2 min]: ethics = group code; legal ≠ ethical (watch the legal-but-unethical quadrant); business ethics = profit vs trust, misbehavior is situational so systems fix it. Next: CSR.")

# ---------------- S2 ----------------
add_divider("Session 2 · Lecture hour 2", "Corporate Social Responsibility (CSR)",
            "HOOK [~5 min]: Ncell funds a school; a bank plants trees; a firm sponsors flood relief — kindness, marketing, or duty? Show of hands for each. Most vote 'marketing' cynically — acknowledge it, then complicate: what if it's all three AND society can legitimately expect it? Nepal actually REQUIRES some CSR. Keep the cynicism alive for the closing discussion.")
add_content("What CSR Means", "S2 · Concept 1  [THEORY]  ·  ~7 min",
    [("CSR = a company's DUTY to act for society's benefit, not only shareholders", 0),
     ("Key word: DUTY — business owes the society that lets it operate (infrastructure, workers, trust)", 0),
     ("Modern consensus shifted from 'business of business is only profit' → profit AND responsibility", 2)],
    notes="Contrast the old Friedman 'only profit' view with the modern consensus. The crucial idea is that society GRANTS a business its licence to operate, so responsibility flows back.")
diag_csr_pyramid()
add_content("Stakeholder Theory", "S2 · Concept 3  [THEORY]  ·  ~8 min",
    [("To whom does a company answer? Narrow answer: shareholders (owners)", 0),
     ("Stakeholder theory widens it: EVERYONE materially affected — customers, staff, suppliers,", 0),
     ("   community, environment, government, AND owners", 1),
     ("Consequence: a decision great for owners but harmful to the community is a FAILURE", 2),
     ("Mini case: a factory profitable for owners but polluting a village's water", 1)],
    notes="Evaluate the factory TWICE, out loud: shareholder lens = success (profit up); stakeholder lens = failure (village harmed; costs hidden not absent — health, lawsuits, lost social licence). Land it: 'ethics fails when we count only the stakeholders on the balance sheet.'")
add_content("Fostering CSR & an Ethical Culture", "S2 · Concept 4  [THEORY]  ·  ~7 min",
    [("Tone from the top — staff copy what leaders DO, not what posters say", 0),
     ("A clear, public code of ethics — turns vague values into specific expectations", 0),
     ("An ethical culture — raising a concern is rewarded, doing right is the easy path", 0),
     ("Social-responsibility reporting — publishing impact creates accountability", 0),
     ("'Good ethics is good business' — trust compounds like interest; a scandal is a withdrawal", 2)],
    notes="Deliver each lever WITH its why. NRB directs Nepali banks to allocate profit to CSR — so part of CSR here is regulated, not optional.")
add_activity("S2 — 'Genuine or greenwashing?'  ·  ~6 min",
    [("Think–Pair–Share (2 min pairs, 4 min sharing)", 0),
     ("Each pair names ONE Nepali company's CSR activity they've seen", 1),
     ("Class judges 2–3 examples: genuine responsibility or image-building?", 1),
     ("Test to apply: does the company's EVERYDAY operation match the CSR message?", 1),
     ("e.g. a bank advertising 'we care about farmers' while its loan terms crush them = greenwashing", 1)],
    notes="This reuses the hook's cynicism productively. The 'does daily operation match the message' test is the takeaway tool.")
add_quiz("S2 — Quick Check  ·  ~5 min",
    [("Q1. The base of Carroll's CSR pyramid is…", "q"),
     ("✅ d) Economic responsibility", "a"),
     ("Q2. Stakeholders include all EXCEPT…", "q"),
     ("✅ c) only the shareholders", "a"),
     ("Discussion: a Nepali company doing visible CSR — genuine or image-building? Apply the test.", "o")],
    notes="APPLICATION [~3 min]: CSR shapes which firms attract talent/investment and which face boycotts; as a manager you can build real vs decorative CSR. SUMMARY [~2 min]: CSR = duty to society (Carroll's 4 layers); stakeholders = all affected; good ethics is built via leadership/culture/codes/transparency. Next: the systems that turn values into behavior.")

# ---------------- S3 ----------------
add_divider("Session 3 · Lecture hour 3", "Improving Business Ethics",
            "HOOK [~5 min]: a company hires consultants, writes a glossy Code of Ethics, frames it in reception… and behaves exactly as before. Ask: WHY did the document change nothing? Take 2–3 answers, then pull together: a value written down ≠ a value practiced. Ethics needs SYSTEMS, not posters.")
add_content("From Values to Behavior", "S3 · Concept 1  [THEORY]  ·  ~7 min",
    [("An ethics program = the SYSTEMS that turn stated values into actual behavior", 0),
     ("It bridges the gap between what a firm SAYS and what it does on a deadline Tuesday", 0),
     ("Two firms, identical code: one where leaders follow it + violations have consequences,", 0),
     ("   one where leaders quietly break it → same document, OPPOSITE cultures", 2)],
    notes="The identical-code contrast is the whole point: the document was never the cause, the system around it was.")
diag_ethics_cycle()
add_content("Components of an Effective Program", "S3 · Concept 2  [THEORY]  ·  ~9 min",
    [("Ethics officer/committee with REAL authority (not an ignored side-duty)", 0),
     ("Board-level oversight — so ethics can't be buried by middle management", 0),
     ("A specific code of conduct — 'no gifts over X', not just 'be good'", 0),
     ("Ongoing training — real scenarios, refreshed (not a one-time slideshow)", 0),
     ("Safe reporting (whistle-blower) channels — anonymous & retaliation-free or no one uses them", 0),
     ("Periodic ethics audits — check rules are FOLLOWED, not just posted", 0),
     ("Misconception: 'a code guarantees ethics.' Unenforced code is WORSE than none — signals rules are for show", 2)],
    notes="Deliver each component WITH its purpose. Nepal: a bank's AML/compliance training + grievance/whistle-blower hotline is exactly this machinery, and it's regulated.")
add_content("Measuring & Sustaining Ethics", "S3 · Concept 3  [THEORY]  ·  ~7 min",
    [("Ethics audits + anonymous culture surveys — surface problems before scandals", 0),
     ("Consistent consequences for violations — at EVERY level, especially the top", 0),
     ("Track reported concerns and act, so reporting feels worthwhile", 0),
     ("Make-or-break factor: CONSISTENCY — selective punishment destroys the program", 2),
     ("Mini case: senior let off, junior punished for the same rule — what does everyone learn?", 1)],
    notes="The unequal-punishment case: visible hypocrisy destroys trust FASTER than having no code at all, because now the rule is openly power-dependent.")
add_activity("S3 — 'Design one enforceable rule'  ·  ~5 min",
    [("In pairs (2 min), then 3 answers aloud (3 min)", 0),
     ("Design ONE concrete rule to make the college exam process more ethical", 1),
     ("AND state how it would be ENFORCED (a rule with no enforcement is a poster)", 1),
     ("Pressure-test each: who checks it? what's the consequence? could a senior dodge it?", 1)],
    notes="Ties straight back to the consistency point — push them on enforcement, not just the rule.")
add_quiz("S3 — Quick Check  ·  ~5 min",
    [("Q1. The most important driver of an ethical culture is…", "q"),
     ("✅ b) Leadership / tone at the top", "a"),
     ("Q2. A safe, retaliation-free way to report wrongdoing is a…", "q"),
     ("✅ b) Whistle-blower / reporting hotline", "a"),
     ("Discussion: which single component would do the most good in a typical Nepali office — why?", "o")],
    notes="APPLICATION [~3 min]: in your first weeks these signals reveal a healthy vs toxic employer (real reporting channel? rules apply to seniors? genuine training?). SUMMARY [~2 min]: values need systems; effective programs = leadership+code+training+reporting+audits; consistent enforcement beats any document. Next: a repeatable decision process.")

# ---------------- S4 ----------------
add_divider("Session 4 · Lecture hour 4", "Ethical Considerations in Decision Making",
            "HOOK [~5 min]: a classmate accidentally shares a Drive folder with TOMORROW'S exam answers. You have them. Don't ask 'what's right' — ask 'HOW would you even decide?' Take gut reactions, then point out everyone jumped to an answer with no METHOD. Today we build the method.")
add_content("The 5-Step Decision Process", "S4 · Concept 1  [THEORY]  ·  ~9 min",
    [("1. Get the facts — most 'dilemmas' dissolve with facts; never reason on assumptions", 0),
     ("2. Identify stakeholders & options — who's affected; what are the REALISTIC options?", 0),
     ("3. Evaluate options using ethical lenses (next slide)", 0),
     ("4. Choose and act — and own it", 0),
     ("5. Reflect — outcome? what would I do differently? (this is how judgment improves)", 0),
     ("Steps 1 & 2 are the ones people skip — and where most bad decisions are born", 2)],
    notes="This 5-step spine returns in EVERY later unit (IP, security, forensics, cyber law). Emphasize that 'get the facts' alone resolves most apparent dilemmas.")
diag_decision_flow()
add_content("Four Ethical Lenses", "S4 · Concept 2  [THEORY]  ·  ~9 min",
    [("Utilitarian — 'which option does the most good for the most people?'", 0),
     ("Rights/duty — 'does this respect everyone's rights & my duties, regardless of outcome?'", 0),
     ("Fairness/justice — 'am I treating similar people alike? is this fair?'", 0),
     ("Virtue/common good — 'what would a person of good character do?'", 0),
     ("Nepal: fuel-queue + ambulance — all four lenses CONVERGE, which is why it feels obvious", 1),
     ("Misconception: 'ethics is just opinion.' Lenses give defensible, comparable reasoning", 2)],
    notes="Walk the ambulance case through every lens (utilitarian: a life outweighs the wait; rights: right to life; fairness: an emergency genuinely IS a different case; virtue: a decent person waves them through). Then: the HARD cases are where lenses disagree — naming the disagreement is half the work.")
add_content("Common Decision Traps", "S4 · Concept 3  [THEORY]  ·  ~7 min",
    [("Rationalization — 'it's not really cheating, everyone uses these notes'", 0),
     ("Groupthink — going along because the whole team is", 0),
     ("Slippery slope — 'just this once' lowers the bar for next time", 0),
     ("Blind obedience — 'my manager told me to' (you're still responsible)", 0),
     ("Mini case: manager says 'fudge this one report, only this once' — name the traps", 2)],
    notes="Deliver each trap as it SOUNDS in real life. The fudge-the-report case = obedience + slippery slope; have students spot which pressure is being used.")
add_activity("S4 — 'Run the dilemma'  ·  ~6 min",
    [("In pairs (3 min), then 2 pairs walk through aloud (3 min)", 0),
     ("Take the leaked-exam-answers hook and work ALL 5 steps explicitly", 1),
     ("Name at least one LENS used in step 3", 1),
     ("Goal isn't unanimity — it's that everyone USED the method and can defend the choice", 1)],
    notes="This rehearses the process on a dilemma students actually feel. Don't force a 'correct' answer; reward defensible reasoning.")
add_quiz("S4 — Quick Check  ·  ~5 min",
    [("Q1. The first step in ethical decision making is to…", "q"),
     ("✅ b) Gather the facts", "a"),
     ("Q2. 'Greatest good for the greatest number' is which lens?", "q"),
     ("✅ c) Utilitarian", "a"),
     ("Discussion: a workplace example where two lenses DISAGREE — how would you decide?", "o")],
    notes="APPLICATION [~3 min]: grey-area calls happen constantly (a gift, a bent rule, a deadline shortcut); a repeatable process protects you and others; interviewers literally ask for it. SUMMARY [~2 min]: 5-step process; four lenses; four traps. Next: why this is sharper and bigger in IT.")

# ---------------- S5 ----------------
add_divider("Session 5 · Lecture hour 5 — CLOSES UNIT 1", "Ethics in Information Technology",
            "HOOK [~5 min]: an app silently sells your location; a loan algorithm rejects applicants by neighborhood (digital redlining); a deepfake goes viral the night before an election. Framing line: a shopkeeper cheats one customer at a time; bad code harms MILLIONS in seconds, permanently. Ask: what makes the IT versions WORSE?")
add_content("Why IT Ethics is Distinct", "S5 · Concept 1  [THEORY]  ·  ~9 min",
    [("Scale — one decision (a default, an algorithm) affects millions at once", 0),
     ("Speed — harm spreads in seconds, before anyone can intervene", 0),
     ("Anonymity — wrongdoers hide; victims targeted from a distance, across borders", 0),
     ("Permanence — leaked/posted data is almost impossible to erase", 0),
     ("Capability–regulation gap — tech outruns the law, so 'is it legal?' often has no answer yet", 0),
     ("'A rumor in a village fades; a screenshot online is forever.'", 2)],
    notes="Deliver each factor with a concrete consequence. Tie back to S1: 'the law is the floor' matters MOST in IT because the floor is often missing — your own ethics is all that's holding.")
diag_papa()
add_content("The IT Professional's Responsibility", "S5 · Concept 3  [THEORY]  ·  ~7 min",
    [("Defining question of IT ethics: you CAN — but SHOULD you?", 2),
     ("It lands on YOU: you'll choose what data to collect, how an algorithm treats people, what to log", 0),
     ("Professional codes (ACM, IEEE) put duty to users & public above employer convenience", 0),
     ("Mini case: you CAN scrape 100,000 users' public posts for a project — SHOULD you?", 1),
     ("Misconception: 'the tool is neutral.' What you build & how an algorithm decides encodes values + bias", 2)],
    notes="Run the scraping case through the S4 5-step process quickly: facts (is 'public' = 'consented'?), stakeholders (the 100,000), lens (rights: did they agree to THIS use?), act, reflect. Show 'can' and 'should' give different answers. A biased-data hiring algorithm discriminates automatically, at scale.")
add_activity("S5 — 'Can vs Should audit'  ·  ~6 min",
    [("In pairs (2 min), then 3 answers aloud (4 min)", 0),
     ("Each pair picks ONE app they use daily", 1),
     ("Identify (a) one PAPA theme it raises a concern about", 1),
     ("and (b) one 'can but maybe shouldn't' thing it does with their data", 1),
     ("Closes the unit on THEIR real digital life", 1)],
    notes="Makes the abstract PAPA framework personal. Good closer — students leave seeing their own apps through the unit's lens.")
add_quiz("S5 — Quick Check  ·  ~5 min",
    [("Q1. What most makes IT ethics distinct?", "q"),
     ("✅ b) The scale, speed & permanence of impact", "a"),
     ("Q2. 'PAPA' = Privacy, Accuracy, Property, and…", "q"),
     ("✅ c) Access", "a"),
     ("Discussion: an IT product you use that raises a concern — which PAPA theme, and why?", "o")],
    notes="APPLICATION [~3 min]: you'll build/manage/buy systems and decide what data to collect and how algorithms treat people; every later unit is a tool for 'can vs should'. SUMMARY [~2 min]: IT amplifies impact (scale/speed/anonymity/permanence); tech is value-laden; PAPA; professionals own can-vs-should. Next unit: ethics for IT workers & users.")

# ---------------- End-of-unit ----------------
add_divider("Assessment", "Unit 1 — End-of-Unit Quiz",
            "12 MCQ + 5 short answer + 2 applied cases + 1 discussion. Full questions and answer key are in Unit1_material.md.")
add_content("Exam pointers — what to master", "Unit 1 · Review",
    [("Ethics vs morals vs law vs etiquette; 'legal ≠ ethical'", 0),
     ("CSR & Carroll's pyramid (4 layers in order); stakeholder theory", 0),
     ("Components of an effective ethics program; tone at the top", 0),
     ("The 5-step decision process and the four ethical lenses", 0),
     ("Why IT ethics is distinct; the PAPA framework; 'can vs should'", 0)],
    notes="These five lines map to the most common exam questions. Full quiz with answer key is in the markdown file.")
add_title("End of Unit 1",
          "Next: Unit 2 — Ethics for IT Workers and IT Users",
          "Generated for IT 246 · diagrams are native shapes (editable) · speaker notes in each slide's Notes pane")

out = "/Users/inventechg1/Desktop/2083_SEM/sixth/IT246_Unit1.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
