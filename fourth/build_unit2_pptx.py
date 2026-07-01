#!/usr/bin/env python3
"""IT220 Unit 2 deck — S5–S8 (ER modelling), rebuilt to the COURSE_MATERIAL_STANDARD.

Reference implementation of the self-contained / PDF-safe style:
- Every concept is a TWO-SLIDE PAIR:
    (A) "Understand it" — definition + how-it-works + labelled diagram + memory hook.
    (B) "Apply it"      — 🌍 real example / ⚠️ common trap / 🎯 exam-ready answer / 🔑 key terms,
                          as colored panels ON THE SLIDE FACE (survives PDF export).
- Speaker notes hold ONLY timing/delivery cues — nothing load-bearing lives in notes.
- Content is distilled from Unit2_material.md (the exhaustive source of truth).

Run: python3 build_unit2_pptx.py  ->  IT220_Unit2.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "images")

# ---- palette (canonical — see COURSE_MATERIAL_STANDARD.md §5) ----
NAVY=RGBColor(0x0C,0x2B,0x4A); BLUE=RGBColor(0x18,0x5F,0xA5); TEAL=RGBColor(0x0F,0x6E,0x56)
AMBER=RGBColor(0x85,0x4F,0x0B); CORAL=RGBColor(0x99,0x3C,0x1D); LIGHT=RGBColor(0xF2,0xF5,0xF8)
PALET=RGBColor(0xE1,0xF5,0xEE); WHITE=RGBColor(0xFF,0xFF,0xFF); DARK=RGBColor(0x1A,0x1A,0x1A)
GREY=RGBColor(0x55,0x60,0x6B)
# light tints for panel cards
BLUE_T=RGBColor(0xE6,0xF1,0xFB); CORAL_T=RGBColor(0xFA,0xEC,0xE7)
TEAL_T=RGBColor(0xE1,0xF5,0xEE); AMBER_T=RGBColor(0xFA,0xEE,0xDA)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

# ================= low-level helpers =================
def _bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def _box(s,l,t,w,h): return s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)).text_frame
def _bar(s,c,l,t,w,h):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=c; sp.line.fill.background(); sp.shadow.inherit=False
def _notes(s,t):
    if t: s.notes_slide.notes_text_frame.text=t

def _header(s,kicker,header,kc=BLUE):
    _bar(s,kc,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=kc
    th=_box(s,0.7,0.66,12,1.0); th.word_wrap=True; r=th.paragraphs[0].add_run(); r.text=header
    r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY

def _card(s,l,t,w,h,accent,fill,label,body,body_sz=14):
    """A colored panel with a bold label + body text — the PDF-safe study unit."""
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    sp.line.color.rgb=accent; sp.line.width=Pt(1); sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP
    tf.margin_left=Inches(0.22); tf.margin_right=Inches(0.20)
    tf.margin_top=Inches(0.10); tf.margin_bottom=Inches(0.08)
    p=tf.paragraphs[0]; r=p.add_run(); r.text=label
    r.font.size=Pt(12.5); r.font.bold=True; r.font.color.rgb=accent
    p2=tf.add_paragraph(); p2.space_before=Pt(3); r2=p2.add_run(); r2.text=body
    r2.font.size=Pt(body_sz); r2.font.color.rgb=DARK
    return sp

# ================= slide builders =================
def add_title(title,subtitle,footer):
    s=prs.slides.add_slide(BLANK); _bg(s,NAVY); _bar(s,AMBER,0.9,3.05,1.6,0.10)
    tf=_box(s,0.9,2.0,11.5,1.0); tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(34); r.font.bold=True; r.font.color.rgb=WHITE
    t2=_box(s,0.9,3.25,11.5,1.2); t2.word_wrap=True
    r=t2.paragraphs[0].add_run(); r.text=subtitle; r.font.size=Pt(18); r.font.color.rgb=RGBColor(0xCF,0xDD,0xEA)
    t3=_box(s,0.9,6.5,11.5,0.7); t3.word_wrap=True; r=t3.paragraphs[0].add_run(); r.text=footer
    r.font.size=Pt(13); r.font.color.rgb=RGBColor(0x9A,0xB2,0xC6); return s

def add_outcomes(header,kicker,intro,items,nxt):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,header)
    tf=_box(s,0.7,1.75,12.0,4.6); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=intro
    r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=NAVY; p.space_after=Pt(8)
    for it in items:
        p=tf.add_paragraph(); r=p.add_run(); r.text="•  "+it
        r.font.size=Pt(14.5); r.font.color.rgb=DARK; p.space_after=Pt(5)
    _card(s,0.7,6.35,12.0,0.85,BLUE,BLUE_T,"➡️ Coming after this batch",nxt,body_sz=13)
    return s

def add_divider(kicker,title,hook,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,BLUE)
    tk=_box(s,0.9,1.9,11.5,0.6); r=tk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=RGBColor(0xBF,0xDC,0xF2)
    tt=_box(s,0.9,2.5,11.5,1.4); tt.word_wrap=True; r=tt.paragraphs[0].add_run(); r.text=title
    r.font.size=Pt(31); r.font.bold=True; r.font.color.rgb=WHITE
    # HOOK on the slide face (readable in PDF), not hidden in notes
    _bar(s,AMBER,0.9,4.05,1.4,0.09)
    th=_box(s,0.9,4.25,11.5,2.4); th.word_wrap=True
    p=th.paragraphs[0]; r=p.add_run(); r.text="🎣  "+hook
    r.font.size=Pt(18); r.font.color.rgb=RGBColor(0xE8,0xF1,0xF9); r.font.italic=True
    _notes(s,timing)
    return s

def concept_understand(kicker,heading,definition,mechanism,image,hook,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker+"  ·  understand",heading)
    tf=_box(s,0.7,1.72,6.2,5.05); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=definition
    r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=NAVY; p.space_after=Pt(10)
    p=tf.add_paragraph(); r=p.add_run(); r.text="HOW IT WORKS"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=BLUE; p.space_after=Pt(5)
    for m in mechanism:
        p=tf.add_paragraph(); r=p.add_run(); r.text="•  "+m
        r.font.size=Pt(14); r.font.color.rgb=DARK; p.space_after=Pt(6)
    if image:
        path=os.path.join(IMG,image)
        if os.path.exists(path):
            pic=s.shapes.add_picture(path,Inches(6.95),Inches(1.9),width=Inches(6.05))
            band_top,band_h=1.9,4.55; h_in=pic.height/914400.0
            pic.top=Inches(band_top+max(0.0,(band_h-h_in)/2.0))
    _bar(s,TEAL,0,6.9,SW.inches,0.6)
    ft=_box(s,0.7,6.96,11.9,0.5); p=ft.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text="🧠  Memory hook:  "+hook; r.font.size=Pt(13.5); r.font.bold=True; r.font.color.rgb=WHITE
    _notes(s,timing)
    return s

def concept_apply(kicker,heading,example,trap,exam,keyterms):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker+"  ·  apply",heading)
    _card(s,0.7,1.66,12.0,1.48,BLUE ,BLUE_T ,"🌍  Real example",example)
    _card(s,0.7,3.26,12.0,1.30,CORAL,CORAL_T,"⚠️  Common trap — and the fix",trap)
    _card(s,0.7,4.68,12.0,1.92,TEAL ,TEAL_T ,"🎯  Exam-ready answer",exam)
    _card(s,0.7,6.72,12.0,0.60,AMBER,AMBER_T,"🔑  Key terms",keyterms,body_sz=11.5)
    return s

def add_activity(title,steps,expected,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,PALET); _bar(s,TEAL,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text="🛠  IN-CLASS ACTIVITY  ·  THINK–PAIR–SHARE"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=TEAL
    th=_box(s,0.7,0.66,12,0.9); r=th.paragraphs[0].add_run(); r.text=title
    r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
    tf=_box(s,0.7,1.75,11.9,3.6); tf.word_wrap=True; first=True
    for text in steps:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text="•  "+text; r.font.size=Pt(15.5); r.font.color.rgb=DARK; p.space_after=Pt(7)
    _card(s,0.7,5.5,11.9,1.7,TEAL,WHITE,"✅  Expected answer",expected)
    _notes(s,timing)
    return s

def add_quiz(title,items,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,LIGHT); _bar(s,AMBER,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text="CHECK FOR UNDERSTANDING"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=AMBER
    th=_box(s,0.7,0.66,12,0.9); r=th.paragraphs[0].add_run(); r.text=title
    r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
    tf=_box(s,0.7,1.75,11.9,5.3); tf.word_wrap=True; first=True
    for text,kind in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=text
        if kind=="q": r.font.size=Pt(17); r.font.bold=True; r.font.color.rgb=NAVY; p.space_before=Pt(9)
        elif kind=="a": r.font.size=Pt(15.5); r.font.color.rgb=TEAL; r.font.bold=True; p.level=1
        else: r.font.size=Pt(15); r.font.color.rgb=DARK; r.font.italic=True; p.level=1; p.space_before=Pt(6)
    _notes(s,timing)
    return s

def add_summary(kicker,takeaways,why,nxt,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,"Summary & Key Takeaways",kc=NAVY)
    tf=_box(s,0.7,1.7,12.0,2.7); tf.word_wrap=True; first=True
    for i,tk in enumerate(takeaways,1):
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=f"{i}.  "+tk; r.font.size=Pt(15); r.font.color.rgb=DARK; p.space_after=Pt(8)
    _card(s,0.7,4.55,12.0,1.35,TEAL,TEAL_T,"💡  Why this matters",why,body_sz=13.5)
    _card(s,0.7,6.02,12.0,1.15,BLUE,BLUE_T,"➡️  Next session",nxt,body_sz=13.5)
    _notes(s,timing)
    return s

def add_roadmap(kicker,title,done,todo):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    _card(s,0.7,1.8,5.85,4.55,TEAL,TEAL_T,"✅  THIS BATCH  ·  S5–S8","\n".join(done),body_sz=14)
    _card(s,6.78,1.8,5.85,4.55,GREY,LIGHT,"⏭️  COMING NEXT  ·  S9–S12","\n".join(todo),body_sz=14)
    return s

def add_solved_problem(kicker,title,scenario,steps,answer,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title,kc=AMBER)
    _card(s,0.7,1.58,12.0,1.0,AMBER,AMBER_T,"📋  Problem",scenario,body_sz=13)
    tf=_box(s,0.85,2.72,11.7,2.95); tf.word_wrap=True; first=True
    for i,st in enumerate(steps,1):
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=f"Step {i}.  "; r.font.size=Pt(12.5); r.font.bold=True; r.font.color.rgb=AMBER
        r2=p.add_run(); r2.text=st; r2.font.size=Pt(12.5); r2.font.color.rgb=DARK; p.space_after=Pt(4)
    _card(s,0.7,5.72,12.0,1.5,TEAL,TEAL_T,"✅  Worked answer",answer,body_sz=13)
    _notes(s,timing)
    return s

def add_cheatsheet(kicker,title,blocks):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    xs=[0.7,6.78]; ys=[1.72,3.36,5.0]; hh=1.55
    for idx,(label,body) in enumerate(blocks[:6]):
        x=xs[idx%2]; y=ys[idx//2]
        _card(s,x,y,5.85,hh,BLUE,BLUE_T,label,body,body_sz=12)
    return s

def add_glossary(kicker,title,terms):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    half=(len(terms)+1)//2
    for col,(x,group) in enumerate([(0.7,terms[:half]),(6.78,terms[half:])]):
        tf=_box(s,x,1.75,5.85,5.3); tf.word_wrap=True; first=True
        for term,defn in group:
            p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            r=p.add_run(); r.text=term+" — "; r.font.size=Pt(11.5); r.font.bold=True; r.font.color.rgb=NAVY
            r2=p.add_run(); r2.text=defn; r2.font.size=Pt(11.5); r2.font.color.rgb=DARK
            p.space_after=Pt(6)
    return s

def _add_page_numbers():
    """Final pass: number every slide, auto-contrasting with its background."""
    total=len(prs.slides._sldIdLst)
    for i,s in enumerate(prs.slides,1):
        try:
            h=str(s.background.fill.fore_color.rgb)
            r,g,b=int(h[0:2],16),int(h[2:4],16),int(h[4:6],16)
            dark=(0.299*r+0.587*g+0.114*b)<128
        except Exception:
            dark=False
        col=RGBColor(0xE8,0xF0,0xF7) if dark else GREY
        tb=_box(s,12.35,0.30,0.9,0.35); p=tb.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
        r=p.add_run(); r.text=f"{i} / {total}"; r.font.size=Pt(10); r.font.color.rgb=col

# ============================================================
#                        BUILD
# ============================================================
add_title("Unit 2 — Data Modelling (ER Model)",
          "IT 220: Database Management System  ·  BIM 4th Semester  ·  Chen notation",
          "Sessions S5–S8 (of S5–S12)  ·  50 min each  ·  Nepal / South Asia context  ·  "
          "Self-contained slides: every concept is complete on the slide face — no teacher or speaker-notes required to learn it.")

add_outcomes("Unit 2 — Learning Outcomes","overview  ·  s5–s8 batch",
    "By the end of this batch (S5–S8), you will be able to:",
    ["Explain the conceptual (ER) model and where it sits in the database design process",
     "Model entity types, entity sets, and all four attribute types, and choose keys",
     "Model relationships: type/set, degree, roles, cardinality ratio, and participation",
     "Model weak entity types with identifying relationships and partial keys"],
    "S9–S12: higher-degree (ternary) relationships, specialization/generalization, ER→table mapping, and the consolidated end-of-unit quiz.")

add_roadmap("Unit 2 — Roadmap","Where each session fits (S5–S12)",
    ["S5   Conceptual models & the design process",
     "S6   Entity types/sets, attributes, keys",
     "S7   Relationships, roles, cardinality, participation",
     "S8   Weak entities, identifying relationship, partial key"],
    ["S9    Full ER diagrams · naming · design pitfalls",
     "S10  Higher-degree (ternary) relationships",
     "S11  Specialization & generalization",
     "S12  Mapping ER → relational tables",
     "+     Consolidated end-of-unit quiz"])

# ============================ S5 ============================
add_divider("Session 5 · Lecture hour 5 (of 8)","Conceptual Data Models & the Design Process",
    "Before Daraz wrote a single line of code, someone drew boxes and lines — Product, Customer, Order — on a whiteboard. Why draw before you build? Because conceptual design comes first.",
    "OPENING HOOK [~5 min]. Don't teach ER symbols yet — just plant 'describe WHAT before HOW'. Agenda: conceptual model → design pipeline → data vs functional requirements.")

concept_understand("S5 · Concept 1 · [THEORY]","What a Conceptual Data Model Is (and Why ER)",
    "A high-level, implementation-independent description of WHAT data a system holds and how the pieces relate — expressed as real-world things, not tables, files, or a specific DBMS.",
    ["The Entity-Relationship (ER) model is the standard notation. Chen: rectangle = entity, oval = attribute, diamond = relationship.",
     "Technology-neutral and readable by non-technical stakeholders — a registrar or finance manager can verify it.",
     "It is a communication + thinking tool: agree on WHAT before committing to costly HOW.",
     "It sits ABOVE tables; turning it into tables is a separate, later step (S12)."],
    "s5_what_vs_how.png","Model the WHAT; the HOW comes later.",
    "~12 min. Airline example (Passenger/Flight/Booking) BEFORE naming Oracle vs MySQL.")
concept_apply("S5 · Concept 1 · [THEORY]","What a Conceptual Data Model Is (and Why ER)",
    "Daraz's whiteboard: boxes CUSTOMER, PRODUCT, ORDER with lines between them — a conceptual model with no mention of MySQL or servers. A startup that skips it and types CREATE TABLE on day one discovers 'a customer has many addresses' too late, and must rewrite live tables.",
    "\"Modelling = drawing the tables.\" No — the conceptual model is ABOVE tables; rows and foreign keys come later (S12). Concept first, tables later.",
    "A conceptual data model is a high-level, implementation-independent description of WHAT data a system stores and how the data relate — in terms of real-world things, not tables or storage. The ER model is the most widely used conceptual notation because it is technology-neutral and easy for non-technical stakeholders (rectangle = entity, oval = attribute, diamond = relationship).",
    "conceptual / high-level model  ·  implementation-independent  ·  Entity-Relationship (ER) model  ·  Chen notation  ·  stakeholder-readable")

concept_understand("S5 · Concept 2 · [THEORY]","The Database Design Process",
    "Database design is a pipeline of four stages, each producing output that feeds the next: Requirements → Conceptual → Logical → Physical (RCLP).",
    ["Requirements collection & analysis — capture what data AND what operations are needed.",
     "Conceptual design — build the DBMS-independent ER diagram (the bulk of Unit 2).",
     "Logical design — MAP the ER diagram to relational tables with keys / foreign keys (this is S12).",
     "Physical design — storage, indexes, file organization for performance (the DBA's job).",
     "Mirrors Unit 1's levels: conceptual → logical → internal schema."],
    "s5_design_pipeline.png","Requirements → Conceptual → Logical → Physical (RCLP).",
    "~12 min. Draw the four-box pipeline and keep it up all session — later concepts hang off it.")
concept_apply("S5 · Concept 2 · [THEORY]","The Database Design Process",
    "eSewa's wallet, stage by stage: finance-team requirements → ER diagram (USER makes TRANSACTION) → relational tables (User, Transaction, a foreign key) → indexes so your statement loads in a tap. Teams that skip requirements discover 'top-up source' was never captured and pay in painful live-data migrations.",
    "Skipping requirements to 'save time' is the #1 cause of costly rework. Every stage you skip, a later stage pays for — with interest.",
    "Database design has four phases: (1) Requirements collection & analysis — what data and what operations; (2) Conceptual design — a DBMS-independent ER diagram; (3) Logical design — map the ER diagram to relational tables (keys, foreign keys); (4) Physical design — storage structures and indexes. Each phase feeds the next; ER sits at the conceptual phase, ER→tables is the logical phase.",
    "database design process  ·  requirements collection  ·  conceptual design (ER)  ·  logical design (mapping)  ·  physical design  ·  rework")

concept_understand("S5 · Concept 3 · [THEORY]","Data vs Functional Requirements",
    "Requirements come in two kinds, captured together: data requirements = WHAT to store (the nouns); functional requirements = WHAT operations to support (the verbs).",
    ["Data requirements → shape the ENTITIES and ATTRIBUTES of the ER diagram.",
     "Functional requirements → shape the OPERATIONS the design must serve quickly and correctly.",
     "A DB that stores the right data but can't do a needed operation is as broken as one missing data.",
     "Capture both — the nouns AND the verbs."],
    "s5_data_vs_functional.png","Data = the nouns to store; functional = the verbs to do.",
    "~11 min. Split the board: STORE (left) vs DO (right). Use Khalti.")
concept_apply("S5 · Concept 3 · [THEORY]","Data vs Functional Requirements",
    "Khalti's two lists — data: users (name, number, balance, KYC) and transactions (amount, date, status); functional: send money, top up, pay a bill, view statement. The 'view statement for a date range' verb is WHY history loads instantly. Capture only the data and you ship a correct but painfully slow statement screen.",
    "\"Requirements are only about data.\" No — miss the operations (verbs) and you can build a perfectly-structured DB that is slow or unable to do what users actually need.",
    "Data requirements specify what data must be stored — the things and their properties (nouns), e.g. Khalti users, balances, transactions. Functional requirements specify what operations the system must support (verbs), e.g. send money, view statement. Data requirements shape entities/attributes; functional requirements shape the operations the design must serve. Both must be captured.",
    "data requirements (nouns → entities/attributes)  ·  functional requirements (verbs → transactions)  ·  requirements collection  ·  capture both")

add_activity("Nouns & Verbs of a Nepali app",
    ["In pairs (2–3 min): pick a Nepali app — eSewa, Khalti, Daraz, Nagarik App, or Foodmandu.",
     "Write TWO data requirements (nouns to store) and TWO functional requirements (operations to support).",
     "Then say which pipeline stage turns your nouns into an ER diagram.",
     "Share aloud (2 min); sort each answer into the STORE column or the DO column."],
    "Daraz: data = products, customers, orders; functional = search product, place order, track delivery. The stage that turns nouns into an ER diagram is CONCEPTUAL design. Students reliably list nouns and forget verbs — push every pair to name at least one operation.",
    "ACTIVITY [~5 min].")
add_quiz("S5 — Quick Check",
    [("Q1.  Which design phase produces the ER diagram?","q"),
     ("a) physical   b) ✅ conceptual   c) implementation   d) testing","a"),
     ("Q2.  A conceptual data model is independent of…","q"),
     ("a) requirements   b) ✅ the specific DBMS / storage   c) the entities   d) stakeholders","a"),
     ("Discussion: pick a Nepali app — two data requirements and two functional requirements.","o")],
    "QUIZ [~5 min].")
add_summary("S5 · Summary  ·  [~2 min]",
    ["A conceptual data model describes WHAT data you hold, not HOW it's stored; the ER model (Chen notation) is the standard, technology-neutral, stakeholder-readable notation.",
     "Design flows through four stages — Requirements → Conceptual (ER) → Logical (tables) → Physical (storage) — each feeding the next.",
     "Requirements come in two kinds: data requirements (nouns to store) and functional requirements (verbs to support). Capture both."],
    "Every data project in Nepal — fintechs, e-commerce, government IT — begins with requirements gathering and a conceptual (ER) model. Skipping to tables is the #1 cause of costly rework. \"Can you draw the ER model before you code?\" is a standard junior-developer expectation.",
    "S6 — the building blocks of the ER model: entities, entity types/sets, attributes, and keys.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S6 ============================
add_divider("Session 6 · Lecture hour 6 (of 8)","Entity Types, Sets, Attributes & Keys",
    "On the Nagarik App, YOU are a citizen and your friend is a citizen — same 'shape' (a citizenship number, a name, a DOB), different data. What's the shape, and what's the data?",
    "OPENING HOOK [~5 min]. Tie straight back to S2 schema-vs-instance: type ≈ schema (design-time), set ≈ instance (a snapshot). Agenda: entity/type/set → attribute types → keys.")

concept_understand("S6 · Concept 1 · [THEORY]","Entity, Entity Type, Entity Set",
    "Entity = one real-world thing. Entity TYPE = the category/template all such things share (design-time; a rectangle in Chen). Entity SET = all entities of that type existing right now (instance-time).",
    ["One entity type has many entity sets over time — the same pattern as 'one schema, many instances' (S2).",
     "The type is stable; the set changes constantly as things come and go.",
     "In Chen notation, an entity type is drawn as a rectangle.",
     "An entity is NOT a table row — the row is only how it is represented later (S12)."],
    "s6_type_vs_set.png","Type = the cutter; entity = one biscuit; set = the tray right now.",
    "~11 min. Draw rectangle STUDENT (type), write 'Sita' beside it (entity), dotted circle round a roster (set).")
concept_apply("S6 · Concept 1 · [THEORY]","Entity, Entity Type, Entity Set",
    "On Daraz, CUSTOMER is an entity type (the category, defined by ID, name, email). YOU, with your ID and order history, are one entity. Everyone registered at this moment — millions — is the entity set: sign-ups enlarge it, deactivations shrink it, yet the TYPE 'CUSTOMER' is unchanged.",
    "\"An entity is just a row in a table.\" That's later thinking — an entity is a real-world thing; the row is only its stored representation, produced when we map to tables in S12.",
    "An entity is a single real-world thing (e.g. the student Sita). An entity type is the category/template all such things share — a design-time definition (e.g. STUDENT, a rectangle in Chen); it is stable. An entity set is all entities of that type at a given moment (e.g. all students enrolled today); it changes constantly. One entity type has many entity sets — analogous to one schema having many instances.",
    "entity  ·  entity type (rectangle, design-time)  ·  entity set (instance-time)  ·  type ≈ schema, set ≈ instance")

concept_understand("S6 · Concept 2 · [THEORY]","Attributes and Their Four Types",
    "An attribute is a property that describes an entity (an oval in Chen notation). Attributes come in four types — simple, composite, derived, multivalued — and classifying them correctly now decides table design later.",
    ["Simple (atomic): one indivisible value, e.g. roll_no — a plain oval.",
     "Composite: made of sub-parts, e.g. name = first/middle/last — the oval branches into child ovals.",
     "Derived: computed, NOT stored, e.g. age from DOB — a dashed oval.",
     "Multivalued: can hold several values, e.g. phone numbers — a double oval.",
     "In S12 this decides mapping: multivalued → own table; composite → flattened columns; derived → left out."],
    "s6_attributes.png","Simple = one; Composite = parts; Derived = computed; Multivalued = many.",
    "~13 min. Draw each: plain / branching / dashed / double oval. Hit the 'age' misconception hard.")
concept_apply("S6 · Concept 2 · [THEORY]","Attributes and Their Four Types",
    "Daraz's saved delivery addresses = a MULTIVALUED attribute (many per customer) that is also COMPOSITE (district/municipality/ward) — drawn as a double oval that branches, and mapped to its own table later. Daraz stores your date of birth, never your age — age is DERIVED and would rot every birthday.",
    "\"Store the age\" / \"cram several phone numbers into one box.\" Age is derived (causes update anomalies); several values must be a multivalued attribute → its own table. Store what you can't compute; derive what you can.",
    "Simple (atomic): a single indivisible value (roll_no). Composite: made of sub-parts (name = first+middle+last). Derived: computed, not stored (age from DOB — dashed oval). Multivalued: several values for one entity (phone_numbers — double oval). In Chen, attributes are ovals; composites branch, derived are dashed, multivalued are doubled. Derived and multivalued especially affect how the entity maps to tables.",
    "attribute (oval)  ·  simple/atomic  ·  composite (branches)  ·  derived (dashed, computed)  ·  multivalued (double)  ·  update anomaly")

concept_understand("S6 · Concept 3 · [THEORY]","Keys",
    "A key is a uniqueness constraint — an attribute (or set of attributes) whose value is unique for every entity, used to tell one entity from another. In Chen notation, a key attribute is underlined.",
    ["Candidate key — any minimal attribute set that could uniquely identify (an entity type may have several).",
     "Primary key — the one candidate key the designer chooses as the main identifier (the underlined one).",
     "Composite key — a key made of more than one attribute together (e.g. roll_no + program).",
     "A good key is unique, stable, and minimal."],
    "s6_keys.png","A key = the roll number that no two share.",
    "~11 min. 'If two students are both Ram Thapa, how does the system tell them apart?' — that's a key.")
concept_apply("S6 · Concept 3 · [THEORY]","Keys",
    "The Nagarik App keys each citizen on their citizenship / national ID, not their name — names repeat across Nepal, but that ID is designed to be unique. It is the primary key; a phone number may be another candidate key. Identify citizens by name and two 'Sita Sharma's merge their land, tax, and benefit records.",
    "\"A name can be a key.\" Names repeat — two 'Ram Thapa's will collide. Test any proposed key: could two different real things ever share this value? If yes, it is not a key.",
    "A key is a uniqueness constraint that uniquely identifies each entity. A candidate key is any minimal attribute set that could serve as identifier (there may be several). The primary key is the chosen main identifier (underlined in Chen). A composite key combines more than one attribute (roll_no + program) when no single attribute is unique alone. Names should not be keys because they repeat.",
    "key / uniqueness constraint  ·  candidate key  ·  primary key (underlined)  ·  composite key  ·  test: unique, stable, minimal")

add_activity("Label the attributes of a Khalti account",
    ["In pairs (3 min): list the attributes of a Khalti user account.",
     "Label each as simple / composite / derived / multivalued.",
     "Then pick the PRIMARY key and name one other CANDIDATE key.",
     "Share (2 min); draw the winning account as a Chen fragment."],
    "user_id (simple, primary key) · mobile_number (candidate key) · full_name (composite) · age (derived — don't store!) · linked_bank_accounts or saved_numbers (multivalued). Watch two reliable slips: storing 'age', and treating several saved numbers as one field.",
    "ACTIVITY [~5 min].")
add_quiz("S6 — Quick Check",
    [("Q1.  'Age' computed from date of birth is a ___ attribute:","q"),
     ("a) stored   b) multivalued   c) ✅ derived   d) composite","a"),
     ("Q2.  All STUDENT entities existing at one moment form the:","q"),
     ("a) entity type   b) ✅ entity set   c) attribute   d) relationship","a"),
     ("Discussion: list the attributes of a Khalti account and label each simple / composite / derived / multivalued.","o")],
    "QUIZ [~5 min].")
add_summary("S6 · Summary  ·  [~2 min]",
    ["Entity type = the category (design-time, a rectangle, like a schema); entity set = all its members right now (instance-time); entity = one real-world thing.",
     "Attributes come in four types — simple, composite, derived, multivalued (ovals in Chen) — and you never store what you can derive.",
     "A key uniquely identifies an entity; the primary key is the chosen identifier (underlined); composite keys combine attributes. Names are not keys."],
    "Correctly classifying attributes — spotting multivalued ones (addresses, phones) and derived ones (age, totals) — decides table design and prevents duplicated, out-of-sync data. Choosing a genuinely unique key (national ID, roll number, customer ID — never a name) is a decision juniors are trusted to get right.",
    "S7 — how entities CONNECT: relationship types and sets, roles, and the structural constraints (cardinality and participation).",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S7 ============================
add_divider("Session 7 · Lecture hour 7 (of 8)","Relationships, Roles & Structural Constraints",
    "A student ENROLLS IN a course; a customer PLACES an order. Those verbs — not the nouns — are where most real modelling decisions, and most mistakes, happen.",
    "OPENING HOOK [~5 min]. Densest session (4 concepts) — keep each tight, lean on the diagrams. Agenda: relationship type/set/degree → roles & recursive → cardinality → participation.")

concept_understand("S7 · Concept 1 · [THEORY]","Relationship Type, Set & Degree",
    "A relationship is an association among entities (a diamond in Chen). Like entities, it has a TYPE (named, design-time), a SET (the current pairings), and a DEGREE (how many entity types take part).",
    ["Relationship type: a named association between entity types, e.g. ENROLLS between STUDENT & COURSE.",
     "Relationship set: the actual pairings right now — (Sita, DBMS), (Hari, Networking).",
     "Degree: number of participating entity types — 2 = binary (most common), 3 = ternary, n = n-ary.",
     "Type is to set as schema is to instance (again)."],
    "s7_relationship.png","Type = the verb; set = today's pairings; degree = how many entities the verb joins.",
    "~9 min. Draw STUDENT ▭ — ◇ENROLLS — ▭ COURSE. Degree = rectangles the diamond touches.")
concept_apply("S7 · Concept 1 · [THEORY]","Relationship Type, Set & Degree",
    "On Daraz, PLACES connects CUSTOMER and ORDER. Its set right now is the ever-changing list of who-placed-which-order this minute; its degree is 2 (binary). The relationship is NOT the customer_id column — that column is only how it's implemented after mapping. 'A customer places an order' exists first, as a diamond.",
    "\"A relationship is just a foreign key.\" At the conceptual level it's a real-world association; a foreign key is one later implementation — and M:N relationships become a whole junction table, not a foreign key.",
    "A relationship type is a named association between entity types (ENROLLS between STUDENT and COURSE; a diamond in Chen). A relationship set is the actual associations at a moment ((Sita, DBMS), (Hari, Networking)). Degree is the number of participating entity types — 2 is binary (most common), 3 ternary, n n-ary. A relationship is a conceptual association, not a foreign key.",
    "relationship (diamond)  ·  relationship type  ·  relationship set  ·  degree (binary / ternary / n-ary)  ·  not a foreign key")

concept_understand("S7 · Concept 2 · [THEORY]","Roles & Recursive (Self) Relationships",
    "A role name is the part an entity plays in a relationship. A recursive (self) relationship is one where the SAME entity type participates more than once — the diamond loops back to the same rectangle.",
    ["In ordinary relationships between two different entity types, roles are implicit and optional.",
     "In recursive relationships, role names are MANDATORY — both ends are the same type and would otherwise be indistinguishable.",
     "Role names carry the direction / meaning of the association."],
    "s7_recursive.png","Same type on both ends → you must name the roles.",
    "~8 min. Draw ONE rectangle EMPLOYEE with a diamond SUPERVISES looping back; label 'supervisor' / 'supervisee'.")
concept_apply("S7 · Concept 2 · [THEORY]","Roles & Recursive (Self) Relationships",
    "A college peer-mentoring scheme: entity type STUDENT, recursive relationship MENTORS looping back to STUDENT. Label one line 'mentor', the other 'mentee' — without them the system can't tell whether Sita mentors Gita or Gita mentors Sita. (Ntc employee supervision works identically: 'supervisor' / 'supervisee'.)",
    "Leaving roles off a recursive relationship makes it ambiguous — the direction (who supervises or mentors whom) is simply lost, because both ends are the same entity type.",
    "A recursive (self) relationship is one where the same entity type participates more than once — the relationship connects an entity type to itself (EMPLOYEE SUPERVISES EMPLOYEE). A role name states the part an entity plays. Recursive relationships require role names (e.g. 'supervisor' and 'supervisee') because both ends are the same type and would otherwise be indistinguishable; the roles carry the direction of the association.",
    "role name  ·  recursive / self relationship  ·  roles are mandatory for recursive relationships (to distinguish the identical ends)")

concept_understand("S7 · Concept 3 · [THEORY]","Cardinality Ratio (1:1, 1:N, M:N)",
    "The cardinality ratio of a binary relationship specifies the MAXIMUM number of relationship instances an entity can take part in — how many of one thing relate to how many of the other.",
    ["1:1 — each side relates to at most one, e.g. CITIZEN–PASSPORT.",
     "1:N — one relates to many, each of those back to one, e.g. DEPARTMENT–EMPLOYEE (read the other way = N:1).",
     "M:N — both sides relate to many, e.g. STUDENT–COURSE; very common.",
     "Technique: read it BOTH ways — fix one side, ask 'how many of the other?', then repeat.",
     "M:N can't be a single foreign key → it becomes a junction table in S12."],
    "s7_cardinality.png","Read the max on both sides; M:N means a junction table later.",
    "~10 min. Most exam-heavy idea. One crisp example per ratio; teach the read-it-both-ways trick.")
concept_apply("S7 · Concept 3 · [THEORY]","Cardinality Ratio (1:1, 1:N, M:N)",
    "Daraz ORDER–PRODUCT is M:N — one order has many products (phone, case, charger), one product is in many orders. So Daraz can't bolt a product_id onto the orders table; later this becomes an ORDER_ITEM junction table (one row per order-product pair). Contrast 1:N DEPARTMENT–EMPLOYEE, where a single department_id on the employee is enough.",
    "\"M:N can be stored directly.\" No — there's no single 'other side' to point to; every M:N must break out into a junction table when mapped (S12).",
    "The cardinality ratio specifies the maximum relationship instances an entity may participate in. 1:1 — at most one each side (CITIZEN–PASSPORT). 1:N — one relates to many, each back to one (DEPARTMENT–EMPLOYEE). M:N — both sides many (STUDENT–COURSE). Determine it by checking the maximum on EACH side. An M:N relationship must become a junction table when mapped.",
    "cardinality ratio  ·  1:1  ·  1:N / N:1  ·  M:N  ·  read-it-both-ways  ·  M:N ⇒ junction table (preview of S12)")

concept_understand("S7 · Concept 4 · [THEORY]","Participation (Total vs Partial)",
    "The participation constraint tells you whether an entity MUST take part at all (the minimum) — the second structural constraint. Total = mandatory (double line); partial = optional (single line).",
    ["Total participation: every entity of the type must participate — drawn as a double line.",
     "Partial participation: an entity may or may not participate — drawn as a single line.",
     "The two sides of one relationship can differ.",
     "Cardinality = 'how many' (max); participation = 'whether at all' (min) — they are independent.",
     "Drives nullable vs mandatory foreign keys later."],
    "s7_participation.png","Total = double line = compulsory; partial = single line = optional.",
    "~8 min. Bank: every LOAN belongs to a CUSTOMER (total); not every CUSTOMER has a loan (partial).")
concept_apply("S7 · Concept 4 · [THEORY]","Participation (Total vs Partial)",
    "At NIC Asia, model CUSTOMER —HAS— LOAN. Every loan MUST belong to a customer (no ownerless loans) → the loan side is total (double line). But a customer may have only savings and no loan → the customer side is partial (single line). Live, the DB refuses a loan with no customer but allows customers with zero loans.",
    "Making both sides partial by mistake allows orphaned loans belonging to nobody — a genuine data-integrity and audit disaster for a bank.",
    "The participation constraint states whether an entity's participation is mandatory. Total participation — every entity of the type must participate (double line); e.g. every LOAN must belong to a CUSTOMER. Partial participation — an entity may or may not participate (single line); e.g. a CUSTOMER need not have any loan. The two sides can differ. Participation gives the minimum; cardinality gives the maximum.",
    "participation constraint  ·  total (mandatory, double line)  ·  partial (optional, single line)  ·  structural constraints: cardinality = max, participation = min")

add_activity("Cardinality & participation on eSewa",
    ["In pairs (3–4 min): on eSewa, model USER and TRANSACTION with a relationship.",
     "Decide the CARDINALITY ratio by reading it both ways.",
     "Decide the PARTICIPATION on each side, and draw the fragment with the right lines.",
     "Share (2 min); compare fragments on the board."],
    "Cardinality 1:N (one user, many transactions; each transaction one user). Participation — TRANSACTION side total (every transaction belongs to a user, double line); USER side partial (a brand-new user may have zero transactions, single line). Two skills drilled: read-it-both-ways, and can-it-exist-without.",
    "ACTIVITY [~6 min].")
add_quiz("S7 — Quick Check",
    [("Q1.  'One department has many employees, each employee in one department' is:","q"),
     ("a) 1:1   b) ✅ 1:N   c) M:N   d) recursive","a"),
     ("Q2.  A double line from an entity to a relationship diamond means:","q"),
     ("a) many-to-many   b) a weak entity   c) ✅ total participation   d) a recursive role","a"),
     ("Discussion: on eSewa, model USER–TRANSACTION — its cardinality ratio and participation on each side.","o")],
    "QUIZ [~5 min].")
add_summary("S7 · Summary  ·  [~2 min]",
    ["A relationship is an association among entities (a diamond in Chen); it has a type, a set, and a degree (binary = 2 is normal).",
     "Roles name the parts entities play and are mandatory in recursive relationships (same entity type on both ends).",
     "Two structural constraints: cardinality ratio (1:1, 1:N, M:N — how many) and participation (total/double line vs partial/single line — whether at all)."],
    "Cardinality and participation directly drive the schema: cardinality decides foreign key (1:N) vs junction table (M:N); participation decides whether that foreign key can be empty or must always be filled. Getting them right is what separates a working database from one that quietly corrupts.",
    "S8 — entities that can't stand alone: weak entity types, their identifying relationships, and partial keys.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S8 ============================
add_divider("Session 8 · Lecture hour 8 (of 8)","Weak Entity Types",
    "An 'order item' on Daraz means nothing on its own — item #2 of WHICH order? Delete the order and its line items simply cease to make sense. Some entities can't exist alone.",
    "OPENING HOOK [~5 min]. Pre-empt the one trap all session: 'weak' does NOT mean 'unimportant' — it means 'identity-dependent'. Agenda: weak vs strong → identifying relationship → partial key.")

concept_understand("S8 · Concept 1 · [THEORY]","Weak vs Strong (Regular) Entity Types",
    "A strong (regular) entity type has its own key and can be identified independently (single rectangle). A weak entity type has NO key of its own and is identified only through an owner entity (double rectangle).",
    ["Defining property: identity dependence — a weak entity's own attributes are unique only WITHIN one owner, not across the whole set.",
     "Strong examples: EMPLOYEE (employee ID), STUDENT (roll number).",
     "Weak example: DEPENDENT — two employees can each have a child 'Ram', so 'Ram' identifies nobody alone.",
     "This is common and important, not an edge case."],
    "s8_weak_vs_strong.png","Weak = can't identify itself — it borrows identity from an owner.",
    "~12 min. Draw a single rectangle (strong) next to a double rectangle (weak). Kill 'weak = unimportant' immediately.")
concept_apply("S8 · Concept 1 · [THEORY]","Weak vs Strong (Regular) Entity Types",
    "ORDER_ITEM on Daraz: 'item 2' means nothing alone — item 2 of WHICH order? It's identified only as 'line 2 of order #88921' → a weak entity (double rectangle) owned by ORDER. Delete the order and its items lose all meaning. (A bank's INSTALLMENT of a LOAN is identical.)",
    "\"Weak = unimportant / optional.\" Wrong — order items and loan installments are critical to the business. 'Weak' means IDENTITY-DEPENDENT (no key of its own), never 'doesn't matter'.",
    "A strong (regular) entity type has its own key and is identified independently (single rectangle; e.g. EMPLOYEE). A weak entity type has no key of its own and is identified only through an owning (identifying) entity (double rectangle; e.g. DEPENDENT, ORDER_ITEM). Its attributes are unique only within one owner. 'Weak' means identity-dependent, not unimportant.",
    "strong entity (own key, single rectangle)  ·  weak entity (no key, double rectangle)  ·  owner / identifying entity  ·  identity dependence")

concept_understand("S8 · Concept 2 · [THEORY]","The Identifying (Owner) Relationship",
    "The identifying (owner) relationship links a weak entity to its owner, supplying the identity the weak entity lacks. In Chen it is a DOUBLE diamond, and the weak side always has total participation (double line).",
    ["The double diamond distinguishes it from ordinary relationships (single diamond).",
     "The weak entity cannot exist without its owner → always total participation on the weak side.",
     "The 'three doubles' signature: double rectangle + double diamond + double line — examiners look for all three."],
    "s8_identifying.png","The double diamond is the owner link — no owner, no weak entity.",
    "~11 min. Draw EMPLOYEE ══ ◈HAS_DEPENDENT ══ ▨DEPENDENT. Give LOAN—has—INSTALLMENT.")
concept_apply("S8 · Concept 2 · [THEORY]","The Identifying (Owner) Relationship",
    "A car loan: LOAN (strong, own loan number) —HAS→ INSTALLMENT (weak). The identifying relationship is a double diamond; the installment side is total (double line). Live, the bank literally cannot record an installment not tied to a loan — 'installment 3' has no identity until you say WHICH loan's installment 3.",
    "Forgetting the owner link (or drawing a single diamond) breaks identity — a weak entity with no identifying relationship cannot be identified at all.",
    "An identifying (owner) relationship connects a weak entity to its owner (identifying) entity, supplying the identity the weak entity lacks. In Chen notation it is drawn as a double diamond. The weak entity's participation is always total (double line) — a weak entity cannot exist without its owner. Example: LOAN —HAS→ INSTALLMENT.",
    "identifying (owner) relationship (double diamond)  ·  owner entity  ·  total participation on the weak side  ·  the 'three doubles' signature")

concept_understand("S8 · Concept 3 · [THEORY]","The Partial Key (Discriminator)",
    "A partial key (discriminator) is an attribute that distinguishes weak entities belonging to the SAME owner. It is NOT unique on its own — only within one owner. Shown with a dashed underline in Chen.",
    ["The full identifier of a weak entity = owner's primary key + the weak entity's partial key.",
     "This is exactly the S6 composite-key idea, applied to weak entities.",
     "The dashed underline contrasts with the solid underline of a strong entity's primary key."],
    "s8_partial_key.png","Partial key alone repeats; owner's key + partial key = the full identity.",
    "~7 min. 'If installment #3 exists for thousands of loans, how do we pick the right one?' Kill 'partial key alone is unique'.")
concept_apply("S8 · Concept 3 · [THEORY]","The Partial Key (Discriminator)",
    "Hotel room '101' exists in hundreds of hotels — the room number is a partial key, unique only WITHIN one hotel. Full identity = hotel ID + room number (Everest Hotel + 101). Key rooms by '101' alone and a Pokhara room collides with a Kathmandu one — booking the wrong guest into the wrong city.",
    "\"The partial key alone is unique.\" No — installment 3 of loan 4521 and installment 3 of loan 6002 are different installments. Partial key + owner's key = the real identity.",
    "A partial key (discriminator) is an attribute of a weak entity that distinguishes weak entities belonging to the same owner; it is unique only within an owner, shown with a dashed underline (e.g. installment_no). A weak entity is uniquely identified by the owner's primary key + its partial key (e.g. loan_id + installment_no) — neither part suffices alone.",
    "partial key / discriminator (dashed underline)  ·  unique only within one owner  ·  full identifier = owner PK + partial key  ·  links to the S6 composite key")

add_activity("Find the weak entity in a Nepali app",
    ["In pairs (3 min): find a weak entity — a comment on a Facebook post, a room in a hotel booking, an item on a Daraz order, an installment on a loan.",
     "Name (a) its OWNER, (b) its PARTIAL KEY, and (c) the FULL identifier (owner key + partial key).",
     "Draw the fragment: double rectangle + double diamond + double line + dashed-underlined partial key.",
     "Share (2 min); check each pair has all 'three doubles' plus the dashed underline."],
    "e.g. COMMENT owned by POST, partial key = comment sequence/timestamp, full id = post_id + that. Verify two most-tested points: participation on the weak side is TOTAL, and the full identifier INCLUDES the owner's key (not the partial key alone).",
    "ACTIVITY [~5 min].")
add_quiz("S8 — Quick Check",
    [("Q1.  A weak entity is one that:","q"),
     ("a) is unimportant   b) ✅ has no key of its own   c) has no attributes   d) is always 1:1","a"),
     ("Q2.  The full identifier of a weak entity is:","q"),
     ("a) its partial key only   b) any attribute   c) ✅ the owner's primary key + its partial key   d) a derived attribute","a"),
     ("Discussion: find a weak entity in a Nepali app and name its owner and its partial key.","o")],
    "QUIZ [~5 min].")
add_summary("S8 · Summary  ·  [~2 min]",
    ["A weak entity has no key of its own (double rectangle) and depends on an owner for identity — 'weak' means identity-dependent, not unimportant.",
     "An identifying relationship (double diamond) ties the weak entity to its owner, and the weak side always has total participation (double line).",
     "A partial key (dashed underline) distinguishes weak entities of the same owner; full identity = owner's primary key + partial key."],
    "Weak entities are everywhere — invoice line-items, loan installments, exam answer-sheets, comments, hotel rooms. Modelling them wrong (treating a partial key as unique, or forgetting the owner link) corrupts integrity exactly where a business can least afford it: billing, accounting, exams.",
    "S9 — putting it together into clean, correctly-named ER diagrams, and avoiding common design pitfalls.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============= CAPSTONE SOLVED PROBLEM (synthesises S5–S8) =============
add_divider("Capstone · S5–S8","Put it all together: model a bank loan system",
    "One scenario, every idea from this batch — entities, attributes, keys, relationships, cardinality, participation, and a weak entity. Work it step by step before you look at the answer.",
    "Pause here and let students attempt it in pairs before revealing the worked answer.")
add_solved_problem("Capstone · Solved Problem","ER model for a bank loan system",
    "A bank stores CUSTOMERS (each with a citizenship number, name, phone numbers) and the LOANS they take. "
    "Each loan is repaid in monthly INSTALLMENTS numbered 1, 2, 3… Model this in Chen notation.",
    ["Entities: CUSTOMER and LOAN are strong (own keys); INSTALLMENT has no key of its own → WEAK (double rectangle).",
     "Attributes/keys: CUSTOMER — citizenship_number (primary key), name (composite), phone_numbers (multivalued); LOAN — loan_id (primary key).",
     "Relationships: CUSTOMER—HAS—LOAN (single diamond); LOAN—HAS—INSTALLMENT is the identifying relationship (double diamond).",
     "Cardinality: CUSTOMER–LOAN is 1:N; LOAN–INSTALLMENT is 1:N.",
     "Participation: a LOAN must have a customer → total (double line); a customer may have none → partial. Every INSTALLMENT must have a loan → total.",
     "Weak identity: partial key = installment_no (dashed underline); full identifier = loan_id + installment_no."],
    "CUSTOMER (strong) —HAS— LOAN (strong) ══HAS══ INSTALLMENT (weak). The 'three doubles' appear on the "
    "weak side: double rectangle (INSTALLMENT) + double diamond (identifying HAS) + double line (total participation). "
    "A single installment is uniquely identified only by loan_id + installment_no — neither part alone.",
    "This is the model to reach for whenever you see 'owner + repeated sub-items' (orders/items, exams/answer-sheets, hotels/rooms).")

# ============= END-OF-UNIT REVISION AIDS =============
add_cheatsheet("S5–S8 · Cheat Sheet","One-page revision reference",
    [("Design pipeline (RCLP)","Requirements → Conceptual (ER) → Logical (tables) → Physical (storage/indexes). Each stage feeds the next."),
     ("Attribute types","Simple (plain oval) · Composite (branches) · Derived (dashed, computed — never stored) · Multivalued (double oval)."),
     ("Keys","Candidate (any minimal unique set) → Primary (chosen, underlined) → Composite (>1 attribute). Names are never keys."),
     ("Cardinality ratio","1:1 (citizen–passport) · 1:N (dept–employee) · M:N (student–course → junction table in S12). Read the max BOTH ways."),
     ("Participation","Total = double line = must participate. Partial = single line = optional. Cardinality = how many; participation = whether at all."),
     ("Weak entity — 'three doubles'","Double rectangle + double diamond (identifying rel.) + double line (total). Full id = owner PK + partial key (dashed underline).")])

add_glossary("S5–S8 · Glossary","Key terms — quick reference",
    [("Conceptual data model","high-level, DBMS-independent description of WHAT data is held."),
     ("ER model / Chen notation","rectangle = entity, oval = attribute, diamond = relationship."),
     ("Database design process","Requirements → Conceptual → Logical → Physical (RCLP)."),
     ("Data requirement","what to store (nouns → entities/attributes)."),
     ("Functional requirement","what operations to support (verbs → transactions)."),
     ("Entity / type / set","one thing / the category (design-time) / all members now (instance-time)."),
     ("Attribute","property of an entity; simple, composite, derived, or multivalued."),
     ("Key","uniqueness constraint; candidate, primary (underlined), or composite."),
     ("Relationship","association among entities (diamond); has type, set, and degree."),
     ("Degree","number of entity types in a relationship (binary = 2)."),
     ("Cardinality ratio","max instances an entity can join: 1:1, 1:N, M:N."),
     ("Participation","total (mandatory, double line) vs partial (optional, single line)."),
     ("Weak entity","no key of its own; identified via an owner (double rectangle)."),
     ("Identifying relationship","double diamond linking a weak entity to its owner."),
     ("Partial key (discriminator)","distinguishes weak entities of the SAME owner (dashed underline)."),
     ("Full identifier (weak)","owner's primary key + the weak entity's partial key.")])

# ---- close ----
add_divider("Batch note","S9–S12 coming next",
    "This deck covers S5–S8. Next: S9 (full ER diagrams, naming & design issues), S10 (ternary / degree > 2), "
    "S11 (specialization & generalization), S12 (ER → relational tables), plus the consolidated end-of-unit quiz.")
add_title("End of S5–S8 (Unit 2)",
          "Next: S9–S12 — higher-degree relationships, specialization/generalization, ER→tables",
          "Built to COURSE_MATERIAL_STANDARD.md · every concept is self-contained on the slide face · "
          "exports cleanly to PDF with no information lost · Chen-notation diagrams are labelled PNG images.")

# ---- final pass: number every slide (auto-contrast) ----
_add_page_numbers()

# ---- save (path derived relative to this script — machine-independent) ----
out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"IT220_Unit2.pptx")
prs.save(out)
print("Saved",out,"with",len(prs.slides._sldIdLst),"slides")
