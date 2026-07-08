#!/usr/bin/env python3
"""IT220 Unit 1 deck — S1–S4 (Database Concepts & Architecture), rebuilt to
COURSE_MATERIAL_STANDARD.md. Self-contained / PDF-safe two-slide concept pairs;
real example / trap / exam-answer / key terms on the slide face; notes = timing
only. Content distilled from Unit1_material.md; diagrams already in images/.
Run: python3 build_unit1_pptx.py -> IT220_Unit1.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "images")
NAVY=RGBColor(0x0C,0x2B,0x4A); BLUE=RGBColor(0x18,0x5F,0xA5); TEAL=RGBColor(0x0F,0x6E,0x56)
AMBER=RGBColor(0x85,0x4F,0x0B); CORAL=RGBColor(0x99,0x3C,0x1D); LIGHT=RGBColor(0xF2,0xF5,0xF8)
PALET=RGBColor(0xE1,0xF5,0xEE); WHITE=RGBColor(0xFF,0xFF,0xFF); DARK=RGBColor(0x1A,0x1A,0x1A)
GREY=RGBColor(0x55,0x60,0x6B)
BLUE_T=RGBColor(0xE6,0xF1,0xFB); CORAL_T=RGBColor(0xFA,0xEC,0xE7)
TEAL_T=RGBColor(0xE1,0xF5,0xEE); AMBER_T=RGBColor(0xFA,0xEE,0xDA)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

def _bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def _box(s,l,t,w,h): return s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)).text_frame
def _bar(s,c,l,t,w,h):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=c; sp.line.fill.background(); sp.shadow.inherit=False
def _notes(s,t):
    if t: s.notes_slide.notes_text_frame.text=t
def _header(s,kicker,header,kc=BLUE):
    _bar(s,kc,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=kc
    th=_box(s,0.7,0.66,12,1.0); th.word_wrap=True; r=th.paragraphs[0].add_run(); r.text=header
    r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
def _card(s,l,t,w,h,accent,fill,label,body,body_sz=14):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    sp.line.color.rgb=accent; sp.line.width=Pt(1); sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP
    tf.margin_left=Inches(0.22); tf.margin_right=Inches(0.20); tf.margin_top=Inches(0.10); tf.margin_bottom=Inches(0.08)
    p=tf.paragraphs[0]; r=p.add_run(); r.text=label; r.font.size=Pt(12.5); r.font.bold=True; r.font.color.rgb=accent
    p2=tf.add_paragraph(); p2.space_before=Pt(3); r2=p2.add_run(); r2.text=body; r2.font.size=Pt(body_sz); r2.font.color.rgb=DARK
    return sp

def add_title(title,subtitle,footer):
    s=prs.slides.add_slide(BLANK); _bg(s,NAVY); _bar(s,AMBER,0.9,3.05,1.6,0.10)
    tf=_box(s,0.9,2.0,11.5,1.0); tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(34); r.font.bold=True; r.font.color.rgb=WHITE
    t2=_box(s,0.9,3.25,11.5,1.2); t2.word_wrap=True
    r=t2.paragraphs[0].add_run(); r.text=subtitle; r.font.size=Pt(18); r.font.color.rgb=RGBColor(0xCF,0xDD,0xEA)
    t3=_box(s,0.9,6.5,11.5,0.7); t3.word_wrap=True; r=t3.paragraphs[0].add_run(); r.text=footer
    r.font.size=Pt(13); r.font.color.rgb=RGBColor(0x9A,0xB2,0xC6); return s

def add_outcomes(header,kicker,intro,items,nxt):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,header)
    tf=_box(s,0.7,1.75,12.0,4.6); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=intro; r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=NAVY; p.space_after=Pt(8)
    for it in items:
        p=tf.add_paragraph(); r=p.add_run(); r.text="•  "+it; r.font.size=Pt(14.5); r.font.color.rgb=DARK; p.space_after=Pt(5)
    _card(s,0.7,6.35,12.0,0.85,BLUE,BLUE_T,"➡️ Next unit",nxt,body_sz=13)
    return s

def add_roadmap(kicker,title,done,todo):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    _card(s,0.7,1.8,5.85,4.55,TEAL,TEAL_T,"✅  THIS UNIT  ·  S1–S4","\n".join(done),body_sz=14)
    _card(s,6.78,1.8,5.85,4.55,GREY,LIGHT,"⏭️  COMING NEXT  ·  Unit 2+","\n".join(todo),body_sz=14)
    return s

def add_divider(kicker,title,hook,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,BLUE)
    tk=_box(s,0.9,1.9,11.5,0.6); r=tk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=RGBColor(0xBF,0xDC,0xF2)
    tt=_box(s,0.9,2.5,11.5,1.4); tt.word_wrap=True; r=tt.paragraphs[0].add_run(); r.text=title
    r.font.size=Pt(31); r.font.bold=True; r.font.color.rgb=WHITE
    _bar(s,AMBER,0.9,4.05,1.4,0.09)
    th=_box(s,0.9,4.25,11.5,2.4); th.word_wrap=True; p=th.paragraphs[0]; r=p.add_run(); r.text="🎣  "+hook
    r.font.size=Pt(18); r.font.color.rgb=RGBColor(0xE8,0xF1,0xF9); r.font.italic=True
    _notes(s,timing); return s

def concept_understand(kicker,heading,definition,mechanism,image,hook,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker+"  ·  understand",heading)
    tw=6.2 if image else 11.9
    tf=_box(s,0.7,1.72,tw,5.05); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=definition; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=NAVY; p.space_after=Pt(10)
    p=tf.add_paragraph(); r=p.add_run(); r.text="HOW IT WORKS"; r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=BLUE; p.space_after=Pt(5)
    for m in mechanism:
        p=tf.add_paragraph(); r=p.add_run(); r.text="•  "+m; r.font.size=Pt(14); r.font.color.rgb=DARK; p.space_after=Pt(6)
    if image:
        path=os.path.join(IMG,image)
        if os.path.exists(path):
            pic=s.shapes.add_picture(path,Inches(6.95),Inches(1.9),width=Inches(6.05))
            band_top,band_h=1.9,4.55; h_in=pic.height/914400.0
            pic.top=Inches(band_top+max(0.0,(band_h-h_in)/2.0))
    _bar(s,TEAL,0,6.9,SW.inches,0.6)
    ft=_box(s,0.7,6.96,11.9,0.5); p=ft.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text="🧠  Memory hook:  "+hook; r.font.size=Pt(13.5); r.font.bold=True; r.font.color.rgb=WHITE
    _notes(s,timing); return s

def concept_apply(kicker,heading,example,trap,exam,keyterms):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker+"  ·  apply",heading)
    _card(s,0.7,1.66,12.0,1.48,BLUE ,BLUE_T ,"🌍  Real example",example)
    _card(s,0.7,3.26,12.0,1.30,CORAL,CORAL_T,"⚠️  Common trap — and the fix",trap)
    _card(s,0.7,4.68,12.0,1.92,TEAL ,TEAL_T ,"🎯  Exam-ready answer",exam)
    _card(s,0.7,6.72,12.0,0.60,AMBER,AMBER_T,"🔑  Key terms",keyterms,body_sz=11.5)
    return s

def add_activity(title,steps,expected,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,PALET); _bar(s,TEAL,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text="🛠  IN-CLASS ACTIVITY  ·  THINK–PAIR–SHARE"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=TEAL
    th=_box(s,0.7,0.66,12,0.9); r=th.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
    tf=_box(s,0.7,1.75,11.9,3.6); tf.word_wrap=True; first=True
    for text in steps:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text="•  "+text; r.font.size=Pt(15.5); r.font.color.rgb=DARK; p.space_after=Pt(7)
    _card(s,0.7,5.5,11.9,1.7,TEAL,WHITE,"✅  Expected answer",expected)
    _notes(s,timing); return s

def add_quiz(title,items,timing="",compact=False):
    s=prs.slides.add_slide(BLANK); _bg(s,LIGHT); _bar(s,AMBER,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text="CHECK FOR UNDERSTANDING"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=AMBER
    th=_box(s,0.7,0.66,12,0.9); r=th.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
    qz,az,sb=(14,12.5,5) if compact else (17,15,9)
    tf=_box(s,0.7,1.75,11.9,5.3); tf.word_wrap=True; first=True
    for text,kind in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=text
        if kind=="q": r.font.size=Pt(qz); r.font.bold=True; r.font.color.rgb=NAVY; p.space_before=Pt(sb)
        elif kind=="a": r.font.size=Pt(az); r.font.color.rgb=TEAL; r.font.bold=True; p.level=1
        else: r.font.size=Pt(az); r.font.color.rgb=DARK; r.font.italic=True; p.level=1; p.space_before=Pt(sb-3)
    _notes(s,timing); return s

def add_summary(kicker,takeaways,why,nxt,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,"Summary & Key Takeaways",kc=NAVY)
    tf=_box(s,0.7,1.7,12.0,2.7); tf.word_wrap=True; first=True
    for i,tk in enumerate(takeaways,1):
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=f"{i}.  "+tk; r.font.size=Pt(15); r.font.color.rgb=DARK; p.space_after=Pt(8)
    _card(s,0.7,4.55,12.0,1.35,TEAL,TEAL_T,"💡  Why this matters",why,body_sz=13.5)
    _card(s,0.7,6.02,12.0,1.15,BLUE,BLUE_T,"➡️  Next",nxt,body_sz=13.5)
    _notes(s,timing); return s

def add_solved_problem(kicker,title,scenario,steps,answer,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title,kc=AMBER)
    _card(s,0.7,1.58,12.0,1.0,AMBER,AMBER_T,"📋  Problem",scenario,body_sz=13)
    tf=_box(s,0.85,2.72,11.7,2.95); tf.word_wrap=True; first=True
    for i,st in enumerate(steps,1):
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=f"Step {i}.  "; r.font.size=Pt(12.5); r.font.bold=True; r.font.color.rgb=AMBER
        r2=p.add_run(); r2.text=st; r2.font.size=Pt(12.5); r2.font.color.rgb=DARK; p.space_after=Pt(4)
    _card(s,0.7,5.72,12.0,1.5,TEAL,TEAL_T,"✅  Worked answer",answer,body_sz=13)
    _notes(s,timing); return s

def add_cheatsheet(kicker,title,blocks):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    xs=[0.7,6.78]; ys=[1.72,3.36,5.0]; hh=1.55
    for idx,(label,body) in enumerate(blocks[:6]):
        _card(s,xs[idx%2],ys[idx//2],5.85,hh,BLUE,BLUE_T,label,body,body_sz=12)
    return s

def add_glossary(kicker,title,terms):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    half=(len(terms)+1)//2
    for x,group in [(0.7,terms[:half]),(6.78,terms[half:])]:
        tf=_box(s,x,1.75,5.85,5.3); tf.word_wrap=True; first=True
        for term,defn in group:
            p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            r=p.add_run(); r.text=term+" — "; r.font.size=Pt(11.5); r.font.bold=True; r.font.color.rgb=NAVY
            r2=p.add_run(); r2.text=defn; r2.font.size=Pt(11.5); r2.font.color.rgb=DARK; p.space_after=Pt(6)
    return s

def _add_page_numbers():
    total=len(prs.slides._sldIdLst)
    for i,s in enumerate(prs.slides,1):
        try:
            h=str(s.background.fill.fore_color.rgb); rr,gg,bb=int(h[0:2],16),int(h[2:4],16),int(h[4:6],16)
            dark=(0.299*rr+0.587*gg+0.114*bb)<128
        except Exception: dark=False
        col=RGBColor(0xE8,0xF0,0xF7) if dark else GREY
        tb=_box(s,12.35,0.30,0.9,0.35); p=tb.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
        r=p.add_run(); r.text=f"{i} / {total}"; r.font.size=Pt(10); r.font.color.rgb=col

# ============================================================
#                        BUILD
# ============================================================
add_title("Unit 1 — Database Concepts & Architecture",
          "IT 220: Database Management System  ·  BIM 4th Semester",
          "Sessions S1–S4  ·  50 min each  ·  Nepal / South Asia context  ·  "
          "Self-contained slides: every concept is complete on the slide face — exports to PDF with no information lost.")

add_outcomes("Unit 1 — Learning Outcomes","overview  ·  s1–s4",
    "By the end of this unit, you will be able to:",
    ["Define database, DBMS, and the roles around them, and say why a DBMS beats flat files (S1)",
     "Distinguish data models, schemas, and instances (S2)",
     "Explain the three-schema architecture and the two kinds of data independence (S2)",
     "Identify database languages/interfaces and the components of the DBMS environment (S3)",
     "Compare centralized vs client/server architectures and classify DBMSs (S4)"],
    "Unit 2 starts DESIGNING databases — the Entity-Relationship (ER) model and converting it to tables.")

add_roadmap("Unit 1 — Roadmap","Where each session fits (S1–S4)",
    ["S1   Database, DBMS, users, DBA, files vs DBMS",
     "S2   Data models · schema/instance · three-schema · data independence",
     "S3   Languages (DDL/DML/DCL/TCL) · interfaces · components · catalog",
     "S4   Centralized vs client/server · classification · synthesis"],
    ["Unit 2   ER modelling & mapping to tables",
     "Unit 3   Relational algebra & calculus",
     "Unit 4   Normalization",
     "Unit 5  SQL  ·  Unit 6  Transactions  ·  Unit 7  Advanced"])

# ============================ S1 ============================
add_divider("Session 1 · Lecture hour 1 (of 4)","Database & DBMS · Users, DBA & Advantages",
    "Your eSewa balance survives when your phone dies; your exam results are still there after the holidays. None of that lives in the app — it lives in a database. Where does your data actually live?",
    "OPENING HOOK [~5 min]. Build curiosity; don't define 'database' formally yet. Agenda: data vs database → DBMS → users → DBA → why databases beat files.")

concept_understand("S1 · Concept 1 · [THEORY]","Data vs Information vs Database",
    "Data = raw facts with no context. Information = data placed in context so it answers a question. A database = an organized, related collection of data stored so it can be accessed, managed, and updated.",
    ["'78' alone is data (a mark? an age? a temperature?); 'Sita scored 78 in DBMS' is information.",
     "A database is organized and related — records link to each other, not a random heap.",
     "Three defining properties: persistent (survives the program), shared (many users/apps at once), structured.",
     "A kirana notebook is a fragile 'database': one reader, destroyed by a spill, no fast querying."],
    "s1_data_ladder.png","Facts → meaning → organized store.",
    "~7 min. Ask 'is 78 a mark, age, temperature, or house number?' — you can't tell; that's raw data.")
concept_apply("S1 · Concept 1 · [THEORY]","Data vs Information vs Database",
    "One eSewa line — 'Rs 500 to Hari · 2081-02-15 · success' — the bare fields are data; shown so you understand 'you paid Hari Rs 500 and it worked' is information; millions of such records, linked to the right accounts and stored so your full history loads in one tap, are the database.",
    "\"Any stored file is a database.\" Only if it's persistent, shared, and structured. A kirana notebook barely qualifies — one reader at a time, fragile, no efficient querying.",
    "Data are raw facts with no context (500, 'Hari'). Information is data given meaning by context ('You sent Rs 500 to Hari'). A database is an organized, persistent, shared collection of related data (eSewa's store of all transactions). Data becomes information through context; many related pieces stored systematically become a database.",
    "data · information · database · persistent · shared · structured")

concept_understand("S1 · Concept 2 · [THEORY]","DBMS: the software that manages the data",
    "A DBMS (Database Management System) is the software that lets people and applications define, create, query, update, and administer databases. The database is the data; the DBMS is the program that manages it.",
    ["Database + DBMS + apps/users together = a database system.",
     "DBMS engines: MySQL, MariaDB, PostgreSQL, Oracle, SQL Server. One engine manages many databases.",
     "A DBMS provides concurrency control, data integrity, a query language, and security.",
     "'Back up the database' (copy the data) ≠ 'upgrade the DBMS' (install a newer engine)."],
    "s1_users_dbms.png","Data is the books; the DBMS is the librarian.",
    "~7 min. 'MySQL is the DBMS — the engine. The college database is the data it manages.' Say it twice.")
concept_apply("S1 · Concept 2 · [THEORY]","DBMS: the software that manages the data",
    "Your result portal doesn't contain your marks — it asks a DBMS (often MySQL) that manages the marks database. On results night thousands log in at once and the DBMS serves everyone the right marks without records clashing; a shared Excel file, by contrast, lets two clerks silently overwrite each other.",
    "\"Excel is a database.\" No — it lacks real concurrency control, enforced integrity, a query language for millions of rows, and proper security. Fine for 200 rows on a laptop; it breaks exactly at multi-user scale.",
    "A DBMS is software used to define, create, query, update, and administer databases (MySQL, Oracle, PostgreSQL). The database is the actual stored data; the DBMS is the program that manages it — providing concurrency control, integrity, querying, and security. Together (database + DBMS + applications) they form a database system.",
    "DBMS · database system · concurrency control · data integrity · query language · security")

concept_understand("S1 · Concept 3 · [THEORY]","Database Users & the DBA",
    "Different people use a database at different depths: naive/end users, application programmers, sophisticated users, and the Database Administrator (DBA) who administers the database itself.",
    ["Naive/end users — use ready-made screens (a bank teller entering a deposit).",
     "Application programmers — build the apps that talk to the database.",
     "Sophisticated users — write their own queries/analyses (an MIS analyst).",
     "DBA — designs/maintains the schema, controls security & permissions, backup & recovery, performance tuning."],
    "s1_db_users.png","The DBA is air-traffic control for data.",
    "~7 min. Use ONE bank to populate all four roles; make the DBA feel high-stakes.")
concept_apply("S1 · Concept 3 · [THEORY]","Database Users & the DBA",
    "One day in a Nepali bank: the teller (naive user) enters your deposit via a guided form; the team that built the screen are application programmers; the MIS officer pulling a custom loan report is a sophisticated user; the DBA keeps the core-banking database backed up, secured, and fast — and gets paged when it crawls at month-end.",
    "\"The DBA just installs the database.\" The DBA is accountable for the whole data layer — schema, security, backups, and performance. Much of this course (normalization, indexing, transactions, recovery) is the DBA's toolkit.",
    "Users: naive/end users (use ready-made screens), application programmers (build the apps), and sophisticated users (write their own queries). The DBA administers the database itself: designing/maintaining the schema, controlling security and permissions, backup & recovery, and performance tuning — the person accountable for the database.",
    "naive/end user · application programmer · sophisticated user · DBA · schema · security · backup & recovery · tuning")

concept_understand("S1 · Concept 4 · [THEORY]","Why a Database beats a File System",
    "Before databases, data lived in separate files owned by different programs — causing redundancy, inconsistency, poor sharing, weak security, and data loss. A DBMS was invented to fix each problem.",
    ["Uncontrolled redundancy (same fact in many files) → controlled redundancy (store once).",
     "Inconsistency (copies drift apart) → integrity constraints keep data valid.",
     "Hard to share safely → concurrent multi-user access.",
     "Anyone can read/change a file → security & access control. A crash loses data → backup & recovery."],
    "s1_file_vs_dbms.png","One source of truth beats many drifting copies.",
    "~6 min. Walk the address-in-5-files example: update 4, forget 1 → inconsistent, can't tell which is right.")
concept_apply("S1 · Concept 4 · [THEORY]","Why a Database beats a File System",
    "An address stored in five department files; the student moves, a clerk updates four and forgets the fifth — now four files disagree and nobody knows which is right (uncontrolled redundancy → inconsistency). A DBMS stores the address once, so an update in one place is instantly correct everywhere.",
    "\"More copies = safer.\" Uncontrolled duplication is dangerous — each copy can drift and contradict. A DBMS keeps controlled, deliberate backups, not scattered loose copies.",
    "Advantages of a DBMS over a file-based system: (1) controlled redundancy — each fact stored once; (2) consistency & integrity — rules keep data valid; (3) concurrent multi-user access — many users safely at once; (4) security/access control — permissions decide who sees what (also backup & recovery). Each directly fixes a file-system weakness.",
    "file-based system · controlled vs uncontrolled redundancy · inconsistency · integrity constraints · concurrent access · access control · backup & recovery")

add_activity("File or DBMS?",
    ["In pairs (2 min): decide file/spreadsheet or DBMS for — a wedding guest list; a bank's accounts; a one-person diary; Daraz's orders.",
     "For each, name WHICH property decides it — concurrency, integrity, scale, or security.",
     "Share aloud (4 min); build the file-vs-DBMS table from the class's answers."],
    "Bank accounts and Daraz orders clearly need a DBMS (many users, strict integrity, huge scale, security); the diary and short guest list don't. That contrast is the lesson — a DBMS is justified by concurrency, integrity, scale, and security needs.",
    "ACTIVITY [~6 min].")
add_quiz("S1 — Quick Check",
    [("Q1.  Which is NOT a function of a DBMS?","q"),
     ("a) define data   b) ✅ print formatted documents   c) query data   d) control access","a"),
     ("Q2.  Who is responsible for granting and revoking user permissions?","q"),
     ("a) naive user   b) app programmer   c) ✅ DBA   d) end user","a"),
     ("Discussion: name a phone app you use and guess what its database stores about you.","o")],
    "QUIZ [~5 min].")
add_summary("S1 · Summary  ·  [~2 min]",
    ["A database is organized, persistent, shared data; a DBMS is the software that manages it.",
     "Four user types interact at different depths; the DBA is accountable for the data layer.",
     "A DBMS beats files through controlled redundancy, integrity, sharing, security, and recovery."],
    "Every fintech and e-commerce employer in Nepal — eSewa, Khalti, Daraz, IME Pay — runs on a DBMS. 'Understands databases' is a baseline requirement on almost every software, data, or QA job, and this vocabulary is what a first-round interviewer expects you to use correctly.",
    "S2 — how we DESCRIBE data: data models, schemas, instances, and data independence.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S2 ============================
add_divider("Session 2 · Lecture hour 2 (of 4)","Data Models, Schemas & the Three-Schema Architecture",
    "Last year your college stored results one way; this year IT moved everything to a faster server with a different internal layout. You — checking marks online — noticed nothing. Why didn't the change break the portal?",
    "OPENING HOOK [~5 min]. Hold the question until data independence (Concept 4). Agenda: data models → schema vs instance → three-schema → data independence.")

concept_understand("S2 · Concept 1 · [THEORY]","Data Models",
    "A data model is a set of concepts describing the structure, constraints, and operations of a database — the 'language' for describing what data looks like and how it behaves. Models exist at levels of abstraction.",
    ["High-level / conceptual — real-world things & relationships (the ER model, Unit 2).",
     "Representational / implementation — the relational model (tables of rows & columns); what most business DBMSs use.",
     "Low-level / physical — how data is actually stored on disk (files, pages, indexes).",
     "The same system, described at several levels of detail."],
    "s2_data_models.png","Idea → tables → disk.",
    "~7 min. Light touch — models are detailed in Unit 2; the point is describing data at different abstraction levels.")
concept_apply("S2 · Concept 1 · [THEORY]","Data Models",
    "Before Daraz was built, someone sketched 'a Customer places Orders; each Order contains Products' — a conceptual (ER) model. Engineers turned it into relational tables, and the ops team decided how to store and index those tables on disk. Same shop, described at three levels.",
    "\"A data model is just the tables.\" Tables are the representational level; above them is the conceptual model (real-world things), below is the physical storage layout.",
    "A data model is a set of concepts describing the structure, constraints, and operations of a database. By abstraction: high-level/conceptual (the ER model), representational/implementation (the relational/table model), and low-level/physical (storage details) — ranging from human-oriented to machine-oriented.",
    "data model · conceptual (high-level) · representational/relational · physical (low-level) · structure, constraints, operations")

concept_understand("S2 · Concept 2 · [THEORY]","Schema vs Instance",
    "A schema is the design/structure of the database (tables, columns, rules) — decided at design time, changes rarely. An instance (database state) is the actual data stored at a given moment — it changes constantly.",
    ["Schema Student(roll, name, program) says what every row must have.",
     "The instance is the actual rows right now — 60 today, 59 after a student drops out tomorrow.",
     "One schema, an endless succession of instances over time.",
     "Adding a phone column changes the SCHEMA (rare, big); adding a row changes the INSTANCE (constant)."],
    "s2_schema_instance.png","Schema = the mould; instance = what's in it now.",
    "~8 min. Write Student(roll, name, program) then 3 fake rows beside it — same schema, many instances.")
concept_apply("S2 · Concept 2 · [THEORY]","Schema vs Instance",
    "Your class WhatsApp group: its structure (name, members, admins) is the schema and barely changes; the actual messages and current member list right now are the instance and change every minute. Two members leave tomorrow — same group (schema), new membership (instance).",
    "\"Schema and instance are the same thing.\" The schema is the stable design; the instance is the volatile data. One schema has many instances over time.",
    "A schema is the design/structure of a database — the table definitions and constraints; it is stable and changes rarely. An instance (database state) is the actual data stored at a given moment; it changes constantly. One schema has many instances over time. Example: schema Student(roll, name, program); an instance = today's actual rows.",
    "schema (structure, stable) · instance / database state (current data, volatile) · one schema → many instances")

concept_understand("S2 · Concept 3 · [THEORY]","The Three-Schema Architecture",
    "The three-schema architecture describes the same database at three levels: external (per-user views), conceptual (the whole logical structure), and internal (physical storage). Mappings between the levels hide detail.",
    ["External level — per-user views (a student sees only their own marks); a database has many views.",
     "Conceptual level — one complete logical schema: all entities, relationships, constraints, no storage detail.",
     "Internal level — physical storage: files, pages, indexes.",
     "Mappings (external/conceptual, conceptual/internal) hide each level from the one above → enable data independence."],
    "s2_three_schema.png","View → whole → storage.",
    "~9 min. The centrepiece — walk the college exam DB through all three levels; the mappings give data independence.")
concept_apply("S2 · Concept 3 · [THEORY]","The Three-Schema Architecture",
    "The college result system through three pairs of eyes: a student sees only their own marksheet (external view); the exam section's system holds every student, subject, and mark with their relationships (conceptual); the IT team knows it's indexed files on a server (internal). Each level is shielded from the others.",
    "\"There are many conceptual schemas.\" No — many external views, exactly ONE conceptual schema, one internal schema. The single conceptual schema is the shared source of truth.",
    "The three-schema architecture separates a database into: external level (per-user views, e.g. a student sees only their own marks), conceptual level (the whole logical database — entities, relationships, constraints), and internal level (physical storage — files and indexes). Mappings between levels hide lower-level detail and enable data independence.",
    "external schema/view · conceptual schema · internal schema · external/conceptual & conceptual/internal mappings · separation of concerns")

concept_understand("S2 · Concept 4 · [THEORY]","Data Independence",
    "Data independence is the ability to change the database at one level without changing the levels above. Logical = change the conceptual schema without breaking external views; physical = change storage without changing the conceptual schema.",
    ["Logical independence — add a column/table without breaking views that don't use it; harder to achieve.",
     "Physical independence — move to SSDs, add an index, reorganize files without changing the logical schema; easier.",
     "It is the payoff of the three-schema architecture and its mappings.",
     "Lets banks, telecoms, and portals upgrade infrastructure without rewriting every app."],
    "s2_data_independence.png","Logical = change the design safely; physical = change the storage safely.",
    "~6 min. Resolve the hook: SSD migration + new index overnight, every app works next morning = physical independence.")
concept_apply("S2 · Concept 4 · [THEORY]","Data Independence",
    "One night the admin migrates the exam database to faster SSDs and adds an index; the next morning every app and user works exactly as before — nobody noticed. That invisible change is physical data independence, and it's why the college could move servers without breaking the student portal.",
    "\"Any change to the database breaks the apps.\" With data independence, storage changes (physical) and even design additions (logical) leave dependent views and apps working untouched.",
    "Logical data independence is the ability to change the conceptual schema (e.g. add a column or table) without changing external views or the applications using them. Physical data independence is the ability to change physical storage (e.g. add an index, move to SSDs) without changing the conceptual schema. The three-schema architecture provides both; logical independence is harder to achieve.",
    "data independence · logical (change the design safely) · physical (change the storage safely)")

add_activity("Draw your college's three schemas",
    ["In pairs (3 min): for the college result system, write what lives at EACH level.",
     "External — what does a student / teacher / admin each see?   Conceptual — the whole database.   Internal — how it's stored.",
     "Then name ONE change at the internal level that students would NOT notice — that's physical data independence.",
     "Share (2 min)."],
    "External: a student sees only their marks, a teacher their class, admin everything. Conceptual: one logical database of all students/subjects/marks. Internal: indexed files on a server. A change students wouldn't notice: switching to SSDs, adding an index, or compressing old data.",
    "ACTIVITY [~5 min].")
add_quiz("S2 — Quick Check",
    [("Q1.  Adding a column without breaking existing user views demonstrates:","q"),
     ("a) ✅ logical data independence   b) physical data independence   c) redundancy   d) an instance change","a"),
     ("Q2.  Which level is closest to physical storage?","q"),
     ("a) external   b) conceptual   c) ✅ internal   d) view","a"),
     ("Discussion: give an everyday example of 'changing the inside without changing the outside.'","o")],
    "QUIZ [~5 min].")
add_summary("S2 · Summary  ·  [~2 min]",
    ["Data models describe data at three abstraction levels (conceptual → representational → physical).",
     "Schema = structure (stable); instance = the current data (always changing).",
     "The three-schema architecture (external / conceptual / internal) gives logical and physical data independence."],
    "Data independence is why banks, telecoms (Ntc/Ncell), and government portals can upgrade storage and infrastructure without rewriting every application. Systems built without it become unmaintainable — a real, expensive business risk.",
    "S3 — the languages and interfaces we use to talk to a DBMS, and what's inside the DBMS box.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S3 ============================
add_divider("Session 3 · Lecture hour 3 (of 4)","Database Languages, Interfaces & the DBMS Environment",
    "SQL is several languages wearing one name: one to BUILD a database, another to USE it daily, another to PROTECT it. Which SQL keywords have you already seen?",
    "OPENING HOOK [~5 min]. Full SQL syntax is Unit 5. Agenda: languages → interfaces → components → data dictionary.")

concept_understand("S3 · Concept 1 · [THEORY]","Database Languages (DDL / DML / DCL / TCL)",
    "SQL is a family of sub-languages, each with a job: DDL defines structure, DML queries/modifies data, DCL controls permissions, TCL controls transactions.",
    ["DDL — CREATE, ALTER, DROP (define / change structure).",
     "DML — SELECT, INSERT, UPDATE, DELETE (query / modify data).",
     "DCL — GRANT, REVOKE (permissions).",
     "TCL — COMMIT, ROLLBACK (confirm or undo a group of changes as a unit)."],
    "s3_languages.png","Define, Manipulate, Control, Transact.",
    "~8 min. One college example: register a course = DDL; enrol a student = DML; grant a TA access = DCL; commit = TCL.")
concept_apply("S3 · Concept 1 · [THEORY]","Database Languages (DDL / DML / DCL / TCL)",
    "Building the result portal: the developer CREATEs the Student and Marks tables (DDL); staff INSERT/UPDATE marks (DML); the admin GRANTs the clerk access then REVOKEs it once published (DCL); each batch is COMMITted, or ROLLBACKed if a mistake is caught (TCL). One workflow, all four, all 'SQL'.",
    "\"SQL is one single language.\" SQL bundles four sub-languages — CREATE is DDL, SELECT is DML, GRANT is DCL, COMMIT is TCL — all 'SQL' but doing very different jobs.",
    "DDL (CREATE/ALTER/DROP) defines structure; DML (SELECT/INSERT/UPDATE/DELETE) queries and modifies data; DCL (GRANT/REVOKE) controls permissions; TCL (COMMIT/ROLLBACK) controls transactions. SQL bundles all four sub-languages into one language.",
    "SQL · DDL (define) · DML (query/modify) · DCL (permissions) · TCL (transactions) · example verbs")

concept_understand("S3 · Concept 2 · [THEORY]","Interfaces to a DBMS",
    "A DBMS can be reached through several interfaces, each suited to a different user: menu-based, form-based, graphical (GUI), natural-language, and application-program (API) / embedded SQL. The principle: match the interface to the user.",
    ["Form-based / GUI — guided, safe, ideal for non-technical staff.",
     "Menu-based — pick from lists (kiosks, older systems).",
     "Natural-language — type a request in ordinary words.",
     "API / embedded SQL — programs talk to the database directly in code (developers)."],
    "s3_interfaces.png","Same data, many doors.",
    "~7 min. A teller and a developer touch the SAME database through different interfaces matched to their skill.")
concept_apply("S3 · Concept 2 · [THEORY]","Interfaces to a DBMS",
    "eSewa's front doors: your grandmother uses big tap-buttons (a form/GUI); a shopkeeper takes payment via QR; eSewa's developers reach the very same database through APIs and SQL. Force the grandmother to type SQL and she'd never pay a bill; force a developer to tap buttons and the app would never get built.",
    "\"Everyone should use the same interface.\" No — match the interface to the user's skill and task; the same database is reached through whichever door fits.",
    "DBMS interfaces include menu-based, form-based, graphical (GUI), natural-language, and application-program (API) / embedded-SQL interfaces. Each suits a different kind of user; the guiding principle is to match the interface to the user's skill and task — the same database is reached through whichever door fits.",
    "interface · menu-based · form-based · GUI · natural-language · API / embedded SQL · match interface to user")

concept_understand("S3 · Concept 3 · [THEORY]","Inside the DBMS — the System Environment",
    "A DBMS is not one monolithic program — it is several cooperating components: query processor/optimizer, storage manager, buffer manager, transaction manager, and the catalog.",
    ["Query processor / optimizer — plans the cheapest, fastest way to run a query.",
     "Storage manager — reads/writes data to disk.   Buffer manager — caches hot data in memory.",
     "Transaction manager — keeps transactions safe and reliable (ACID, Unit 6).",
     "Catalog / metadata — stores information about the database itself (next concept)."],
    "s3_dbms_components.png","A DBMS is a kitchen of cooperating roles.",
    "~8 min. Restaurant analogy: waiter (query processor), kitchen (storage), manager (transactions), receipt book (catalog).")
concept_apply("S3 · Concept 3 · [THEORY]","Inside the DBMS — the System Environment",
    "Tap 'View Result': the query processor finds the fastest way to fetch your marks; the buffer manager checks whether they're already in memory; if not, the storage manager reads disk; the transaction manager ensures you never see half-updated marks; the catalog confirms the Marks table exists — five components cooperate in under a second.",
    "\"The DBMS is one black box.\" It's cooperating components — so 'it's slow' becomes a diagnosable checklist: bad execution plan? starved cache? failing disk? waiting on a lock?",
    "Main DBMS components: query processor/optimizer (plans the most efficient execution), storage manager (reads/writes data to disk), buffer manager (caches frequently used data in memory), transaction manager (ensures safe, reliable transactions — ACID), and the catalog/data dictionary (stores metadata). They cooperate to answer queries quickly and safely.",
    "query processor/optimizer · storage manager · buffer manager · transaction manager · catalog/metadata")

concept_understand("S3 · Concept 4 · [THEORY]","Data Dictionary / Catalog",
    "The data dictionary (system catalog) stores metadata — data about the data: table names, columns and their types, constraints, users and permissions, and indexes. It makes the DBMS self-describing.",
    ["When the DBMS checks 'does Student exist and what columns?', it reads the catalog.",
     "Metadata = data that describes other data (a phone book's 'sorted by surname' note).",
     "Tools like MySQL Workbench / phpMyAdmin list your tables by reading the catalog.",
     "Self-describing → the DBMS can validate a query (does this column exist?) before running it."],
    "s3_data_dictionary.png","Metadata = data about data.",
    "~6 min. When you ask 'what columns does Student have?', the answer comes FROM the catalog.")
concept_apply("S3 · Concept 4 · [THEORY]","Data Dictionary / Catalog",
    "Open the college database in MySQL Workbench or phpMyAdmin and it instantly lists every table and each table's columns — read straight from the system catalog, the database's own metadata. Delete a column and the catalog updates; every tool then reflects the change automatically.",
    "\"The catalog holds the actual data.\" No — it holds metadata (structure, types, constraints, permissions, indexes): data ABOUT the data, not the rows themselves.",
    "The data dictionary / system catalog is the part of the DBMS that stores metadata — data about the database: table names, column names and types, constraints, registered users and permissions, and indexes. The DBMS uses it to validate and process queries, making the database self-describing.",
    "data dictionary / system catalog · metadata · table names, column types, constraints, users, indexes · self-describing")

add_activity("Sort the command",
    ["Rapid-fire (class or pairs): classify each as DDL / DML / DCL / TCL —",
     "CREATE TABLE · SELECT · GRANT · DROP · UPDATE · COMMIT · REVOKE · INSERT.",
     "Then: 'enrol a student in a brand-new course and make it permanent' — which sub-languages, in order?",
     "Discuss any disagreements."],
    "DDL: CREATE TABLE, DROP. DML: SELECT, UPDATE, INSERT. DCL: GRANT, REVOKE. TCL: COMMIT. The chained task = DDL (create the course) → DML (enrol) → TCL (commit) — one real task touches several sub-languages.",
    "ACTIVITY [~6 min].")
add_quiz("S3 — Quick Check",
    [("Q1.  CREATE TABLE belongs to which sub-language?","q"),
     ("a) ✅ DDL   b) DML   c) DCL   d) TCL","a"),
     ("Q2.  Which component decides the cheapest way to execute a query?","q"),
     ("a) storage manager   b) buffer manager   c) ✅ query processor / optimizer   d) catalog","a"),
     ("Discussion: which DBMS interface would you build for a non-technical shopkeeper, and why?","o")],
    "QUIZ [~5 min].")
add_summary("S3 · Summary  ·  [~2 min]",
    ["DDL / DML / DCL / TCL each have a distinct job; SQL bundles all four.",
     "A DBMS offers many interface styles for many kinds of users — match the interface to the user.",
     "A DBMS is several cooperating components (query processor, storage/buffer/transaction managers) plus a catalog of metadata."],
    "Knowing DDL / DML / DCL is the literal day-one skill in any backend, data, or QA role. When a job ad says 'SQL required', it is asking for exactly these sub-languages — and saying which command does which job marks you as someone who understands databases, not just syntax.",
    "S4 — WHERE the database runs: centralized vs client/server, and how we classify DBMSs.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S4 ============================
add_divider("Session 4 · Lecture hour 4 (of 4) · closes Unit 1","Centralized vs Client/Server · Classifying DBMSs",
    "Results are out; 5,000 students hit refresh in the same minute. What stops the system collapsing? Architecture — WHERE the database and apps actually run.",
    "OPENING HOOK [~5 min]. Agenda: centralized → client/server (2-tier/3-tier) → classification → Unit 1 synthesis.")

concept_understand("S4 · Concept 1 · [THEORY]","Centralized Architecture",
    "In a centralized architecture the data, the DBMS, and the application all run on one central machine; users connect through simple terminals that mostly just display results.",
    ["All the real work happens in one place.",
     "Advantage: simple to set up, manage, and secure — only one machine to look after.",
     "Disadvantage: a single point of failure — if it goes down, everything stops.",
     "It also struggles to scale to large numbers of simultaneous users."],
    "s4_centralized.png","One machine = one point of failure.",
    "~7 min. Emphasise the single point of failure — the weakness that motivates everything that follows.")
concept_apply("S4 · Concept 1 · [THEORY]","Centralized Architecture",
    "A Kathmandu trading firm ran its entire inventory and billing on one back-office PC; during the Dashain rush its disk failed with no backup, and the whole business stopped for two days. That is the defining weakness of a centralized architecture: everything on one machine, so that machine is a single point of failure.",
    "\"Centralized is safest because it's all in one place.\" That one place is a single point of failure and can't scale — its simplicity is exactly its fragility.",
    "In a centralized architecture, the data, DBMS, and application all run on one central machine, with terminals only displaying results. It is simple to set up and manage, but its main drawback is being a single point of failure (if the central machine fails, everything stops) together with poor scalability.",
    "centralized architecture · single point of failure · scalability · terminal")

concept_understand("S4 · Concept 2 · [THEORY]","Client/Server Architecture",
    "In a client/server architecture the work is split across machines over a network. Two-tier: client ↔ database server. Three-tier: client (UI) ↔ application server (business logic) ↔ database server.",
    ["Two-tier: the client (UI + application logic) talks directly to a database server.",
     "Three-tier: an application server (rules, processing, security) sits in the middle — the modern web/mobile standard.",
     "Each tier can be scaled, secured, and updated independently.",
     "Real data lives server-side: for security, one authoritative copy, and access from any device."],
    "s4_client_server.png","Three tiers: face, logic, data.",
    "~9 min. Three-tier is the modern standard; 'why your phone doesn't hold the full ledger' lands the server-side idea.")
concept_apply("S4 · Concept 2 · [THEORY]","Client/Server Architecture",
    "On Daraz your phone runs the client (screens); application servers handle logins, cart, and orders; database servers store products, orders, and accounts. During a Dashain sale Daraz adds more application servers without touching the database design — three tiers scaling independently.",
    "\"Your phone stores your full bank ledger.\" No — the real data lives server-side (security if the phone is lost, one authoritative copy, access from any new device); the phone is just a client showing a view.",
    "Two-tier: the client (UI + application logic) communicates directly with a database server. Three-tier: an application server holding the business logic sits between the client (UI only) and the database server. Three-tier scales better, is easier to secure, and is the standard for modern web and mobile applications.",
    "client/server · two-tier · three-tier · application server / business logic · server-side")

concept_understand("S4 · Concept 3 · [THEORY]","Classification of DBMSs",
    "A DBMS can be classified along several independent axes: by data model, by number of users, by number of sites/distribution, and by purpose.",
    ["Data model: relational (dominant), object-oriented, object-relational, legacy hierarchical/network, NoSQL.",
     "Number of users: single-user vs multi-user (almost all business systems).",
     "Distribution: centralized (one site) vs distributed (data across many sites).",
     "Purpose: general-purpose vs special-purpose."],
    "s4_classification.png","Model · users · distribution.",
    "~7 min. Relational still DOMINATES business systems; NoSQL is an additional tool for specific needs.")
concept_apply("S4 · Concept 3 · [THEORY]","Classification of DBMSs",
    "Classify the apps in your pocket: a bank app's DBMS is relational, multi-user, and distributed; a small offline note app is relational, single-user, centralized; Facebook mixes relational + NoSQL, is massively multi-user, and heavily distributed. The same three axes describe any system.",
    "\"NoSQL means SQL is dead.\" It means 'not only SQL.' Relational databases still dominate business systems; NoSQL is an additional family for huge scale or flexible, changing data — it complements relational, not replaces it.",
    "DBMSs are classified by data model (relational, object, object-relational, NoSQL, legacy hierarchical/network), by number of users (single-user vs multi-user), by number of sites/distribution (centralized vs distributed), and by purpose (general- vs special-purpose). Note: 'NoSQL' means 'not only SQL', not 'no SQL'.",
    "classification axes · data model · number of users · distribution · purpose · relational · NoSQL ('not only SQL')")

concept_understand("S4 · Concept 4 · [THEORY]","Unit 1 Synthesis — the whole unit in one picture",
    "Users reach a DBMS through interfaces & languages; the DBMS (cooperating components + catalog) manages a database described by a schema at three levels (giving data independence), running on a centralized or client/server architecture.",
    ["Users (naive, programmers, sophisticated, DBA) + interfaces/languages (DDL/DML/DCL/TCL; forms, GUIs, APIs).",
     "DBMS components: query processor, storage & buffer managers, transaction manager, catalog.",
     "Database structure = the three-schema architecture → data independence.",
     "Deployment = centralized or (usually) client/server; classify by model, users, distribution."],
    "s4_synthesis.png","User → interface → DBMS → schema → architecture.",
    "~6 min. Walk the mind-map; ask which session each branch came from. This is the slide to photograph for revision.")
concept_apply("S4 · Concept 4 · [THEORY]","Unit 1 Synthesis — the whole unit in one picture",
    "One eSewa tap: you (a user) tap a button (an interface) that sends a query (DML) to eSewa's DBMS (its components + catalog); it reads your balance from a database whose design is a schema at three levels — which is why eSewa can upgrade storage without breaking your app (data independence) — all on a three-tier client/server, relational, multi-user, distributed DBMS. One tap exercises the whole unit.",
    "\"These are five unrelated topics.\" They're one system: users → interface → DBMS → schema/architecture → classification. Each session is one branch of the same picture.",
    "Users reach a DBMS through interfaces and languages; the DBMS (built of cooperating components plus a catalog) manages a database described by a schema at three levels (giving data independence) and runs on a centralized or client/server architecture; any such system can be classified by its data model, number of users, and distribution.",
    "database · DBMS · users · DBA · schema · instance · three-schema · data independence · DDL/DML/DCL/TCL · components · catalog · centralized vs client/server · classification")

add_activity("Classify a Nepali app",
    ["In pairs (3 min): pick a Nepali app — eSewa, Daraz, Nagarik App, a bank app.",
     "Sketch whether it's two-tier or three-tier, and label what lives in each tier.",
     "Classify its likely DBMS on the three axes (data model, number of users, distribution) with a one-line reason each.",
     "Share (2 min)."],
    "Most modern Nepali apps are three-tier (client app · application servers · database servers), multi-user, and relational (sometimes with some NoSQL), and larger ones are distributed. The reasoning matters more than being 'right'.",
    "ACTIVITY [~5 min].")
add_quiz("S4 — Quick Check",
    [("Q1.  The main weakness of a purely centralized architecture is:","q"),
     ("a) too many servers   b) ✅ single point of failure   c) too secure   d) no data storage","a"),
     ("Q2.  A three-tier architecture adds which layer between the client and the database?","q"),
     ("a) a second client   b) a backup disk   c) ✅ an application / business-logic server   d) a second database","a"),
     ("Discussion: pick a Nepali app and sketch whether it's two-tier or three-tier, and why.","o")],
    "QUIZ [~5 min].")
add_summary("S4 · Summary  ·  [~2 min]",
    ["Centralized = simple but fragile (a single point of failure).",
     "Client/server (especially three-tier) is the scalable modern standard: client · application server · database server.",
     "DBMSs are classified by data model, number of users, and distribution."],
    "Every scalable product you use — e-commerce, fintech, the Nagarik App, government portals — is built client/server, usually three-tier. Architecture is the map; the rest of IT 220 (ER modelling, normalization, SQL) fills in how the database layer is designed and queried.",
    "Unit 2 — DESIGNING databases: the Entity-Relationship (ER) model and converting it to tables.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============= CAPSTONE SOLVED PROBLEM (synthesises Unit 1) =============
add_solved_problem("Unit 1 · Capstone Solved Problem","Classify & architect a nationwide banking system",
    "A nationwide bank asks you to describe its database system: classify its DBMS on the standard axes and "
    "choose a deployment architecture, justifying each decision from Unit 1's concepts.",
    ["Data model → relational: banking needs strict integrity and safe transactions (ACID), which the relational model + SQL provide.",
     "Number of users → multi-user: thousands of tellers, apps, and customers access it at the same time.",
     "Distribution → distributed: data replicated across data centres for scale and disaster recovery.",
     "Architecture → three-tier client/server: client (teller/app screen) · application server (business logic, security) · database server (the data).",
     "Reject centralized: one machine would be a single point of failure and couldn't scale to the load.",
     "Data independence → storage can be upgraded (SSDs, indexes) without rewriting apps (physical independence)."],
    "The system is relational, multi-user, and distributed, deployed as a three-tier client/server; centralized is "
    "rejected because of the single-point-of-failure risk. There is one authoritative server-side copy of the data, "
    "and data independence keeps future storage upgrades cheap — no application rewrites.",
    "Pause and let students justify each axis before revealing. This ties S1–S4 into one exam-style answer.")

# ============= END-OF-UNIT REVISION AIDS =============
add_cheatsheet("Unit 1 · Cheat Sheet","One-page revision reference",
    [("Data · Information · Database","Raw facts → +context → an organized, persistent, shared, structured collection of related data."),
     ("DBMS vs database","DBMS = the software (MySQL). Database = the data. A DBMS adds concurrency, integrity, querying, security."),
     ("Users & DBA","Naive user · application programmer · sophisticated user. DBA = schema, security, backup & recovery, tuning."),
     ("Schema / instance / three-schema","Schema = design (stable); instance = data (volatile). External · conceptual · internal → data independence."),
     ("SQL sub-languages","DDL (define) · DML (query/modify) · DCL (permissions) · TCL (transactions). SQL bundles all four."),
     ("Architecture & classification","Centralized (single point of failure) vs client/server (2- / 3-tier). Classify by model · users · distribution.")])

add_glossary("Unit 1 · Glossary","Key terms — quick reference",
    [("Data / information","raw facts / data given meaning by context."),
     ("Database","organized, persistent, shared collection of related data."),
     ("DBMS","software that defines, creates, queries, and administers databases."),
     ("Database system","database + DBMS + applications and users."),
     ("DBA","administers the database: schema, security, backup, tuning."),
     ("Data model","concepts describing structure, constraints, operations."),
     ("Schema / instance","the design (stable) / the current data (volatile)."),
     ("Three-schema architecture","external (views) · conceptual · internal (storage)."),
     ("Data independence","change one level without changing the level above."),
     ("DDL / DML","define structure / query & modify data."),
     ("DCL / TCL","control permissions / control transactions."),
     ("Interface","menu, form, GUI, natural-language, API — matched to the user."),
     ("Query processor","plans the cheapest way to run a query."),
     ("Buffer / storage manager","cache hot data in memory / read & write disk."),
     ("Catalog / metadata","data about the data; makes the DBMS self-describing."),
     ("Centralized / client-server","one machine (SPOF) / work split across tiers.")])

# ============= CONSOLIDATED END-OF-UNIT QUIZ =============
add_divider("Unit 1 · Revision","End-of-unit quiz (S1–S4)",
    "Twelve MCQs across the whole unit (answers shown), then short-answer, applied, and discussion questions to work from the concept slides.",
    "Use as a 15–20 min in-class quiz or take-home review.")
add_quiz("Section A — Multiple choice (answers shown)",
    [("1.  A database is best described as   →  ✅ an organized, persistent, shared collection of related data","a"),
     ("2.  The software that manages a database is the   →  ✅ DBMS","a"),
     ("3.  A file-system problem a DBMS solves   →  ✅ uncontrolled redundancy / inconsistency","a"),
     ("4.  Granting permissions is the responsibility of the   →  ✅ DBA","a"),
     ("5.  The actual data at a given moment is the   →  ✅ instance","a"),
     ("6.  Changing storage without affecting the conceptual schema is   →  ✅ physical data independence","a"),
     ("7.  The external level refers to   →  ✅ per-user views","a"),
     ("8.  CREATE TABLE is part of   →  ✅ DDL","a"),
     ("9.  Finds the most efficient way to run a query   →  ✅ query processor / optimizer","a"),
     ("10.  Main weakness of centralized architecture   →  ✅ single point of failure","a"),
     ("11.  The middle tier in three-tier is the   →  ✅ application / business-logic server","a"),
     ("12.  'NoSQL' most accurately means   →  ✅ not only SQL","a")],
    "Consolidated quiz Section A.",compact=True)
add_quiz("Sections B–D — short answer, applied & discussion",
    [("Section B — Short answer","q"),
     ("13. Define schema & instance (example each for Student).   14. Four advantages of a DBMS over files.","o"),
     ("15. Logical vs physical data independence (one sentence each).   16. Name DDL/DML/DCL/TCL + one command each.   17. What is metadata, and where is it stored?","o"),
     ("Section C — Applied / diagramming","q"),
     ("18. Draw & label the three-schema architecture; say what each level hides from the one above.","o"),
     ("19. Sketch a three-tier architecture for Daraz; label each tier.   20. Classify an online banking system on three axes with a one-line reason each.","o"),
     ("Section D — Discussion","q"),
     ("21. 'A kirana shop's paper notebook is technically a database.' Argue for and against, then say what a DBMS would add.","o")],
    "Consolidated quiz Sections B–D. Model answers live in the concept slides and Unit1_material.md.",compact=True)

# ---- close ----
add_title("End of Unit 1  ·  IT 220",
          "S1–S4 complete: database & DBMS · users & DBA · data models, schema/instance, three-schema & data independence · languages, interfaces & components · architectures & classification",
          "Built to COURSE_MATERIAL_STANDARD.md · self-contained slides · exports cleanly to PDF · "
          "Next unit (Unit 2): the Entity-Relationship (ER) model.")

_add_page_numbers()
out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"IT220_Unit1.pptx")
prs.save(out)
print("Saved",out,"with",len(prs.slides._sldIdLst),"slides")
