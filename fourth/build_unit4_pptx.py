#!/usr/bin/env python3
"""IT220 Unit 4 deck — S18–S21 (Database Normalization), built to COURSE_MATERIAL_STANDARD.md.
Self-contained / PDF-safe two-slide concept pairs + solved-problem decomposition walkthroughs.
Content distilled from Unit4_material.md; diagrams in images/. Runs the ENROLL example 1NF→4NF.
Run: python3 build_unit4_pptx.py -> IT220_Unit4.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "images")
NAVY=RGBColor(0x0C,0x2B,0x4A); BLUE=RGBColor(0x18,0x5F,0xA5); TEAL=RGBColor(0x0F,0x6E,0x56)
AMBER=RGBColor(0x85,0x4F,0x0B); CORAL=RGBColor(0x99,0x3C,0x1D); LIGHT=RGBColor(0xF2,0xF5,0xF8)
PALET=RGBColor(0xE1,0xF5,0xEE); WHITE=RGBColor(0xFF,0xFF,0xFF); DARK=RGBColor(0x1A,0x1A,0x1A)
GREY=RGBColor(0x55,0x60,0x6B)
BLUE_T=RGBColor(0xE6,0xF1,0xFB); CORAL_T=RGBColor(0xFA,0xEC,0xE7)
TEAL_T=RGBColor(0xE1,0xF5,0xEE); AMBER_T=RGBColor(0xFA,0xEE,0xDA)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

def _bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def _box(s,l,t,w,h): return s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)).text_frame
def _bar(s,c,l,t,w,h):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=c; sp.line.fill.background(); sp.shadow.inherit=False
def _notes(s,t):
    if t: s.notes_slide.notes_text_frame.text=t
def _header(s,kicker,header,kc=BLUE):
    _bar(s,kc,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=kc
    th=_box(s,0.7,0.66,12,1.0); th.word_wrap=True; r=th.paragraphs[0].add_run(); r.text=header
    r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
def _card(s,l,t,w,h,accent,fill,label,body,body_sz=14):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    sp.line.color.rgb=accent; sp.line.width=Pt(1); sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP
    tf.margin_left=Inches(0.22); tf.margin_right=Inches(0.20); tf.margin_top=Inches(0.10); tf.margin_bottom=Inches(0.08)
    p=tf.paragraphs[0]; r=p.add_run(); r.text=label; r.font.size=Pt(12.5); r.font.bold=True; r.font.color.rgb=accent
    p2=tf.add_paragraph(); p2.space_before=Pt(3); r2=p2.add_run(); r2.text=body; r2.font.size=Pt(body_sz); r2.font.color.rgb=DARK
    return sp

def add_title(title,subtitle,footer):
    s=prs.slides.add_slide(BLANK); _bg(s,NAVY); _bar(s,AMBER,0.9,3.05,1.6,0.10)
    tf=_box(s,0.9,2.0,11.5,1.0); tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(34); r.font.bold=True; r.font.color.rgb=WHITE
    t2=_box(s,0.9,3.25,11.5,1.2); t2.word_wrap=True
    r=t2.paragraphs[0].add_run(); r.text=subtitle; r.font.size=Pt(18); r.font.color.rgb=RGBColor(0xCF,0xDD,0xEA)
    t3=_box(s,0.9,6.5,11.5,0.7); t3.word_wrap=True; r=t3.paragraphs[0].add_run(); r.text=footer
    r.font.size=Pt(13); r.font.color.rgb=RGBColor(0x9A,0xB2,0xC6); return s

def add_outcomes(header,kicker,intro,items,nxt):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,header)
    tf=_box(s,0.7,1.75,12.0,4.6); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=intro; r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=NAVY; p.space_after=Pt(8)
    for it in items:
        p=tf.add_paragraph(); r=p.add_run(); r.text="•  "+it; r.font.size=Pt(14.5); r.font.color.rgb=DARK; p.space_after=Pt(5)
    _card(s,0.7,6.35,12.0,0.85,BLUE,BLUE_T,"↩️ Builds on",nxt,body_sz=13)
    return s

def add_roadmap(kicker,title,done,todo):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    _card(s,0.7,1.8,5.85,4.55,TEAL,TEAL_T,"✅  THIS UNIT  ·  S18–S21","\n".join(done),body_sz=14)
    _card(s,6.78,1.8,5.85,4.55,GREY,LIGHT,"⏭️  THEN","\n".join(todo),body_sz=14)
    return s

def add_divider(kicker,title,hook,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,BLUE)
    tk=_box(s,0.9,1.9,11.5,0.6); r=tk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=RGBColor(0xBF,0xDC,0xF2)
    tt=_box(s,0.9,2.5,11.5,1.4); tt.word_wrap=True; r=tt.paragraphs[0].add_run(); r.text=title
    r.font.size=Pt(31); r.font.bold=True; r.font.color.rgb=WHITE
    _bar(s,AMBER,0.9,4.05,1.4,0.09)
    th=_box(s,0.9,4.25,11.5,2.4); th.word_wrap=True; p=th.paragraphs[0]; r=p.add_run(); r.text="🎣  "+hook
    r.font.size=Pt(18); r.font.color.rgb=RGBColor(0xE8,0xF1,0xF9); r.font.italic=True
    _notes(s,timing); return s

def concept_understand(kicker,heading,definition,mechanism,image,hook,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker+"  ·  understand",heading)
    tw=6.2 if image else 11.9
    tf=_box(s,0.7,1.72,tw,5.05); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=definition; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=NAVY; p.space_after=Pt(10)
    p=tf.add_paragraph(); r=p.add_run(); r.text="HOW IT WORKS"; r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=BLUE; p.space_after=Pt(5)
    for m in mechanism:
        p=tf.add_paragraph(); r=p.add_run(); r.text="•  "+m; r.font.size=Pt(14); r.font.color.rgb=DARK; p.space_after=Pt(6)
    if image:
        path=os.path.join(IMG,image)
        if os.path.exists(path):
            pic=s.shapes.add_picture(path,Inches(6.95),Inches(1.9),width=Inches(6.05))
            band_top,band_h=1.9,4.55; h_in=pic.height/914400.0
            pic.top=Inches(band_top+max(0.0,(band_h-h_in)/2.0))
    _bar(s,TEAL,0,6.9,SW.inches,0.6)
    ft=_box(s,0.7,6.96,11.9,0.5); p=ft.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text="🧠  Memory hook:  "+hook; r.font.size=Pt(13.5); r.font.bold=True; r.font.color.rgb=WHITE
    _notes(s,timing); return s

def concept_apply(kicker,heading,example,trap,exam,keyterms):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker+"  ·  apply",heading)
    _card(s,0.7,1.66,12.0,1.48,BLUE ,BLUE_T ,"🌍  Real example",example)
    _card(s,0.7,3.26,12.0,1.30,CORAL,CORAL_T,"⚠️  Common trap — and the fix",trap)
    _card(s,0.7,4.68,12.0,1.92,TEAL ,TEAL_T ,"🎯  Exam-ready answer",exam)
    _card(s,0.7,6.72,12.0,0.60,AMBER,AMBER_T,"🔑  Key terms",keyterms,body_sz=11.5)
    return s

def add_activity(title,steps,expected,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,PALET); _bar(s,TEAL,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text="🛠  IN-CLASS ACTIVITY  ·  THINK–PAIR–SHARE"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=TEAL
    th=_box(s,0.7,0.66,12,0.9); r=th.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
    tf=_box(s,0.7,1.75,11.9,3.6); tf.word_wrap=True; first=True
    for text in steps:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text="•  "+text; r.font.size=Pt(15.5); r.font.color.rgb=DARK; p.space_after=Pt(7)
    _card(s,0.7,5.5,11.9,1.7,TEAL,WHITE,"✅  Expected answer",expected)
    _notes(s,timing); return s

def add_quiz(title,items,timing="",compact=False):
    s=prs.slides.add_slide(BLANK); _bg(s,LIGHT); _bar(s,AMBER,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text="CHECK FOR UNDERSTANDING"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=AMBER
    th=_box(s,0.7,0.66,12,0.9); r=th.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
    qz,az,sb=(14,12.5,5) if compact else (17,15,9)
    tf=_box(s,0.7,1.75,11.9,5.3); tf.word_wrap=True; first=True
    for text,kind in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=text
        if kind=="q": r.font.size=Pt(qz); r.font.bold=True; r.font.color.rgb=NAVY; p.space_before=Pt(sb)
        elif kind=="a": r.font.size=Pt(az); r.font.color.rgb=TEAL; r.font.bold=True; p.level=1
        else: r.font.size=Pt(az); r.font.color.rgb=DARK; r.font.italic=True; p.level=1; p.space_before=Pt(sb-3)
    _notes(s,timing); return s

def add_summary(kicker,takeaways,why,nxt,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,"Summary & Key Takeaways",kc=NAVY)
    tf=_box(s,0.7,1.7,12.0,2.7); tf.word_wrap=True; first=True
    for i,tk in enumerate(takeaways,1):
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=f"{i}.  "+tk; r.font.size=Pt(15); r.font.color.rgb=DARK; p.space_after=Pt(8)
    _card(s,0.7,4.55,12.0,1.35,TEAL,TEAL_T,"💡  Why this matters",why,body_sz=13.5)
    _card(s,0.7,6.02,12.0,1.15,BLUE,BLUE_T,"➡️  Next",nxt,body_sz=13.5)
    _notes(s,timing); return s

def add_solved_problem(kicker,title,scenario,steps,answer,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title,kc=AMBER)
    _card(s,0.7,1.58,12.0,1.0,AMBER,AMBER_T,"📋  Problem",scenario,body_sz=13)
    tf=_box(s,0.85,2.72,11.7,2.95); tf.word_wrap=True; first=True
    for i,st in enumerate(steps,1):
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=f"Step {i}.  "; r.font.size=Pt(12.5); r.font.bold=True; r.font.color.rgb=AMBER
        r2=p.add_run(); r2.text=st; r2.font.size=Pt(12.5); r2.font.color.rgb=DARK; p.space_after=Pt(4)
    _card(s,0.7,5.72,12.0,1.5,TEAL,TEAL_T,"✅  Worked answer",answer,body_sz=13)
    _notes(s,timing); return s

def add_cheatsheet(kicker,title,blocks):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    xs=[0.7,6.78]; ys=[1.72,3.36,5.0]; hh=1.55
    for idx,(label,body) in enumerate(blocks[:6]):
        _card(s,xs[idx%2],ys[idx//2],5.85,hh,BLUE,BLUE_T,label,body,body_sz=12)
    return s

def add_glossary(kicker,title,terms):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    half=(len(terms)+1)//2
    for x,group in [(0.7,terms[:half]),(6.78,terms[half:])]:
        tf=_box(s,x,1.75,5.85,5.3); tf.word_wrap=True; first=True
        for term,defn in group:
            p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            r=p.add_run(); r.text=term+" — "; r.font.size=Pt(11.5); r.font.bold=True; r.font.color.rgb=NAVY
            r2=p.add_run(); r2.text=defn; r2.font.size=Pt(11.5); r2.font.color.rgb=DARK; p.space_after=Pt(6)
    return s

def _add_page_numbers():
    total=len(prs.slides._sldIdLst)
    for i,s in enumerate(prs.slides,1):
        try:
            h=str(s.background.fill.fore_color.rgb); rr,gg,bb=int(h[0:2],16),int(h[2:4],16),int(h[4:6],16)
            dark=(0.299*rr+0.587*gg+0.114*bb)<128
        except Exception: dark=False
        col=RGBColor(0xE8,0xF0,0xF7) if dark else GREY
        tb=_box(s,12.35,0.30,0.9,0.35); p=tb.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
        r=p.add_run(); r.text=f"{i} / {total}"; r.font.size=Pt(10); r.font.color.rgb=col

# ============================================================
#                        BUILD
# ============================================================
add_title("Unit 4 — Database Normalization",
          "IT 220: Database Management System  ·  BIM 4th Semester  ·  Sessions S18–S21",
          "Running example: one bad ENROLL table, decomposed 1NF → 2NF → 3NF → BCNF → 4NF.  "
          "Self-contained slides — every concept complete on the slide face; exports to PDF with no information lost.")

add_outcomes("Unit 4 — Learning Outcomes","overview  ·  s18–s21",
    "By the end of this unit, you will be able to:",
    ["Identify insertion, deletion, and update anomalies in a bad design (S18)",
     "Define and detect functional dependencies (FDs) and the keys they imply (S18)",
     "Apply 1NF, 2NF, 3NF to decompose a table step by step, with justification (S19–S20)",
     "Define BCNF vs 3NF, and apply 4NF to multivalued dependencies (S21)",
     "Check a decomposition is lossless and dependency-preserving (S21)"],
    "Unit 2 designed the tables (ER→relational); Unit 3 queried them (algebra/calculus). Unit 4 makes the tables GOOD, so Unit 5's SQL runs on a clean schema.")

add_roadmap("Unit 4 — Roadmap","Where each session fits (S18–S21)",
    ["S18   Design guidelines · anomalies · functional dependencies",
     "S19   Normalization ladder · First Normal Form (1NF)",
     "S20   Second (2NF) & Third (3NF) Normal Forms",
     "S21   BCNF · MVD & 4NF · lossless join & dependency preservation"],
    ["Unit 5   SQL — build & query these clean tables",
     "Unit 6   Transactions & concurrency",
     "Unit 7   Advanced topics",
     "(Unit 4 threads ONE ENROLL table through every rung.)"])

# ============================ S18 ============================
add_divider("Session 18 · Unit 4","Design Guidelines & Functional Dependencies",
    "A college tracks enrollments in one giant spreadsheet. A student changes their phone — and we edit it in 11 rows. Miss one, and the database now believes two different numbers. What went wrong?",
    "OPENING HOOK [~5 min]. Agenda: what makes a schema good → the three anomalies → functional dependencies & keys.")

concept_understand("S18 · Concept 1 · [THEORY]","What Makes a Schema 'Good'",
    "Informal design guidelines: each relation should represent ONE clear thing; minimise redundancy (never store the same fact twice); avoid excessive nulls; and avoid designs that produce spurious tuples when tables are re-joined.",
    ["Don't cram student, course, and instructor facts into one table — that's three things in one.",
     "Storing a fact once (a single source of truth) is the core goal.",
     "The all-in-one ENROLL sheet breaks the first two guidelines badly.",
     "Good design trades a few more tables for far less redundancy."],
    None,"One table = one thing.",
    "~10 min. Contrast the all-in-one ENROLL sheet with a clean split; motivates the whole unit.")
concept_apply("S18 · Concept 1 · [THEORY]","What Makes a Schema 'Good'",
    "Every time a Nepali fintech or e-commerce team designs a new feature table, splitting entities cleanly up front — users vs orders vs merchants — is what prevents data corruption later. The cost of a bad first design compounds for years.",
    "\"Fewer tables = simpler = better.\" No — fewer tables usually means MORE redundancy and MORE anomalies. The right number of tables is 'one per real thing'.",
    "Good relational design: each relation models one entity with clear semantics; minimise redundant storage; minimise nulls; and avoid designs that generate spurious tuples on join. It trades a few more tables for far less redundancy and fewer anomalies.",
    "design guidelines · one-relation-one-entity · redundancy · nulls · spurious tuples")

concept_understand("S18 · Concept 2 · [THEORY]","Redundancy & the Three Anomalies",
    "An anomaly is an update problem caused by REDUNDANT storage. Redundancy in the ENROLL table produces three kinds of anomaly.",
    ["Insertion anomaly: can't add a new course (code, title) until some student enrolls in it.",
     "Deletion anomaly: if the last student drops a course, the course & instructor facts vanish.",
     "Update anomaly: an instructor's department changes → must fix every row; a partial fix = inconsistent data.",
     "All three come from the same root cause: the same fact stored many times."],
    "s18_anomalies.png","Redundancy → insert / delete / update anomalies.",
    "~13 min. Point each anomaly callout at specific rows; highlight the repeated cells.")
concept_apply("S18 · Concept 2 · [THEORY]","Redundancy & the Three Anomalies",
    "eSewa storing a merchant's commission rate copied into every transaction row: when the rate changes, some rows are updated and some aren't — now the database holds two 'truths' for one merchant. That's an update anomaly, live.",
    "\"Anomalies are rare edge cases.\" They're the everyday cost of redundancy — every duplicated fact is a future insert/delete/update anomaly waiting to corrupt data.",
    "An anomaly is an update problem caused by redundancy. Insertion: can't add a course until a student enrolls. Deletion: dropping the last student loses the course/instructor facts. Update: changing an instructor's department needs many-row edits, risking inconsistency. All stem from storing a fact more than once.",
    "anomaly · insertion / deletion / update anomaly · redundancy · inconsistency")

concept_understand("S18 · Concept 3 · [THEORY]","Functional Dependencies (FDs) & Keys",
    "A functional dependency X → Y means 'if you know X, you know exactly one Y' — X DETERMINES Y. A key is an attribute set that functionally determines EVERY attribute; a candidate key is a minimal such set.",
    ["FDs come from real-world rules, not from today's data, and are ONE-directional.",
     "In ENROLL: StudentRoll → StudentName; CourseCode → CourseTitle; InstructorID → InstructorName, InstructorDept.",
     "{StudentRoll, CourseCode} → Grade — the whole key determines the grade.",
     "The key here is {StudentRoll, CourseCode}. These FDs are the raw material for every normal form."],
    "s18_fd.png","X → Y = knowing X pins down one Y.",
    "~12 min. Draw the FD arrows; contrast NationalID→person (works) with Name→person (fails).")
concept_apply("S18 · Concept 3 · [THEORY]","Functional Dependencies (FDs) & Keys",
    "NationalID → citizen: the ID pins down exactly one person, so it can be a key. A NAME cannot — there are two people named 'Anish'. That's why real systems key on IDs, not names, and why FDs reveal the right key.",
    "\"X → Y means Y → X too.\" FDs are one-directional: Roll → Name, but Name does not determine Roll (duplicate names). Reversing an FD is the classic beginner error.",
    "A functional dependency X → Y holds if each X value is associated with exactly one Y value (X determines Y); FDs are one-directional and come from business rules. A key is a minimal attribute set that functionally determines all attributes; here the key is {StudentRoll, CourseCode}.",
    "functional dependency (X → Y) · one-directional · candidate key · determines-all-attributes")

add_activity("Find an FD in your marksheet",
    ["In pairs (3 min): look at your own college marksheet or result slip.",
     "Find one functional dependency that holds (X → Y: one X gives exactly one Y).",
     "Find one pair of columns where NEITHER determines the other.",
     "Share (2 min)."],
    "e.g. FD: SubjectCode → SubjectName (one code, one name). Non-FD: StudentName and Grade — a name doesn't determine a grade, and a grade doesn't determine a name. FDs come from the real-world rule, not the sample rows.",
    "ACTIVITY [~5 min].")
add_quiz("S18 — Quick Check",
    [("Q1.  Being unable to add a new course because no student has enrolled yet is which anomaly?","q"),
     ("a) update   b) ✅ insertion   c) deletion   d) join","a"),
     ("Q2.  CourseCode → CourseTitle means:","q"),
     ("a) titles determine codes   b) ✅ each code maps to exactly one title   c) codes are unique rows   d) nothing","a"),
     ("Discussion: an FD that holds in your marksheet, and a column pair where neither determines the other.","o")],
    "QUIZ [~5 min].")
add_summary("S18 · Summary  ·  [~2 min]",
    ["Good schemas minimise redundancy — one relation per real thing.",
     "Redundancy causes insertion, deletion, and update anomalies.",
     "Functional dependencies (X → Y) formalise the rules and reveal the key ({StudentRoll, CourseCode})."],
    "Spotting FDs and anomalies FIRST — before writing a line of code — is what prevents data corruption when a Nepali fintech or e-commerce team ships a new feature table.",
    "S19 — turning these FDs into formal normal forms, starting with First Normal Form (1NF).",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S19 ============================
add_divider("Session 19 · Unit 4","Normal Forms & First Normal Form (1NF)",
    "One cell in our enrollment sheet says StudentPhone = '9801, 9842, 9818'. Is that one value or three? Can you query 'who has 9842'? This single fat cell is why 1NF exists.",
    "OPENING HOOK [~5 min]. Agenda: the normalization ladder → 1NF → fixing the running table to 1NF.")

concept_understand("S19 · Concept 1 · [THEORY]","Normalization: the Ladder",
    "Normalization is a STEP-BY-STEP process of refining tables to reduce redundancy and anomalies, judged against the FDs and keys. The normal forms form a ladder — 1NF ⊂ 2NF ⊂ 3NF ⊂ BCNF ⊂ 4NF.",
    ["Each rung fixes a SPECIFIC class of anomaly.",
     "You climb one rung at a time — you don't fix everything in one action.",
     "The ladder is based on the FDs and the key you found in S18.",
     "Higher rungs are stricter: satisfying 3NF means you already satisfy 1NF and 2NF."],
    None,"Normalize one rung at a time.",
    "~13 min. Preview the ladder; each session climbs one or two rungs of the running example.")
concept_apply("S19 · Concept 1 · [THEORY]","Normalization: the Ladder",
    "A schema review on any product team is exactly this: 'is this table in 3NF? which dependency is still wrong?' — done incrementally, table by table, not as one giant rewrite.",
    "\"Normalization is one big action.\" It's incremental — you climb one normal form at a time, checking the specific anomaly each rung removes before moving up.",
    "Normalization is a stepwise process of decomposing relations to reduce redundancy and anomalies, guided by FDs and keys, through a ladder of normal forms (1NF → 2NF → 3NF → BCNF → 4NF); each rung removes a specific class of anomaly and each is stricter than the last.",
    "normalization · normal form · ladder (1NF→2NF→3NF→BCNF→4NF) · incremental")

concept_understand("S19 · Concept 2 · [THEORY]","First Normal Form (1NF)",
    "A relation is in 1NF if every attribute holds a single, ATOMIC value — no repeating groups, no multivalued cells, no nested tables. Each cell = one value.",
    ["Violations in ENROLL: StudentPhone = '9801,9842,9818' and repeating Course1, Course2, Course3 columns.",
     "Fix the phone: pull it into STUDENT_PHONE(StudentRoll, Phone) — one row per phone.",
     "Fix the courses: flatten into one row per (StudentRoll, CourseCode).",
     "Result: ENROLL_1NF with key {StudentRoll, CourseCode} — but redundancy still remains (sets up S20)."],
    "s19_1nf.png","1NF = one atomic value per cell.",
    "~12 min. Show the fat multivalued cell splitting into atomic rows.")
concept_apply("S19 · Concept 2 · [THEORY]","First Normal Form (1NF)",
    "A REST API or CSV export breaks the moment a field secretly holds a list. 1NF discipline is what keeps Daraz product attributes individually queryable and filterable — you can search 'phone = 9842' only if each phone is its own atomic value.",
    "\"1NF just means no duplicate rows.\" No — 1NF is about atomic, single-valued cells and no repeating groups; row uniqueness is a different idea (the key).",
    "A relation is in 1NF if every attribute value is atomic (single-valued) with no repeating groups or nested relations. Multivalued/list cells are split into separate rows or a separate table — e.g. StudentPhone becomes STUDENT_PHONE(StudentRoll, Phone).",
    "1NF · atomic value · repeating group · multivalued cell · one-value-per-cell")

add_solved_problem("S19 · Solved Problem","Bring ENROLL to 1NF",
    "Given ENROLL with StudentPhone = '9801, 9842, 9818' and repeating Course1/Course2/Course3 columns, "
    "rewrite it so it satisfies First Normal Form.",
    ["Spot the 1NF violations: a list in the StudentPhone cell, and repeating course columns (a repeating group).",
     "Multivalued phone → its own table, one atomic value per row: STUDENT_PHONE(StudentRoll, Phone).",
     "Repeating course columns → flatten to one row per (StudentRoll, CourseCode).",
     "The main table becomes ENROLL_1NF(StudentRoll, StudentName, CourseCode, CourseTitle, InstructorID, InstructorName, InstructorDept, Grade).",
     "Key = {StudentRoll, CourseCode}. Every cell is now atomic — but the row still repeats student/course/instructor facts."],
    "STUDENT_PHONE(StudentRoll, Phone) + ENROLL_1NF(StudentRoll, StudentName, CourseCode, CourseTitle, InstructorID, "
    "InstructorName, InstructorDept, Grade), key {StudentRoll, CourseCode}. 1NF removed the fat cells; the leftover "
    "redundancy is what 2NF and 3NF attack next.",
    "Let students find the two violations before revealing; stress 1NF fixes CELLS, not yet redundancy.")

add_quiz("S19 — Quick Check",
    [("Q1.  A column storing '9801, 9842, 9818' violates which normal form?","q"),
     ("✅ 1NF","a"),
     ("Q2.  1NF primarily requires:","q"),
     ("a) no redundancy   b) ✅ atomic, single-valued attributes   c) one table only   d) a foreign key","a"),
     ("Discussion: find a field on an online form that would be a 1NF violation if dumped into one column.","o")],
    "QUIZ [~5 min].")
add_summary("S19 · Summary  ·  [~2 min]",
    ["Normalization is a step-by-step ladder (1NF → 2NF → 3NF → BCNF → 4NF).",
     "1NF demands atomic, single-valued cells — no lists, no repeating groups.",
     "We split ENROLL into 1NF, but student/course/instructor facts still repeat."],
    "1NF discipline keeps fields queryable and filterable — every API, CSV export, and search filter depends on values being atomic, not secret lists.",
    "S20 — 2NF and 3NF, attacking the redundancy 1NF left behind (partial and transitive dependencies).",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S20 ============================
add_divider("Session 20 · Unit 4","Second (2NF) & Third (3NF) Normal Forms",
    "Our 1NF table has key {StudentRoll, CourseCode}. But StudentName depends only on StudentRoll — half the key. And InstructorDept depends on InstructorID, not the key at all. These two 'wrong dependencies' are exactly what 2NF and 3NF kill.",
    "OPENING HOOK [~5 min]. Agenda: partial dependency → 2NF → transitive dependency → 3NF → the final design.")

concept_understand("S20 · Concept 1 · [THEORY]","Partial Dependency & 2NF",
    "A relation is in 2NF if it is in 1NF AND no non-key attribute depends on only PART of a composite key — no partial dependency. This only matters when the key is composite.",
    ["In ENROLL_1NF (key {StudentRoll, CourseCode}): StudentRoll → StudentName and CourseCode → CourseTitle are partial.",
     "Fix by moving each partially-dependent attribute to a table keyed on the part it depends on.",
     "→ STUDENT(StudentRoll, StudentName), COURSE(CourseCode, CourseTitle, InstructorID, InstructorName, InstructorDept), ENROLL(StudentRoll, CourseCode, Grade).",
     "A table with a single-attribute key is automatically in 2NF."],
    "s20_2nf.png","2NF: depend on the WHOLE key, not half.",
    "~12 min. Colour the partial FDs red; split into three tables.")
concept_apply("S20 · Concept 1 · [THEORY]","Partial Dependency & 2NF",
    "A shared grocery bill where some items belong to only one roommate: you split those out so nobody pays for what isn't theirs. 2NF does the same — StudentName belongs with the student, CourseTitle with the course, not on the enrollment row.",
    "\"Every table must be forced to 2NF separately.\" A table whose key is a single attribute is already in 2NF — partial dependency is only possible against a COMPOSITE key.",
    "A relation is in 2NF if it is in 1NF and has no partial dependency — no non-key attribute depends on only part of a composite key. Fix by moving each partially-dependent attribute into a table keyed on the part it depends on (StudentRoll→StudentName gives a STUDENT table).",
    "2NF · partial dependency · composite key · non-key attribute")

concept_understand("S20 · Concept 2 · [THEORY]","Transitive Dependency & 3NF",
    "A relation is in 3NF if it is in 2NF AND no non-key attribute depends on ANOTHER non-key attribute — no transitive dependency (a key → non-key → non-key chain).",
    ["In the COURSE table: CourseCode → InstructorID → InstructorName, InstructorDept is transitive.",
     "The instructor's facts repeat for every course they teach — redundancy again.",
     "Fix: split into COURSE(CourseCode, CourseTitle, InstructorID) and INSTRUCTOR(InstructorID, InstructorName, InstructorDept).",
     "Now each fact is determined directly by a key, not through a middle non-key attribute."],
    "s20_3nf.png","3NF: no non-key → non-key chains.",
    "~12 min. Colour the transitive chain orange; split off INSTRUCTOR.")
concept_apply("S20 · Concept 2 · [THEORY]","Transitive Dependency & 3NF",
    "A Khalti Transaction table with MerchantID → MerchantCity → CityRegion: the region repeats in every transaction row until you break the chain into its own table. 3NF is the everyday fix that keeps lookup facts in one place.",
    "\"3NF means absolutely no redundancy ever.\" 3NF removes the common transitive redundancy, but a few odd cases survive — which is exactly why BCNF (S21) exists.",
    "A relation is in 3NF if it is in 2NF and has no transitive dependency — no non-key attribute determines another non-key attribute. Fix by moving the dependent non-key attributes into a table keyed on their determinant (InstructorID → its instructor facts).",
    "3NF · transitive dependency · non-key determinant · key→non-key→non-key chain")

concept_understand("S20 · Concept 3 · [THEORY]","The Fully Decomposed Design",
    "The running ENROLL table is now five clean tables in 3NF — and every S18 anomaly is gone.",
    ["STUDENT(StudentRoll, StudentName) · STUDENT_PHONE(StudentRoll, Phone).",
     "COURSE(CourseCode, CourseTitle, InstructorID) · INSTRUCTOR(InstructorID, InstructorName, InstructorDept).",
     "ENROLL(StudentRoll, CourseCode, Grade).",
     "Add a course with no students ✓; drop the last student, keep the course ✓; change an instructor's dept in ONE place ✓."],
    "s20_final.png","Each fact in exactly one place.",
    "~6 min. Re-show the S18 anomalies checklist, now all ticked.")
concept_apply("S20 · Concept 3 · [THEORY]","The Fully Decomposed Design",
    "This exact 1NF→2NF→3NF split is the everyday schema review on any Nepali product team — it's what keeps a user's profile, their orders, and merchant data from corrupting each other as the product grows.",
    "\"More tables makes it slower/harder.\" The five tables re-join losslessly on their shared keys; correctness comes first, and indexes/joins handle performance. A corrupt single-table design is far costlier.",
    "After decomposition the design is STUDENT, STUDENT_PHONE, COURSE, INSTRUCTOR, ENROLL — each fact stored once. Course facts live in COURSE (no insert/delete anomaly), instructor facts in INSTRUCTOR (a dept change is one row), and the tables re-join losslessly.",
    "decomposed design · one-fact-one-place · anomalies eliminated")

add_solved_problem("S20 · Solved Problem","Decompose ENROLL_1NF → 2NF → 3NF",
    "Starting from ENROLL_1NF (key {StudentRoll, CourseCode}), decompose it to 3NF, justifying each step.",
    ["List the offending FDs: StudentRoll → StudentName (partial); CourseCode → CourseTitle (partial); InstructorID → InstructorName, InstructorDept.",
     "2NF — remove partial deps: pull StudentRoll→StudentName into STUDENT; pull CourseCode→(CourseTitle, instructor facts) into COURSE; leave ENROLL(StudentRoll, CourseCode, Grade).",
     "Now inspect COURSE: CourseCode → InstructorID → InstructorName, InstructorDept is transitive.",
     "3NF — remove the transitive dep: split COURSE into COURSE(CourseCode, CourseTitle, InstructorID) and INSTRUCTOR(InstructorID, InstructorName, InstructorDept).",
     "Re-check: every non-key attribute now depends on the whole key and only the key."],
    "Five 3NF tables: STUDENT(Roll, Name) · STUDENT_PHONE(Roll, Phone) · COURSE(Course, Title, InstrID) · "
    "INSTRUCTOR(InstrID, Name, Dept) · ENROLL(Roll, Course, Grade). Partial deps removed at 2NF, the transitive "
    "dep removed at 3NF — all three anomalies gone.",
    "This is the single most exam-tested skill in Unit 4 — let students drive each split before revealing.")

add_quiz("S20 — Quick Check",
    [("Q1.  A non-key attribute depending on only part of a composite key is a:","q"),
     ("a) transitive dependency   b) ✅ partial dependency   c) multivalued dependency   d) trivial FD","a"),
     ("Q2.  CourseCode → InstructorID → InstructorDept is a ___ dependency, fixed by 3NF:","q"),
     ("✅ transitive","a"),
     ("Discussion: on the 5-table design, re-run 'instructor changes department' — how many rows change now, and why is that better?","o")],
    "QUIZ [~5 min].")
add_summary("S20 · Summary  ·  [~2 min]",
    ["2NF removes partial dependencies on a composite key (StudentRoll→StudentName).",
     "3NF removes transitive dependencies through non-key attributes (CourseCode→InstructorID→Dept).",
     "The running table is now anomaly-free across five clean tables."],
    "The 1NF→2NF→3NF split is the everyday schema review on a product team — it keeps a user's profile, orders, and merchant data from corrupting each other.",
    "S21 — the stricter BCNF, multivalued dependencies/4NF, and how to know a decomposition is SAFE.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S21 ============================
add_divider("Session 21 · Unit 4 (closes the unit)","BCNF · MVD & 4NF · Properties of Decomposition",
    "Our 3NF design looks clean — but here's a course where every instructor uses one fixed room, and we still get a weird repeating glitch. 3NF says it's fine; reality says it isn't. Welcome to BCNF — and to knowing when a split is actually safe.",
    "OPENING HOOK [~5 min]. Agenda: BCNF → MVD & 4NF → lossless join & dependency preservation → synthesis.")

concept_understand("S21 · Concept 1 · [THEORY]","Boyce–Codd Normal Form (BCNF)",
    "A relation is in BCNF if, for every non-trivial FD X → Y, X is a SUPERKEY. This is stricter than 3NF, which tolerates a non-superkey determinant in some prime-attribute cases.",
    ["Example: SECTION(CourseCode, Instructor, Room) where Instructor → Room, but Instructor is NOT a superkey.",
     "That is 3NF-OK yet BCNF-violating (a non-superkey determines something).",
     "Fix: split into ROOM_OF(Instructor, Room) and TEACHES(CourseCode, Instructor).",
     "3NF and BCNF coincide most of the time; they differ with overlapping candidate keys / a non-superkey determinant."],
    "s21_bcnf.png","BCNF: every determinant is a superkey.",
    "~11 min. Walk the Instructor→Room counterexample and its two-table fix.")
concept_apply("S21 · Concept 1 · [THEORY]","Boyce–Codd Normal Form (BCNF)",
    "BCNF is the strict inspector who flags the one corner the 'good-enough' 3NF building code lets slide. In the SECTION example, 3NF passes it but the room fact still repeats for every course an instructor teaches — BCNF removes that last redundancy.",
    "\"3NF and BCNF are the same.\" They're equal most of the time, but differ when a non-superkey determines another attribute (often with overlapping candidate keys) — that case is 3NF-OK but not BCNF.",
    "A relation is in BCNF if, for every non-trivial FD X → Y, X is a superkey. It is stricter than 3NF (which permits a non-superkey determinant of a prime attribute). They coincide unless a non-superkey determines something; then decompose so every determinant is a key.",
    "BCNF · superkey · determinant · stricter-than-3NF · overlapping candidate keys")

concept_understand("S21 · Concept 2 · [THEORY]","Multivalued Dependency (MVD) & 4NF",
    "A multivalued dependency X ↠ Y means: for each X there is a SET of independent Y values, independent of the other attributes. A relation is in 4NF if it is in BCNF and has no non-trivial MVD (unless X is a superkey).",
    ["MVDs cause a multiplying 'cartesian product' redundancy even when there is no FD problem.",
     "Example: STUDENT_ACTIVITY(StudentRoll, Course, Hobby) — courses and hobbies are independent.",
     "So every course pairs with every hobby (2 courses × 2 hobbies = 4 rows) — a row explosion.",
     "Fix: split into (StudentRoll, Course) and (StudentRoll, Hobby)."],
    "s21_mvd.png","4NF: don't cross-multiply independent multivalued facts.",
    "~12 min. Show the courses×hobbies explosion grid, then the clean split.")
concept_apply("S21 · Concept 2 · [THEORY]","Multivalued Dependency (MVD) & 4NF",
    "Storing a person's multiple PHONES and multiple EMAILS in one table is the classic 4NF violation — every phone gets paired with every email, so 2 phones × 3 emails = 6 meaningless rows. Split the two independent facts into separate tables.",
    "\"4NF is the same as 1NF — both are about multiple values.\" 1NF bans lists INSIDE a cell; 4NF bans two INDEPENDENT multivalued facts combined in one table (the cross-multiply).",
    "A multivalued dependency X ↠ Y holds when X determines a set of Y values independent of the rest of the tuple. A relation is in 4NF if it is in BCNF and has no non-trivial MVD (unless the determinant is a superkey). Fix by splitting the independent multivalued attributes into separate tables.",
    "multivalued dependency (X ↠ Y) · 4NF · independent multivalued facts · row explosion")

concept_understand("S21 · Concept 3 · [THEORY]","Lossless Join & Dependency Preservation",
    "A decomposition must be CORRECT, judged by two properties: it must be lossless, and ideally dependency-preserving.",
    ["Lossless (non-additive) join: re-joining the pieces gives back EXACTLY the original — no missing, no spurious rows.",
     "Rule: a binary split is lossless if the shared attribute is a SUPERKEY of at least one resulting table.",
     "Dependency preservation: every original FD can still be enforced WITHOUT re-joining tables.",
     "A split can be lossless but NOT dependency-preserving — sometimes you keep 3NF (not BCNF) to preserve a dependency."],
    "s21_lossless.png","Lossless = rebuild exactly; preserve = check FDs without re-joining.",
    "~12 min. Show a lossy split inventing spurious tuples vs a lossless one.")
concept_apply("S21 · Concept 3 · [THEORY]","Lossless Join & Dependency Preservation",
    "Before a bank or telecom database migration, the litmus test is exactly this: are tables in (at least) 3NF/BCNF, and was every historical split lossless — so no records were silently duplicated or lost in the move?",
    "\"Any split into smaller tables is safe.\" A careless split is LOSSY — re-joining invents spurious tuples. Losslessness must be verified (shared attribute is a superkey of one piece), never assumed.",
    "Two properties of a good decomposition: lossless (non-additive) join — the natural join of the pieces reproduces the original exactly (guaranteed if the common attribute is a superkey of one piece); and dependency preservation — all original FDs remain checkable without joining. A split must be lossless; dependency preservation is desirable and can conflict with BCNF.",
    "lossless (non-additive) join · spurious tuples · shared attribute = superkey · dependency preservation · BCNF-vs-3NF trade-off")

concept_understand("S21 · Concept 4 · [THEORY]","Unit 4 Synthesis — the whole ladder",
    "The complete toolkit: anomalies → FDs → 1NF → 2NF → 3NF → BCNF → 4NF, with the running ENROLL example at each rung, and a final gate: is the decomposition lossless and dependency-preserving?",
    ["1NF atomic values → 2NF no partial dep → 3NF no transitive dep.",
     "BCNF every determinant a superkey → 4NF no independent MVD.",
     "Each rung removes a specific class of anomaly; higher rungs are stricter.",
     "Final gate on every split: lossless join (mandatory) + dependency preservation (desirable)."],
    "s21_ladder.png","Anomalies → FDs → climb the ladder → check it's lossless.",
    "~7 min. This is the slide to photograph — the whole unit as one ladder into a safe design.")
concept_apply("S21 · Concept 4 · [THEORY]","Unit 4 Synthesis — the whole ladder",
    "Given any messy table in an exam or a real design review, the procedure is always the same: write the FDs and key, climb 1NF→2NF→3NF→BCNF→4NF removing one anomaly class per rung, and verify each split is lossless. That's a complete, repeatable method.",
    "\"Normalize as far as possible, always to the top.\" Not always — you sometimes stop at 3NF to keep a decomposition dependency-preserving. The goal is a correct, maintainable design, not the highest rung for its own sake.",
    "Normalization procedure: detect anomalies → write FDs and the key → 1NF (atomic) → 2NF (no partial dep) → 3NF (no transitive dep) → BCNF (every determinant a superkey) → 4NF (no independent MVD), checking each decomposition is lossless and ideally dependency-preserving.",
    "normal-form ladder · lossless & dependency-preserving gate · anomaly-driven design")

add_solved_problem("S21 · Capstone Solved Problem","Decompose ENROLL all the way to BCNF/4NF",
    "Take the original ENROLL (with multivalued StudentPhone) and normalize it fully, checking correctness.",
    ["FDs & key: Roll→Name; Course→Title; InstrID→InstrName,InstrDept; {Roll,Course}→Grade; MVD Roll ↠ Phone. Key = {Roll, Course}.",
     "1NF: pull multivalued Phone → STUDENT_PHONE(Roll, Phone); flatten repeating courses.",
     "2NF: remove partial deps → STUDENT(Roll, Name), COURSE(Course, Title, InstrID, InstrName, InstrDept), ENROLL(Roll, Course, Grade).",
     "3NF: remove transitive dep in COURSE → COURSE(Course, Title, InstrID) + INSTRUCTOR(InstrID, InstrName, InstrDept).",
     "BCNF check: every determinant (Roll, Course, InstrID) is now a key of its table → already BCNF. 4NF: Phone is isolated in its own table → no independent-MVD explosion.",
     "Correctness: each split shares a superkey (Roll, Course, InstrID) → LOSSLESS; all FDs are checkable in one table each → dependency-preserving."],
    "Final design (BCNF & 4NF): STUDENT(Roll, Name) · STUDENT_PHONE(Roll, Phone) · COURSE(Course, Title, InstrID) · "
    "INSTRUCTOR(InstrID, Name, Dept) · ENROLL(Roll, Course, Grade). Lossless (shared superkeys) and "
    "dependency-preserving — the exam Part-C answer in full.",
    "This is exactly the consolidated end-of-unit applied task; pause and let students drive every rung.")

add_quiz("S21 — Quick Check",
    [("Q1.  A table is in 3NF but a non-superkey determines another attribute. It violates:","q"),
     ("a) 1NF   b) 2NF   c) ✅ BCNF   d) nothing","a"),
     ("Q2.  Re-joining decomposed tables produces extra, incorrect rows. The decomposition is:","q"),
     ("a) lossless   b) ✅ lossy / has spurious tuples   c) dependency-preserving   d) in 4NF","a"),
     ("Discussion: a real 'independent multivalued facts' case (phones and skills) and how you'd 4NF-split it.","o")],
    "QUIZ [~5 min].")
add_summary("S21 · Summary  ·  [~2 min]",
    ["BCNF tightens 3NF so every determinant is a superkey.",
     "4NF removes independent multivalued redundancy (don't cross-multiply independent multivalued facts).",
     "A good decomposition must be lossless and ideally dependency-preserving."],
    "When a bank or telecom audits its database before a migration, the test is exactly this — are tables in 3NF/BCNF, and was every historical split lossless so no records were silently duplicated or lost?",
    "Unit 5 — SQL: building and querying these clean, normalized tables in practice.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============= END-OF-UNIT REVISION AIDS =============
add_cheatsheet("Unit 4 · Cheat Sheet","One-page revision reference",
    [("The 3 anomalies","From redundancy: INSERTION (can't add w/o a related row) · DELETION (lose facts) · UPDATE (fix many rows)."),
     ("FD & key","X → Y = X determines exactly one Y (one-directional). Key = minimal set that determines all attributes."),
     ("1NF","Atomic, single-valued cells; no repeating groups or lists in a cell."),
     ("2NF / 3NF","2NF = 1NF + no PARTIAL dep on a composite key. 3NF = 2NF + no TRANSITIVE dep (non-key→non-key)."),
     ("BCNF / 4NF","BCNF = every determinant is a superkey (stricter 3NF). 4NF = BCNF + no independent MVD (X ↠ Y)."),
     ("Safe decomposition","LOSSLESS: shared attribute is a superkey of one piece (else spurious tuples). Plus DEPENDENCY-PRESERVING.")])

add_glossary("Unit 4 · Glossary","Key terms — quick reference",
    [("Anomaly","update problem caused by redundant storage."),
     ("Insertion / deletion / update","can't add / lose facts / must fix many rows."),
     ("Functional dependency","X → Y: X determines exactly one Y."),
     ("Candidate key","minimal attribute set determining all attributes."),
     ("Superkey","any attribute set that determines all attributes."),
     ("Partial dependency","non-key depends on part of a composite key."),
     ("Transitive dependency","non-key depends on another non-key."),
     ("1NF","atomic values; no repeating groups."),
     ("2NF","1NF + no partial dependency."),
     ("3NF","2NF + no transitive dependency."),
     ("BCNF","every determinant is a superkey."),
     ("Multivalued dependency (↠)","X gives a set of independent Y values."),
     ("4NF","BCNF + no non-trivial MVD."),
     ("Lossless join","re-join reproduces the original exactly."),
     ("Spurious tuples","extra wrong rows from a lossy split."),
     ("Dependency preservation","every FD checkable without re-joining.")])

# ============= CONSOLIDATED END-OF-UNIT QUIZ =============
add_divider("Unit 4 · Revision","Consolidated end-of-unit quiz (S18–S21)",
    "Ten MCQs across the whole unit (answers shown), then short-answer and a full applied decomposition task to work from the concept slides and Unit4_material.md.",
    "Use for revision/assessment after S21.")
add_quiz("Part A — Multiple choice (answers shown)",
    [("1.  Can't record a new course until a student enrolls   →  ✅ insertion anomaly","a"),
     ("2.  InstructorID → InstructorName is a   →  ✅ functional dependency","a"),
     ("3.  A candidate key is a minimal set that   →  ✅ functionally determines all attributes","a"),
     ("4.  A cell storing '9801,9842,9818' violates   →  ✅ 1NF","a"),
     ("5.  2NF is only at risk with a   →  ✅ composite / multi-attribute key","a"),
     ("6.  CourseCode → InstructorID → InstructorDept is a   →  ✅ transitive dependency","a"),
     ("7.  BCNF requires every non-trivial determinant X to be a   →  ✅ superkey","a"),
     ("8.  Independent courses & hobbies crammed in one table violate   →  ✅ 4NF","a"),
     ("9.  A decomposition producing spurious tuples on re-join is   →  ✅ lossy","a"),
     ("10.  Property letting every FD be checked without re-joining   →  ✅ dependency preservation","a")],
    "Consolidated quiz Part A.",compact=True)
add_quiz("Parts B–C — short answer & applied decomposition",
    [("Part B — Short answer","q"),
     ("1. List the three anomalies with a one-line example of each.   2. Partial vs transitive dependency, and the NF that removes each.","o"),
     ("3. Why can a decomposition be lossless but NOT dependency-preserving? State the BCNF-vs-3NF trade-off.","o"),
     ("Part C — Applied decomposition (the running ENROLL table)","q"),
     ("1. Write the FDs and any MVD.   2. Identify the primary key.","o"),
     ("3. Decompose step by step to BCNF/4NF, showing the table set at each stage (1NF→2NF→3NF→BCNF/4NF).","o"),
     ("4. Confirm the final decomposition is lossless, and state whether it is dependency-preserving.","o")],
    "Consolidated quiz Parts B–C. Full worked answer = the S21 capstone; details in Unit4_material.md.",compact=True)

# ---- close ----
add_title("End of Unit 4  ·  IT 220",
          "S18–S21 complete: anomalies & functional dependencies → 1NF → 2NF → 3NF → BCNF → 4NF, plus lossless join & dependency preservation — one ENROLL table threaded through every rung",
          "Built to COURSE_MATERIAL_STANDARD.md · self-contained slides · exports cleanly to PDF · "
          "Next unit (Unit 5): SQL — building and querying these clean tables.")

_add_page_numbers()
out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"IT220_Unit4.pptx")
prs.save(out)
print("Saved",out,"with",len(prs.slides._sldIdLst),"slides")
