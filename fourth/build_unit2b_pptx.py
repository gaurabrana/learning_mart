#!/usr/bin/env python3
"""IT220 Unit 2 deck — S9–S12 batch (ER diagrams, ternary, specialization/
generalization, ER→table mapping). Built to COURSE_MATERIAL_STANDARD.md.

Same self-contained / PDF-safe style as build_unit2_pptx.py: every concept is a
two-slide pair (understand + apply); real example / trap / exam-answer / key terms
live on the slide face; notes hold only timing cues. Content distilled from
Unit2b_material.md. Run: python3 build_unit2b_pptx.py -> IT220_Unit2b.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "images")

NAVY=RGBColor(0x0C,0x2B,0x4A); BLUE=RGBColor(0x18,0x5F,0xA5); TEAL=RGBColor(0x0F,0x6E,0x56)
AMBER=RGBColor(0x85,0x4F,0x0B); CORAL=RGBColor(0x99,0x3C,0x1D); LIGHT=RGBColor(0xF2,0xF5,0xF8)
PALET=RGBColor(0xE1,0xF5,0xEE); WHITE=RGBColor(0xFF,0xFF,0xFF); DARK=RGBColor(0x1A,0x1A,0x1A)
GREY=RGBColor(0x55,0x60,0x6B)
BLUE_T=RGBColor(0xE6,0xF1,0xFB); CORAL_T=RGBColor(0xFA,0xEC,0xE7)
TEAL_T=RGBColor(0xE1,0xF5,0xEE); AMBER_T=RGBColor(0xFA,0xEE,0xDA)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

# ---------- low-level ----------
def _bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def _box(s,l,t,w,h): return s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)).text_frame
def _bar(s,c,l,t,w,h):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=c; sp.line.fill.background(); sp.shadow.inherit=False
def _notes(s,t):
    if t: s.notes_slide.notes_text_frame.text=t
def _header(s,kicker,header,kc=BLUE):
    _bar(s,kc,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=kc
    th=_box(s,0.7,0.66,12,1.0); th.word_wrap=True; r=th.paragraphs[0].add_run(); r.text=header
    r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
def _card(s,l,t,w,h,accent,fill,label,body,body_sz=14):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    sp.line.color.rgb=accent; sp.line.width=Pt(1); sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP
    tf.margin_left=Inches(0.22); tf.margin_right=Inches(0.20); tf.margin_top=Inches(0.10); tf.margin_bottom=Inches(0.08)
    p=tf.paragraphs[0]; r=p.add_run(); r.text=label; r.font.size=Pt(12.5); r.font.bold=True; r.font.color.rgb=accent
    p2=tf.add_paragraph(); p2.space_before=Pt(3); r2=p2.add_run(); r2.text=body; r2.font.size=Pt(body_sz); r2.font.color.rgb=DARK
    return sp

# ---------- slide builders ----------
def add_title(title,subtitle,footer):
    s=prs.slides.add_slide(BLANK); _bg(s,NAVY); _bar(s,AMBER,0.9,3.05,1.6,0.10)
    tf=_box(s,0.9,2.0,11.5,1.0); tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(34); r.font.bold=True; r.font.color.rgb=WHITE
    t2=_box(s,0.9,3.25,11.5,1.2); t2.word_wrap=True
    r=t2.paragraphs[0].add_run(); r.text=subtitle; r.font.size=Pt(18); r.font.color.rgb=RGBColor(0xCF,0xDD,0xEA)
    t3=_box(s,0.9,6.5,11.5,0.7); t3.word_wrap=True; r=t3.paragraphs[0].add_run(); r.text=footer
    r.font.size=Pt(13); r.font.color.rgb=RGBColor(0x9A,0xB2,0xC6); return s

def add_outcomes(header,kicker,intro,items,nxt):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,header)
    tf=_box(s,0.7,1.75,12.0,4.6); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=intro; r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=NAVY; p.space_after=Pt(8)
    for it in items:
        p=tf.add_paragraph(); r=p.add_run(); r.text="•  "+it; r.font.size=Pt(14.5); r.font.color.rgb=DARK; p.space_after=Pt(5)
    _card(s,0.7,6.35,12.0,0.85,BLUE,BLUE_T,"↩️ Builds on",nxt,body_sz=13)
    return s

def add_roadmap(kicker,title,done,todo):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    _card(s,0.7,1.8,5.85,4.55,TEAL,TEAL_T,"✅  THIS BATCH  ·  S9–S12","\n".join(done),body_sz=14)
    _card(s,6.78,1.8,5.85,4.55,GREY,LIGHT,"✔️  ALREADY COVERED  ·  S5–S8","\n".join(todo),body_sz=14)
    return s

def add_divider(kicker,title,hook,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,BLUE)
    tk=_box(s,0.9,1.9,11.5,0.6); r=tk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=RGBColor(0xBF,0xDC,0xF2)
    tt=_box(s,0.9,2.5,11.5,1.4); tt.word_wrap=True; r=tt.paragraphs[0].add_run(); r.text=title
    r.font.size=Pt(31); r.font.bold=True; r.font.color.rgb=WHITE
    _bar(s,AMBER,0.9,4.05,1.4,0.09)
    th=_box(s,0.9,4.25,11.5,2.4); th.word_wrap=True; p=th.paragraphs[0]; r=p.add_run(); r.text="🎣  "+hook
    r.font.size=Pt(18); r.font.color.rgb=RGBColor(0xE8,0xF1,0xF9); r.font.italic=True
    _notes(s,timing); return s

def add_reference(kicker,title,intro,image,timing=""):
    """Full-width reference slide (used for the notation legend)."""
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    tb=_box(s,0.7,1.62,11.9,0.7); tb.word_wrap=True; r=tb.paragraphs[0].add_run(); r.text=intro
    r.font.size=Pt(14); r.font.color.rgb=GREY
    path=os.path.join(IMG,image)
    if os.path.exists(path):
        pic=s.shapes.add_picture(path,Inches(3.0),Inches(2.4),height=Inches(4.7))
        pic.left=Inches(max(0.2,(SW.inches-pic.width/914400.0)/2.0))
    _notes(s,timing); return s

def concept_understand(kicker,heading,definition,mechanism,image,hook,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker+"  ·  understand",heading)
    tw=6.2 if image else 11.9
    tf=_box(s,0.7,1.72,tw,5.05); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=definition; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=NAVY; p.space_after=Pt(10)
    p=tf.add_paragraph(); r=p.add_run(); r.text="HOW IT WORKS"; r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=BLUE; p.space_after=Pt(5)
    for m in mechanism:
        p=tf.add_paragraph(); r=p.add_run(); r.text="•  "+m; r.font.size=Pt(14); r.font.color.rgb=DARK; p.space_after=Pt(6)
    if image:
        path=os.path.join(IMG,image)
        if os.path.exists(path):
            pic=s.shapes.add_picture(path,Inches(6.95),Inches(1.9),width=Inches(6.05))
            band_top,band_h=1.9,4.55; h_in=pic.height/914400.0
            pic.top=Inches(band_top+max(0.0,(band_h-h_in)/2.0))
    _bar(s,TEAL,0,6.9,SW.inches,0.6)
    ft=_box(s,0.7,6.96,11.9,0.5); p=ft.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text="🧠  Memory hook:  "+hook; r.font.size=Pt(13.5); r.font.bold=True; r.font.color.rgb=WHITE
    _notes(s,timing); return s

def concept_apply(kicker,heading,example,trap,exam,keyterms):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker+"  ·  apply",heading)
    _card(s,0.7,1.66,12.0,1.48,BLUE ,BLUE_T ,"🌍  Real example",example)
    _card(s,0.7,3.26,12.0,1.30,CORAL,CORAL_T,"⚠️  Common trap — and the fix",trap)
    _card(s,0.7,4.68,12.0,1.92,TEAL ,TEAL_T ,"🎯  Exam-ready answer",exam)
    _card(s,0.7,6.72,12.0,0.60,AMBER,AMBER_T,"🔑  Key terms",keyterms,body_sz=11.5)
    return s

def add_activity(title,steps,expected,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,PALET); _bar(s,TEAL,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text="🛠  IN-CLASS ACTIVITY  ·  THINK–PAIR–SHARE"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=TEAL
    th=_box(s,0.7,0.66,12,0.9); r=th.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
    tf=_box(s,0.7,1.75,11.9,3.6); tf.word_wrap=True; first=True
    for text in steps:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text="•  "+text; r.font.size=Pt(15.5); r.font.color.rgb=DARK; p.space_after=Pt(7)
    _card(s,0.7,5.5,11.9,1.7,TEAL,WHITE,"✅  Expected answer",expected)
    _notes(s,timing); return s

def add_quiz(title,items,timing="",compact=False):
    s=prs.slides.add_slide(BLANK); _bg(s,LIGHT); _bar(s,AMBER,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text="CHECK FOR UNDERSTANDING"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=AMBER
    th=_box(s,0.7,0.66,12,0.9); r=th.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
    qz,az,sb=(14,12.5,5) if compact else (17,15,9)
    tf=_box(s,0.7,1.75,11.9,5.3); tf.word_wrap=True; first=True
    for text,kind in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=text
        if kind=="q": r.font.size=Pt(qz); r.font.bold=True; r.font.color.rgb=NAVY; p.space_before=Pt(sb)
        elif kind=="a": r.font.size=Pt(az); r.font.color.rgb=TEAL; r.font.bold=True; p.level=1
        else: r.font.size=Pt(az); r.font.color.rgb=DARK; r.font.italic=True; p.level=1; p.space_before=Pt(sb-3)
    _notes(s,timing); return s

def add_summary(kicker,takeaways,why,nxt,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,"Summary & Key Takeaways",kc=NAVY)
    tf=_box(s,0.7,1.7,12.0,2.7); tf.word_wrap=True; first=True
    for i,tk in enumerate(takeaways,1):
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=f"{i}.  "+tk; r.font.size=Pt(15); r.font.color.rgb=DARK; p.space_after=Pt(8)
    _card(s,0.7,4.55,12.0,1.35,TEAL,TEAL_T,"💡  Why this matters",why,body_sz=13.5)
    _card(s,0.7,6.02,12.0,1.15,BLUE,BLUE_T,"➡️  Next",nxt,body_sz=13.5)
    _notes(s,timing); return s

def add_solved_problem(kicker,title,scenario,steps,answer,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title,kc=AMBER)
    _card(s,0.7,1.58,12.0,1.0,AMBER,AMBER_T,"📋  Problem",scenario,body_sz=13)
    tf=_box(s,0.85,2.72,11.7,2.95); tf.word_wrap=True; first=True
    for i,st in enumerate(steps,1):
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=f"Step {i}.  "; r.font.size=Pt(12.5); r.font.bold=True; r.font.color.rgb=AMBER
        r2=p.add_run(); r2.text=st; r2.font.size=Pt(12.5); r2.font.color.rgb=DARK; p.space_after=Pt(4)
    _card(s,0.7,5.72,12.0,1.5,TEAL,TEAL_T,"✅  Worked answer",answer,body_sz=13)
    _notes(s,timing); return s

def add_cheatsheet(kicker,title,blocks):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    xs=[0.7,6.78]; ys=[1.72,3.36,5.0]; hh=1.55
    for idx,(label,body) in enumerate(blocks[:6]):
        _card(s,xs[idx%2],ys[idx//2],5.85,hh,BLUE,BLUE_T,label,body,body_sz=12)
    return s

def add_glossary(kicker,title,terms):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    half=(len(terms)+1)//2
    for x,group in [(0.7,terms[:half]),(6.78,terms[half:])]:
        tf=_box(s,x,1.75,5.85,5.3); tf.word_wrap=True; first=True
        for term,defn in group:
            p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            r=p.add_run(); r.text=term+" — "; r.font.size=Pt(11.5); r.font.bold=True; r.font.color.rgb=NAVY
            r2=p.add_run(); r2.text=defn; r2.font.size=Pt(11.5); r2.font.color.rgb=DARK; p.space_after=Pt(6)
    return s

def _add_page_numbers():
    total=len(prs.slides._sldIdLst)
    for i,s in enumerate(prs.slides,1):
        try:
            h=str(s.background.fill.fore_color.rgb); rr,gg,bb=int(h[0:2],16),int(h[2:4],16),int(h[4:6],16)
            dark=(0.299*rr+0.587*gg+0.114*bb)<128
        except Exception: dark=False
        col=RGBColor(0xE8,0xF0,0xF7) if dark else GREY
        tb=_box(s,12.35,0.30,0.9,0.35); p=tb.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
        r=p.add_run(); r.text=f"{i} / {total}"; r.font.size=Pt(10); r.font.color.rgb=col

# ============================================================
#                        BUILD
# ============================================================
add_title("Unit 2 — Data Modelling (ER Model)  ·  S9–S12",
          "IT 220: Database Management System  ·  BIM 4th Semester  ·  Chen / EER notation",
          "Sessions S9–S12 — closes Unit 2  ·  50 min each  ·  Nepal / South Asia context  ·  "
          "Self-contained slides: every concept is complete on the slide face — exports to PDF with no information lost.")

add_outcomes("Unit 2 — Outcomes for this batch","overview  ·  s9–s12 batch",
    "By the end of S9–S12, you will be able to:",
    ["Draw clean, correctly-named ER diagrams and avoid common design pitfalls (S9)",
     "Handle relationships of degree higher than two — ternary and beyond (S10)",
     "Apply specialization/generalization with disjoint/overlap and total/partial constraints (S11)",
     "Convert (map) a complete ER diagram into relational tables using the standard algorithm (S12)"],
    "This batch assumes S5–S8: conceptual models, entities/attributes/keys, relationships & constraints, and weak entities.")

add_roadmap("Unit 2 — Roadmap","Where each session fits (S5–S12)",
    ["S9    Full ER diagrams · naming · design pitfalls",
     "S10  Higher-degree (ternary) relationships",
     "S11  Specialization & generalization",
     "S12  Mapping ER → relational tables  (closes Unit 2)"],
    ["S5   Conceptual models & the design process",
     "S6   Entity types/sets, attributes, keys",
     "S7   Relationships, roles, cardinality, participation",
     "S8   Weak entities, identifying rel., partial key"])

# ============================ S9 ============================
add_divider("Session 9 · Unit 2","ER Diagrams, Naming Conventions & Design Issues",
    "Here are two ER diagrams of the same college system — one is clear, one is a spaghetti mess. What makes the difference? Conventions and good design choices, not just correct symbols.",
    "OPENING HOOK [~5 min]. Agenda: notation legend → naming conventions → attribute-vs-entity design decisions.")

add_reference("S9 · Concept 1 · [THEORY]","ER Notation — one consolidated legend",
    "The complete Chen/EER symbol set from S5–S8 in one place. Read any diagram by naming each shape aloud — the habit that catches most errors before they reach the tables.",
    "s9_legend.png",
    "~10 min. This is the alphabet of ER — keep it visible while reading later diagrams.")

concept_understand("S9 · Concept 2 · [THEORY]","Naming Conventions",
    "Good names make a diagram self-explanatory: entity types are singular nouns, relationships are verbs, so a fragment reads left-to-right as a sentence.",
    ["Entity types = singular nouns (STUDENT, not STUDENTS) — a rectangle is the type, one specimen of the kind.",
     "Relationships = verbs (ENROLLS, TEACHES): STUDENT — ENROLLS — COURSE reads as a sentence.",
     "Use consistent case and meaningful names; never TBL1 / REL2.",
     "Names become table/column names after mapping — bad names breed miscommunication and errors."],
    "s9_naming.png","Entities are singular nouns; relationships are verbs.",
    "~12 min. Contrast REL1 vs TEACHES — which can a registrar verify?")
concept_apply("S9 · Concept 2 · [THEORY]","Naming Conventions",
    "A relationship named REL1 tells the registrar nothing; renamed TEACHES (TEACHER—TEACHES—COURSE) it reads like plain English and the registrar can confirm or correct it on the spot.",
    "\"Naming doesn't matter, it's just a diagram.\" Wrong — meaningless names (E1, R3) hide what each construct is, so mapping errors go unnoticed. Singular-noun entities + verb relationships make errors visible.",
    "Entity types are singular nouns (STUDENT); relationship types are verbs (ENROLLS) so the diagram reads as a sentence; names must be consistent and meaningful (never TBL1/REL2). Good names let stakeholders verify the model and become clean table/column names after mapping.",
    "singular-noun entity  ·  verb relationship  ·  reads-as-a-sentence  ·  meaningful / consistent naming")

concept_understand("S9 · Concept 3 · [THEORY]","Design Issue: Attribute vs Entity",
    "The hard modelling choices are about WHAT should be WHAT. Key rule: if a concept has its own attributes or its own relationships, model it as an entity; otherwise an attribute is enough.",
    ["'Colour' of a product is just a value → attribute.",
     "'Category' on Daraz has a name, a parent, an icon, and groups many products → its own ENTITY.",
     "Attribute vs relationship: a connection between two things is a relationship; a plain property is an attribute.",
     "Beware over-modelling — more entities is not 'better'; model exactly what the requirements need."],
    "s9_attr_vs_entity.png","Own attributes/relationships → entity; plain value → attribute.",
    "~13 min. Decision flowchart: does 'X' have its own attributes/relationships?")
concept_apply("S9 · Concept 3 · [THEORY]","Design Issue: Attribute vs Entity",
    "On Daraz, 'category' looks like a text attribute of PRODUCT — but because categories have their own data (name, parent, icon) and group many products, Daraz models CATEGORY as a separate entity. 'Colour', with no data of its own, stays an attribute.",
    "\"More entities = better design.\" No — over-modelling adds needless complexity. Model something as its own entity only when it genuinely carries its own attributes or relationships.",
    "Model a concept as its own entity when it has its own attributes and/or participates in its own relationships (e.g. CATEGORY, MERCHANT); keep it as an attribute when it is a plain value with no data of its own (e.g. colour, status). Avoid over-modelling.",
    "attribute-vs-entity decision  ·  attribute-vs-relationship  ·  over-modelling  ·  'has its own attributes/relationships' test")

add_activity("Attribute or entity?",
    ["In pairs (3 min): for a Foodmandu-style app, classify each of — cuisine, restaurant, delivery-address, rating.",
     "Decide attribute vs its own entity using the 'own attributes/relationships' test.",
     "Justify each choice in one sentence.",
     "Share (2 min)."],
    "restaurant = entity (has name, location, menu, relationships); delivery-address = entity or multivalued composite attribute of CUSTOMER; cuisine = attribute (or a small entity if it groups data); rating = attribute of a review.",
    "ACTIVITY [~5 min].")
add_quiz("S9 — Quick Check",
    [("Q1.  Best name for the entity of course-takers:","q"),
     ("a) ENROLLS   b) STUDENTS   c) ✅ STUDENT   d) TBL1","a"),
     ("Q2.  Model X as its own entity (not an attribute) when:","q"),
     ("✅ it has its own attributes / relationships","a"),
     ("Discussion: for an eSewa-like app, is 'merchant' an attribute of TRANSACTION or its own entity? Justify.","o")],
    "QUIZ [~5 min].")
add_summary("S9 · Summary  ·  [~2 min]",
    ["One consistent notation legend is the alphabet of ER — name every shape aloud to read a diagram.",
     "Entity types = singular nouns; relationship types = verbs (the fragment reads as a sentence).",
     "Model a concept as an entity only when it has its own attributes/relationships; don't over-model."],
    "Clean, well-named ER diagrams are the shared language between analysts, developers, and DBAs in every Nepali IT company; the attribute-vs-entity call decides how many tables you get.",
    "S10 — when two entities aren't enough: higher-degree (ternary) relationships.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S10 ============================
add_divider("Session 10 · Unit 2","Relationships of Degree Higher Than Two",
    "A supplier supplies a product for a project. Three things in one fact — can a normal two-sided arrow capture that?",
    "OPENING HOOK [~5 min]. Agenda: degree revisited / ternary → ternary ≠ 3 binaries → associative entities.")

concept_understand("S10 · Concept 1 · [THEORY]","Degree Revisited: Binary, Ternary, n-ary",
    "The degree of a relationship is the number of entity types it connects (from S7). Binary = 2 (most common), ternary = 3, n-ary = n. A ternary captures a single fact involving three things at once.",
    ["SUPPLIES(SUPPLIER, PART, PROJECT) = 'supplier S supplied part P FOR project J' — one three-way fact.",
     "Drawn as ONE diamond touching three rectangles.",
     "A timetable fact TEACHES(TEACHER, COURSE, ROOM/TIMESLOT) is another natural ternary.",
     "Most real relationships are binary; ternary is occasional, n>3 is rare."],
    "s10_ternary.png","Degree = how many entity types the diamond touches.",
    "~10 min. Draw one diamond connected to three rectangles.")
concept_apply("S10 · Concept 1 · [THEORY]","Degree Revisited: Binary, Ternary, n-ary",
    "A logistics/procurement system records 'this supplier delivered this part for this project' — a single fact about three entities at once. That is a ternary relationship; splitting it apart loses the point.",
    "\"Higher-degree relationships are everywhere.\" Actually most designs are binary; reach for ternary only when a fact genuinely ties three entity types together in one instance.",
    "The degree of a relationship is the number of entity types participating. Binary = 2, ternary = 3, n-ary = n. A ternary (e.g. SUPPLIES(SUPPLIER, PART, PROJECT)) associates three entity types in one fact, drawn as one diamond connected to three rectangles.",
    "degree  ·  binary (2)  ·  ternary (3)  ·  n-ary  ·  one diamond, three rectangles")

concept_understand("S10 · Concept 2 · [THEORY]","Why a Ternary Is NOT Three Binaries",
    "Replacing a ternary with three binary relationships often LOSES information, because the binaries record pairwise facts but not the combination of all three.",
    ["'S supplies P' + 'P used in J' + 'S works on J' are three true pairwise facts…",
     "…but they still do NOT say S supplied P specifically TO J.",
     "Only the ternary SUPPLIES(S, P, J) records the three-way truth.",
     "So three binaries ≠ a ternary whenever the meaning lives in the combination."],
    "s10_ternary_vs_binary.png","Three binaries lose the combined fact.",
    "~13 min. Walk the supplier/part/project counter-example slowly.")
concept_apply("S10 · Concept 2 · [THEORY]","Why a Ternary Is NOT Three Binaries",
    "In a delivery system, 'customer–restaurant', 'restaurant–rider', 'rider–customer' as three binaries can't reconstruct WHICH rider carried WHICH order from WHICH restaurant to WHICH customer on one trip — that combined fact needs one relationship.",
    "\"Any ternary can be split into three binaries without loss.\" False — pairwise facts can't capture that three entities were associated together in one instance. Use a ternary (or an associative entity).",
    "Three binary relationships record only pairwise facts; they cannot capture that the three entities were associated together in one instance (S–P, P–J, S–J do not imply S supplied P to J). When the meaning is in the combination, a single ternary or an associative entity is required.",
    "ternary ≠ three binaries  ·  loss of the combined fact  ·  pairwise vs three-way association")

concept_understand("S10 · Concept 3 · [THEORY]","Constraints & the Associative Entity",
    "A cleaner alternative to a bare n-ary diamond is the associative (intermediate) entity: replace the diamond with a rectangle linked to each participant — especially when the relationship has its own attributes.",
    ["Use it when the relationship carries its own data (a delivery has time, fee, status).",
     "Also clarifies higher-degree associations that get hard to read as one diamond.",
     "n-ary relationships can carry cardinality constraints, but they're hard to read.",
     "Rule of thumb: prefer binary; use ternary/associative entity only when needed."],
    "s10_associative.png","Ternary + its own attributes → make it an associative entity.",
    "~12 min. Foodmandu DELIVERY linking CUSTOMER, RESTAURANT, RIDER.")
concept_apply("S10 · Concept 3 · [THEORY]","Constraints & the Associative Entity",
    "Foodmandu/Pathao Food model a DELIVERY that links CUSTOMER, RESTAURANT, and RIDER and stores delivery time, fee, and status — an associative entity, cleaner than a bare ternary diamond.",
    "Forcing a relationship-with-its-own-attributes to stay a diamond makes the model awkward and hard to map. When it has its own data, promote it to an associative entity.",
    "An associative (intermediate) entity replaces a relationship (often ternary/n-ary) with a rectangle linked to each participant. Use it when the relationship needs its own attributes or relationships, or to make a higher-degree association clearer — e.g. DELIVERY linking CUSTOMER, RESTAURANT, RIDER with time/fee/status.",
    "associative / intermediate entity  ·  n-ary cardinality  ·  'prefer binary'  ·  relationship with its own attributes")

add_activity("Ternary or associative?",
    ["In pairs (3 min): model a Pathao ride involving RIDER, PASSENGER, and ROUTE.",
     "Decide whether a ternary is justified or an associative RIDE entity is cleaner.",
     "List two attributes that would push you toward the entity.",
     "Share (2 min)."],
    "If the ride has its own data — fare, distance, time, rating — an associative RIDE entity wins over a bare ternary diamond, because a relationship can't cleanly hold those attributes.",
    "ACTIVITY [~5 min].")
add_quiz("S10 — Quick Check",
    [("Q1.  A relationship among three entity types is:","q"),
     ("a) binary   b) recursive   c) ✅ ternary   d) weak","a"),
     ("Q2.  Splitting a ternary into three binaries is risky because:","q"),
     ("✅ it can lose the combined-fact meaning","a"),
     ("Discussion: model a Pathao ride (RIDER, PASSENGER, ROUTE) — ternary or associative entity?","o")],
    "QUIZ [~5 min].")
add_summary("S10 · Summary  ·  [~2 min]",
    ["Degree = number of participating entity types; ternary = 3, drawn as one diamond on three rectangles.",
     "A ternary ≠ three binaries when the meaning is in the combination of all three.",
     "Prefer binary; use a ternary or an associative entity only when needed (especially if it has its own attributes)."],
    "Logistics, delivery, and booking systems (common in Nepal's startup scene) frequently need ternary thinking; getting it wrong silently loses business facts like 'who delivered what to whom'.",
    "S11 — modelling 'kinds of': specialization and generalization.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S11 ============================
add_divider("Session 11 · Unit 2","Specialization & Generalization (EER)",
    "On the Nagarik App a USER can be a citizen, a government employee, or a business owner — same person-record, but each kind stores extra things. How do we model 'kinds of'?",
    "OPENING HOOK [~5 min]. Agenda: specialization/generalization → the constraints (d/o × total/partial) → attribute inheritance.")

concept_understand("S11 · Concept 1 · [THEORY]","Specialization & Generalization",
    "The EER constructs for 'is-a' hierarchies. Specialization is top-down (split a superclass into subclasses); generalization is bottom-up (combine similar entities into a superclass). Same idea, opposite directions.",
    ["Specialization: EMPLOYEE → ENGINEER, SECRETARY, TECHNICIAN (subclasses add their own attributes).",
     "Generalization: CAR and TRUCK → VEHICLE (factor out shared attributes).",
     "Both create a superclass/subclass hierarchy with an 'is-a' relationship.",
     "Both give subclasses attribute inheritance."],
    "s11_specialization.png","Specialization = split down; generalization = combine up; both are is-a.",
    "~12 min. Draw superclass, a circle, subset lines to subclasses.")
concept_apply("S11 · Concept 1 · [THEORY]","Specialization & Generalization",
    "A bank ACCOUNT is specialized into SAVINGS and CURRENT; a platform USER into CUSTOMER and MERCHANT. Each subclass is still 'an account' / 'a user' (is-a) but stores extra fields. Two near-identical tables signal you should generalize upward into a shared superclass.",
    "\"Specialization and generalization are different tools.\" They're inverse directions of the SAME idea — both produce an is-a hierarchy with inheritance; only the direction of thinking differs.",
    "Specialization is top-down: splitting a superclass into subclasses that inherit its attributes and add specific ones (EMPLOYEE → ENGINEER, SECRETARY). Generalization is bottom-up: combining similar entity types into a superclass (CAR, TRUCK → VEHICLE). Both create an is-a hierarchy with inheritance; they are inverse directions of one idea.",
    "superclass / subclass  ·  specialization (top-down)  ·  generalization (bottom-up)  ·  is-a  ·  inheritance  ·  EER")

concept_understand("S11 · Concept 2 · [THEORY]","Constraints: Disjoint/Overlap × Total/Partial",
    "Two independent constraints govern a specialization: disjointness (how many subclasses an entity can join) and completeness (whether it must join one). Together they give four combinations.",
    ["Disjoint (d): an entity is in at most ONE subclass (ACCOUNT is SAVINGS or CURRENT, not both).",
     "Overlap (o): an entity may be in SEVERAL subclasses (a PERSON can be STUDENT and STAFF).",
     "Total (double line): every superclass member MUST be in some subclass.",
     "Partial: a member may be in none. → four combos: d/o × total/partial."],
    "s11_constraints_2x2.png","d/o = how many subclasses; total/partial = must you join one.",
    "~13 min. Walk the 2×2 with one example per cell.")
concept_apply("S11 · Concept 2 · [THEORY]","Constraints: Disjoint/Overlap × Total/Partial",
    "Bank ACCOUNT → SAVINGS/CURRENT is disjoint, total (exactly one type; every account has one). College PERSON → STUDENT/STAFF is overlap, partial (a TA is both; a visitor is neither).",
    "\"Subclasses are always mutually exclusive.\" Not so — overlap is allowed and common (a teaching assistant is both student and staff). Choose disjoint vs overlap from the real rules.",
    "Two constraints. Disjointness: disjoint (d) — an entity is in at most one subclass; overlap (o) — it may be in several. Completeness: total — every superclass entity must be in some subclass (double line); partial — some need not be. They are independent → four combinations, e.g. ACCOUNT→SAVINGS/CURRENT is disjoint-total.",
    "disjoint (d)  ·  overlap (o)  ·  total (double line)  ·  partial  ·  completeness vs disjointness  ·  four combinations")

concept_understand("S11 · Concept 3 · [THEORY]","Attribute Inheritance & Specific Attributes",
    "A subclass inherits all attributes (and relationships) of its superclass, and then adds its own specific (local) attributes. This is why factoring shared attributes up into a superclass avoids repetition.",
    ["SAVINGS inherits ACCOUNT's acc_no, holder, balance — and adds interest_rate.",
     "CURRENT inherits the same — and adds overdraft_limit.",
     "Inheritance flows through relationships too: if ACCOUNT relates to CUSTOMER, every subclass does.",
     "Shared fields live once, in the superclass."],
    "s11_inheritance.png","Inherit the shared, add the specific.",
    "~10 min. Superclass ACCOUNT with shared attrs; subclasses add local attrs.")
concept_apply("S11 · Concept 3 · [THEORY]","Attribute Inheritance & Specific Attributes",
    "In an HR system every EMPLOYEE has id, name, salary (superclass); an ENGINEER additionally has a licence number, a SECRETARY a typing speed — specific attributes layered on top of the inherited ones.",
    "Repeating the shared columns in every subclass table is the smell that you skipped generalization — factor the common attributes up into the superclass instead.",
    "A subclass automatically receives all attributes and relationships of its superclass (inheritance) and adds its own specific attributes. E.g. SAVINGS and CURRENT inherit ACCOUNT's acc_no/holder/balance; SAVINGS adds interest_rate, CURRENT adds overdraft_limit.",
    "attribute inheritance  ·  specific / local attribute  ·  inherited relationships  ·  factoring shared attributes upward")

add_activity("Classify the specialization",
    ["In pairs (3 min): for a platform USER → CUSTOMER and SELLER…",
     "Decide the disjoint/overlap constraint and the total/partial constraint.",
     "List one specific attribute for each subclass.",
     "Share (2 min)."],
    "Overlap (an account can both buy and sell) and partial (a brand-new user may be neither yet). CUSTOMER: default_address; SELLER: store_name, PAN. So it's an overlap, partial specialization.",
    "ACTIVITY [~5 min].")
add_quiz("S11 — Quick Check",
    [("Q1.  'Every employee must be permanent or contract' is a ___ constraint:","q"),
     ("a) overlap   b) disjoint   c) ✅ total   d) partial","a"),
     ("Q2.  A subclass receives the superclass's attributes via:","q"),
     ("✅ inheritance","a"),
     ("Discussion: on Daraz, USER → CUSTOMER and SELLER — disjoint or overlap? total or partial? Justify.","o")],
    "QUIZ [~5 min].")
add_summary("S11 · Summary  ·  [~2 min]",
    ["Specialization (down) and generalization (up) are inverses sharing 'is-a' and inheritance.",
     "Disjoint vs overlap controls multi-subclass membership; total vs partial controls mandatory membership — four combinations.",
     "Subclasses inherit shared attributes and add their own specific ones."],
    "Banking (account types), HR (employee types), and e-commerce (user roles) all rely on specialization; choosing disjoint/total correctly prevents impossible records (an account of both types) and missing ones (a member in no subclass).",
    "S12 — turning all of these ER constructs into real relational tables.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S12 ============================
add_divider("Session 12 · Unit 2 (closes the unit)","Converting ER Diagrams to Tables",
    "We've drawn beautiful diagrams. But MySQL doesn't understand diamonds and ovals — it understands tables. How do we translate?",
    "OPENING HOOK [~5 min]. Agenda: the 7-step algorithm → 1:N (FK) → M:N (junction) → common mistakes → capstone mapping.")

concept_understand("S12 · Concept 1 · [THEORY]","The 7-Step Mapping Algorithm",
    "A mechanical algorithm converts any ER diagram into relational tables. Learn the seven cases and mapping becomes routine.",
    ["1 Strong entity → table + simple attributes; pick a primary key.",
     "2 Weak entity → table incl. owner's PK (owner PK + partial key).",
     "3 1:1 → FK on one side (prefer the total side).   4 1:N → FK on the N (many) side.",
     "5 M:N → NEW junction table with both PKs.   6 Multivalued attr → separate table.",
     "7 n-ary → table with keys of ALL participants.  Composite → flatten; derived → not stored."],
    "s12_mapping_table.png","FK on the N side; junction for M:N; separate table for multivalued/weak.",
    "~11 min. Keep the 7-row table on screen for the rest of the session.")
concept_apply("S12 · Concept 1 · [THEORY]","The 7-Step Mapping Algorithm",
    "This is exactly the step every developer does between an approved ER design and the first CREATE TABLE (Unit 5) — turning a diagram into a running schema, one construct at a time.",
    "Mapping by intuition instead of the algorithm is where schemas go wrong. Apply the seven cases mechanically and the tables fall out correctly every time.",
    "(1) strong entity → table + PK; (2) weak entity → table with owner PK + partial key; (3) 1:1 → FK on one (total) side; (4) 1:N → FK on the N side; (5) M:N → junction table with both PKs; (6) multivalued attribute → separate table; (7) n-ary → table with all participants' keys. Composite attributes flatten; derived attributes are not stored.",
    "mapping algorithm  ·  foreign key  ·  junction table  ·  composite → flatten  ·  derived → omit")

concept_understand("S12 · Concept 2 · [THEORY]","Mapping 1:N — Foreign Key on the N Side",
    "For a 1:N relationship the foreign key goes on the MANY side — each 'many' row points to its one partner. Putting it on the 'one' side would need a multi-valued column.",
    ["DEPARTMENT (1) — has — EMPLOYEE (N)  →  EMPLOYEE gets a dept_id foreign key.",
     "Tables: DEPARTMENT(dept_id PK, name) and EMPLOYEE(emp_id PK, name, dept_id FK).",
     "FK on the '1' side would force one row to hold many ids in one column.",
     "That breaks atomicity / first normal form (Unit 4)."],
    "s12_1n_fk.png","1:N → FK on the many side.",
    "~9 min. Draw the two tables; point at where dept_id lives.")
concept_apply("S12 · Concept 2 · [THEORY]","Mapping 1:N — Foreign Key on the N Side",
    "On Daraz, CUSTOMER (1) — places — ORDER (N): the ORDER table carries a customer_id foreign key; you never store a list of order-ids inside CUSTOMER.",
    "\"Put the FK on the 1 side.\" Then the one row must hold many foreign values in a single column — non-atomic, breaks 1NF, and makes queries and updates unreliable. FK always goes on the N side.",
    "In a 1:N relationship the foreign key goes on the N (many) side, because each many-side row relates to exactly one entity on the one side, so a single FK column suffices. Placing it on the one side would require a multi-valued (non-atomic) column.",
    "1:N mapping  ·  foreign key placement  ·  N-side FK  ·  atomic column (1NF preview)")

concept_understand("S12 · Concept 3 · [THEORY]","Mapping M:N — the Junction Table",
    "An M:N relationship cannot be a single foreign key on either side. It maps to a NEW junction table whose rows are pairs of the two entities' primary keys.",
    ["The pair of keys is the junction table's composite primary key (each half is a FK).",
     "STUDENT (M) — ENROLLS — COURSE (N)  →  ENROLLS(roll_no FK, course_id FK).",
     "Any attributes of the relationship (date, grade) go INSIDE the junction table.",
     "One row per (student, course) pair."],
    "s12_junction.png","M:N → a junction table of paired keys.",
    "~9 min. Show STUDENT, COURSE, and the ENROLLS junction table.")
concept_apply("S12 · Concept 3 · [THEORY]","Mapping M:N — the Junction Table",
    "Daraz ORDER (M) — contains — PRODUCT (N) becomes an ORDER_ITEM(order_id FK, product_id FK, quantity, price) junction table — one row per product in an order, with per-item data held right there.",
    "\"Store M:N as a comma-separated list of ids in a column.\" Then you can't query/aggregate one product cleanly and can't attach per-item data like quantity. The junction table solves both.",
    "An M:N relationship maps to a new junction table whose composite primary key is the two participating entities' primary keys (each a foreign key); any attributes of the relationship become columns of the junction table (e.g. ENROLLS date/grade). Example: STUDENT–COURSE → ENROLLS(roll_no, course_id).",
    "junction / relationship table  ·  composite primary key  ·  relationship attributes → junction  ·  one row per pair")

concept_understand("S12 · Concept 4 · [THEORY]","Common Mapping Mistakes",
    "Three recurring errors turn a clean design into a buggy schema. Knowing them is half of getting mapping right.",
    ["(a) Storing an M:N as a 'list of ids in a column' instead of a junction table — non-atomic, unqueryable.",
     "(b) Putting the 1:N foreign key on the '1' side instead of the N side.",
     "(c) Storing derived attributes (age, order total) instead of computing them — they drift out of date.",
     "Mapping specialization: options are one-table-per-subclass or a single table with a type field (detail in Unit 4)."],
    None,"No lists in a column; FK on the N side; never store derived.",
    "~6 min. These three mistakes cause most 'why is the report wrong/slow' bugs.")
concept_apply("S12 · Concept 4 · [THEORY]","Common Mapping Mistakes",
    "A slow, buggy reporting screen is often traced back to one of these three mapping mistakes made months earlier — a list-in-a-column M:N, a misplaced FK, or a stored derived total that no longer matches its parts.",
    "All three feel convenient at the keyboard and cost dearly later. Follow the algorithm: junction table for M:N, FK on the N side, and omit derived attributes.",
    "Common mistakes: (1) M:N stored as a multi-valued column instead of a junction table; (2) 1:N foreign key placed on the one side instead of the N side; (3) storing derived attributes rather than computing them. Correct: junction table for M:N, FK on the N side, omit derived attributes.",
    "non-atomic column  ·  FK-placement error  ·  storing-derived error  ·  subclass mapping options")

add_solved_problem("S12 · Capstone Solved Problem","Map a college ER diagram to tables",
    "Map a college ER: STUDENT and COURSE (M:N ENROLLS, with a grade), DEPARTMENT (1:N with STUDENT), "
    "and a student's multiple phone numbers. Write every table with its primary and foreign keys.",
    ["Strong entities → their own tables: DEPARTMENT(dept_id PK, name); COURSE(course_id PK, title, credits).",
     "1:N DEPARTMENT–STUDENT → FK on the N side: STUDENT(roll_no PK, name, dept_id FK).",
     "M:N STUDENT–COURSE → new junction table with both PKs: ENROLLS(roll_no FK, course_id FK).",
     "Relationship attribute 'grade' lives INSIDE the junction: ENROLLS(roll_no, course_id, grade), PK = (roll_no, course_id).",
     "Multivalued 'phone' → its own table: STUDENT_PHONE(roll_no FK, phone), PK = (roll_no, phone)."],
    "Five tables: DEPARTMENT(dept_id, name) · STUDENT(roll_no, name, dept_id→DEPARTMENT) · COURSE(course_id, title, "
    "credits) · ENROLLS(roll_no, course_id, grade) with composite PK · STUDENT_PHONE(roll_no, phone). Every M:N became a "
    "junction, every 1:N a FK on the N side, the multivalued attribute its own table — derived data stored nowhere.",
    "Pause and let students attempt each step before revealing. This is the single most exam-tested S12 skill.")

add_quiz("S12 — Quick Check",
    [("Q1.  An M:N relationship maps to:","q"),
     ("a) a FK on either side   b) ✅ a new junction table   c) nothing   d) a derived attribute","a"),
     ("Q2.  In a 1:N relationship the foreign key is placed on the:","q"),
     ("✅ N (many) side","a"),
     ("Discussion: take the eSewa USER–TRANSACTION model (S7) and write the tables (keys + FKs) it maps to.","o")],
    "QUIZ [~5 min].")
add_summary("S12 · Summary  ·  [~2 min]",
    ["A 7-step algorithm mechanically turns any ER diagram into relational tables.",
     "FK on the N side; a junction table for M:N; a separate table for multivalued attributes and weak entities.",
     "Flatten composite attributes into components; never store derived data."],
    "This mapping is the exact bridge between design and CREATE TABLE (Unit 5) — the step every developer performs to turn an approved diagram into a running database in any Nepali tech job.",
    "Unit 3 — querying these tables formally with relational algebra and calculus.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============= END-OF-UNIT REVISION AIDS (whole Unit 2) =============
add_cheatsheet("Unit 2 · Cheat Sheet (S9–S12)","One-page revision reference",
    [("Naming","Entities = singular nouns (STUDENT). Relationships = verbs (ENROLLS). Reads L→R as a sentence."),
     ("Attribute vs entity","Own attributes/relationships → its own ENTITY (CATEGORY). Plain value → attribute (colour). Don't over-model."),
     ("Ternary (degree 3)","One diamond, three entities. ≠ three binaries when the combination matters. Own attributes → associative entity."),
     ("Specialization / generalization","Top-down split / bottom-up combine. 'is-a' + attribute inheritance (shared up, specific down)."),
     ("Spec. constraints","Disjoint(d)/overlap(o) × total/partial = 4 combos. d,total: ACCOUNT→SAVINGS/CURRENT."),
     ("ER → tables","FK on the N side · M:N → junction table (both PKs) · multivalued/weak → own table · derived → not stored.")])

add_glossary("Unit 2 · Glossary (S9–S12)","Key terms — quick reference",
    [("Notation legend","the full Chen/EER symbol set for reading any diagram."),
     ("Naming convention","entities = singular nouns; relationships = verbs."),
     ("Attribute-vs-entity","model as entity only if it has its own attributes/relationships."),
     ("Associative entity","a relationship promoted to a rectangle (has its own attributes)."),
     ("Degree","number of entity types in a relationship (binary/ternary/n-ary)."),
     ("Ternary","a relationship among three entity types."),
     ("Specialization","top-down split of a superclass into subclasses."),
     ("Generalization","bottom-up combining of entities into a superclass."),
     ("Superclass / subclass","the parent type and its 'is-a' children."),
     ("Disjoint (d) / overlap (o)","subclass membership: at most one / possibly several."),
     ("Total / partial","completeness: every member must join a subclass / need not."),
     ("Attribute inheritance","subclass gets the superclass's attributes + relationships."),
     ("Mapping algorithm","7-step ER→relational conversion."),
     ("Foreign key","an attribute referencing another table's primary key."),
     ("Junction table","the table an M:N relationship maps to (composite PK)."),
     ("Composite PK","a primary key made of more than one column.")])

# ============= CONSOLIDATED END-OF-UNIT QUIZ =============
add_divider("Unit 2 · Revision","Consolidated end-of-unit quiz (S5–S12)",
    "Twelve MCQs across the whole unit, then short-answer, diagramming, and discussion questions. Answers to the MCQs are shown; work the rest from the concept slides and Unit2b_material.md.",
    "Use for revision/assessment after S12.")
add_quiz("Part A — Multiple choice (answers shown)",
    [("1.  ER model is used mainly at which phase?   →  ✅ conceptual","a"),
     ("2.  All students enrolled today form a/an   →  ✅ entity set","a"),
     ("3.  'Age' from date of birth is a/an   →  ✅ derived attribute","a"),
     ("4.  Multiple saved phone numbers is a/an   →  ✅ multivalued attribute","a"),
     ("5.  'One department, many employees' is   →  ✅ 1:N","a"),
     ("6.  A double line to a relationship denotes   →  ✅ total participation","a"),
     ("7.  A weak entity is one that   →  ✅ has no key of its own","a"),
     ("8.  Full identifier of a weak entity =   →  ✅ owner's PK + partial key","a"),
     ("9.  A relationship among three entity types is   →  ✅ ternary","a"),
     ("10.  'Savings or current, never both' is   →  ✅ disjoint","a"),
     ("11.  An M:N relationship maps to   →  ✅ a junction table","a"),
     ("12.  In a 1:N the FK goes on the   →  ✅ N side","a")],
    "Consolidated quiz Part A.",compact=True)
add_quiz("Parts B–D — short answer, diagramming & discussion",
    [("Part B — Short answer","q"),
     ("1. Entity type vs entity set (one example each).   2. The four attribute types with a Nepali example each.","o"),
     ("3. Cardinality ratio vs participation — how do they differ?   4. Why does a weak entity need an identifying relationship + partial key?   5. Disjoint vs overlap.","o"),
     ("Part C — Diagramming / applied","q"),
     ("1. ER for a college library (MEMBER, BOOK, BORROWS+date, COPIES weak) — mark keys, partial key, cardinality, participation.","o"),
     ("2. Map a Daraz fragment (CUSTOMER 1:N ORDER, ORDER M:N PRODUCT, multiple ADDRESSES) to tables with PKs and FKs.","o"),
     ("3. Model a banking USER specialization (CUSTOMER, EMPLOYEE) — justify d/o and total/partial; list specific attributes.","o"),
     ("Part D — Discussion","q"),
     ("1. Food-delivery fact as three binaries — what is lost, and when is a ternary right?   2. Daraz 'category' — attribute or entity? Argue both sides.","o")],
    "Consolidated quiz Parts B–D. Model answers live in the concept slides and Unit2b_material.md.",compact=True)

# ---- close ----
add_title("End of Unit 2  ·  IT 220",
          "S5–S12 complete: conceptual models → ER building blocks → relationships → weak entities → ER diagrams → ternary → specialization → ER-to-tables",
          "Built to COURSE_MATERIAL_STANDARD.md · self-contained slides · exports cleanly to PDF · "
          "Next unit (Unit 3): Relational Algebra & Relational Calculus.")

_add_page_numbers()
out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"IT220_Unit2b.pptx")
prs.save(out)
print("Saved",out,"with",len(prs.slides._sldIdLst),"slides")
