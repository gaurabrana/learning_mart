#!/usr/bin/env python3
"""IT220 Unit 3 deck — S13–S17 (Relational Algebra & Relational Calculus), built to
COURSE_MATERIAL_STANDARD.md. Self-contained / PDF-safe two-slide concept pairs; RA
operators shown as symbol + plain-English word + worked example; solved-problem slides
for expression-writing. Content distilled from Unit3_material.md; diagrams in images/.
Run: python3 build_unit3_pptx.py -> IT220_Unit3.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "images")
NAVY=RGBColor(0x0C,0x2B,0x4A); BLUE=RGBColor(0x18,0x5F,0xA5); TEAL=RGBColor(0x0F,0x6E,0x56)
AMBER=RGBColor(0x85,0x4F,0x0B); CORAL=RGBColor(0x99,0x3C,0x1D); LIGHT=RGBColor(0xF2,0xF5,0xF8)
PALET=RGBColor(0xE1,0xF5,0xEE); WHITE=RGBColor(0xFF,0xFF,0xFF); DARK=RGBColor(0x1A,0x1A,0x1A)
GREY=RGBColor(0x55,0x60,0x6B)
BLUE_T=RGBColor(0xE6,0xF1,0xFB); CORAL_T=RGBColor(0xFA,0xEC,0xE7)
TEAL_T=RGBColor(0xE1,0xF5,0xEE); AMBER_T=RGBColor(0xFA,0xEE,0xDA)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

def _bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def _box(s,l,t,w,h): return s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)).text_frame
def _bar(s,c,l,t,w,h):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=c; sp.line.fill.background(); sp.shadow.inherit=False
def _notes(s,t):
    if t: s.notes_slide.notes_text_frame.text=t
def _header(s,kicker,header,kc=BLUE):
    _bar(s,kc,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=kc
    th=_box(s,0.7,0.66,12,1.0); th.word_wrap=True; r=th.paragraphs[0].add_run(); r.text=header
    r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
def _card(s,l,t,w,h,accent,fill,label,body,body_sz=14):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    sp.line.color.rgb=accent; sp.line.width=Pt(1); sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP
    tf.margin_left=Inches(0.22); tf.margin_right=Inches(0.20); tf.margin_top=Inches(0.10); tf.margin_bottom=Inches(0.08)
    p=tf.paragraphs[0]; r=p.add_run(); r.text=label; r.font.size=Pt(12.5); r.font.bold=True; r.font.color.rgb=accent
    p2=tf.add_paragraph(); p2.space_before=Pt(3); r2=p2.add_run(); r2.text=body; r2.font.size=Pt(body_sz); r2.font.color.rgb=DARK
    return sp

def add_title(title,subtitle,footer):
    s=prs.slides.add_slide(BLANK); _bg(s,NAVY); _bar(s,AMBER,0.9,3.05,1.6,0.10)
    tf=_box(s,0.9,2.0,11.5,1.0); tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(34); r.font.bold=True; r.font.color.rgb=WHITE
    t2=_box(s,0.9,3.25,11.5,1.2); t2.word_wrap=True
    r=t2.paragraphs[0].add_run(); r.text=subtitle; r.font.size=Pt(18); r.font.color.rgb=RGBColor(0xCF,0xDD,0xEA)
    t3=_box(s,0.9,6.5,11.5,0.7); t3.word_wrap=True; r=t3.paragraphs[0].add_run(); r.text=footer
    r.font.size=Pt(13); r.font.color.rgb=RGBColor(0x9A,0xB2,0xC6); return s

def add_outcomes(header,kicker,intro,items,nxt):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,header)
    tf=_box(s,0.7,1.75,12.0,4.6); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=intro; r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=NAVY; p.space_after=Pt(8)
    for it in items:
        p=tf.add_paragraph(); r=p.add_run(); r.text="•  "+it; r.font.size=Pt(14.5); r.font.color.rgb=DARK; p.space_after=Pt(5)
    _card(s,0.7,6.35,12.0,0.85,BLUE,BLUE_T,"↩️ Builds on",nxt,body_sz=13)
    return s

def add_roadmap(kicker,title,done,todo):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    _card(s,0.7,1.8,5.85,4.55,TEAL,TEAL_T,"✅  THIS UNIT  ·  S13–S17","\n".join(done),body_sz=14)
    _card(s,6.78,1.8,5.85,4.55,GREY,LIGHT,"⏭️  THEN","\n".join(todo),body_sz=14)
    return s

def add_divider(kicker,title,hook,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,BLUE)
    tk=_box(s,0.9,1.9,11.5,0.6); r=tk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=RGBColor(0xBF,0xDC,0xF2)
    tt=_box(s,0.9,2.5,11.5,1.4); tt.word_wrap=True; r=tt.paragraphs[0].add_run(); r.text=title
    r.font.size=Pt(31); r.font.bold=True; r.font.color.rgb=WHITE
    _bar(s,AMBER,0.9,4.05,1.4,0.09)
    th=_box(s,0.9,4.25,11.5,2.4); th.word_wrap=True; p=th.paragraphs[0]; r=p.add_run(); r.text="🎣  "+hook
    r.font.size=Pt(18); r.font.color.rgb=RGBColor(0xE8,0xF1,0xF9); r.font.italic=True
    _notes(s,timing); return s

def concept_understand(kicker,heading,definition,mechanism,image,hook,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker+"  ·  understand",heading)
    tw=6.2 if image else 11.9
    tf=_box(s,0.7,1.72,tw,5.05); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=definition; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=NAVY; p.space_after=Pt(10)
    p=tf.add_paragraph(); r=p.add_run(); r.text="HOW IT WORKS"; r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=BLUE; p.space_after=Pt(5)
    for m in mechanism:
        p=tf.add_paragraph(); r=p.add_run(); r.text="•  "+m; r.font.size=Pt(14); r.font.color.rgb=DARK; p.space_after=Pt(6)
    if image:
        path=os.path.join(IMG,image)
        if os.path.exists(path):
            pic=s.shapes.add_picture(path,Inches(6.95),Inches(1.9),width=Inches(6.05))
            band_top,band_h=1.9,4.55; h_in=pic.height/914400.0
            pic.top=Inches(band_top+max(0.0,(band_h-h_in)/2.0))
    _bar(s,TEAL,0,6.9,SW.inches,0.6)
    ft=_box(s,0.7,6.96,11.9,0.5); p=ft.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text="🧠  Memory hook:  "+hook; r.font.size=Pt(13.5); r.font.bold=True; r.font.color.rgb=WHITE
    _notes(s,timing); return s

def concept_apply(kicker,heading,example,trap,exam,keyterms):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker+"  ·  apply",heading)
    _card(s,0.7,1.66,12.0,1.48,BLUE ,BLUE_T ,"🌍  Real example",example)
    _card(s,0.7,3.26,12.0,1.30,CORAL,CORAL_T,"⚠️  Common trap — and the fix",trap)
    _card(s,0.7,4.68,12.0,1.92,TEAL ,TEAL_T ,"🎯  Exam-ready answer",exam)
    _card(s,0.7,6.72,12.0,0.60,AMBER,AMBER_T,"🔑  Key terms",keyterms,body_sz=11.5)
    return s

def add_activity(title,steps,expected,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,PALET); _bar(s,TEAL,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text="🛠  IN-CLASS ACTIVITY  ·  THINK–PAIR–SHARE"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=TEAL
    th=_box(s,0.7,0.66,12,0.9); r=th.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
    tf=_box(s,0.7,1.75,11.9,3.6); tf.word_wrap=True; first=True
    for text in steps:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text="•  "+text; r.font.size=Pt(15.5); r.font.color.rgb=DARK; p.space_after=Pt(7)
    _card(s,0.7,5.5,11.9,1.7,TEAL,WHITE,"✅  Expected answer",expected)
    _notes(s,timing); return s

def add_quiz(title,items,timing="",compact=False):
    s=prs.slides.add_slide(BLANK); _bg(s,LIGHT); _bar(s,AMBER,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text="CHECK FOR UNDERSTANDING"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=AMBER
    th=_box(s,0.7,0.66,12,0.9); r=th.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
    qz,az,sb=(14,12.5,5) if compact else (17,15,9)
    tf=_box(s,0.7,1.75,11.9,5.3); tf.word_wrap=True; first=True
    for text,kind in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=text
        if kind=="q": r.font.size=Pt(qz); r.font.bold=True; r.font.color.rgb=NAVY; p.space_before=Pt(sb)
        elif kind=="a": r.font.size=Pt(az); r.font.color.rgb=TEAL; r.font.bold=True; p.level=1
        else: r.font.size=Pt(az); r.font.color.rgb=DARK; r.font.italic=True; p.level=1; p.space_before=Pt(sb-3)
    _notes(s,timing); return s

def add_summary(kicker,takeaways,why,nxt,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,"Summary & Key Takeaways",kc=NAVY)
    tf=_box(s,0.7,1.7,12.0,2.7); tf.word_wrap=True; first=True
    for i,tk in enumerate(takeaways,1):
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=f"{i}.  "+tk; r.font.size=Pt(15); r.font.color.rgb=DARK; p.space_after=Pt(8)
    _card(s,0.7,4.55,12.0,1.35,TEAL,TEAL_T,"💡  Why this matters",why,body_sz=13.5)
    _card(s,0.7,6.02,12.0,1.15,BLUE,BLUE_T,"➡️  Next",nxt,body_sz=13.5)
    _notes(s,timing); return s

def add_solved_problem(kicker,title,scenario,steps,answer,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title,kc=AMBER)
    _card(s,0.7,1.58,12.0,1.0,AMBER,AMBER_T,"📋  Problem",scenario,body_sz=13)
    tf=_box(s,0.85,2.72,11.7,2.95); tf.word_wrap=True; first=True
    for i,st in enumerate(steps,1):
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=f"Step {i}.  "; r.font.size=Pt(12.5); r.font.bold=True; r.font.color.rgb=AMBER
        r2=p.add_run(); r2.text=st; r2.font.size=Pt(12.5); r2.font.color.rgb=DARK; p.space_after=Pt(4)
    _card(s,0.7,5.72,12.0,1.5,TEAL,TEAL_T,"✅  Worked answer",answer,body_sz=13)
    _notes(s,timing); return s

def add_cheatsheet(kicker,title,blocks):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    xs=[0.7,6.78]; ys=[1.72,3.36,5.0]; hh=1.55
    for idx,(label,body) in enumerate(blocks[:6]):
        _card(s,xs[idx%2],ys[idx//2],5.85,hh,BLUE,BLUE_T,label,body,body_sz=12)
    return s

def add_glossary(kicker,title,terms):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    half=(len(terms)+1)//2
    for x,group in [(0.7,terms[:half]),(6.78,terms[half:])]:
        tf=_box(s,x,1.75,5.85,5.3); tf.word_wrap=True; first=True
        for term,defn in group:
            p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            r=p.add_run(); r.text=term+" — "; r.font.size=Pt(11.5); r.font.bold=True; r.font.color.rgb=NAVY
            r2=p.add_run(); r2.text=defn; r2.font.size=Pt(11.5); r2.font.color.rgb=DARK; p.space_after=Pt(6)
    return s

def _add_page_numbers():
    total=len(prs.slides._sldIdLst)
    for i,s in enumerate(prs.slides,1):
        try:
            h=str(s.background.fill.fore_color.rgb); rr,gg,bb=int(h[0:2],16),int(h[2:4],16),int(h[4:6],16)
            dark=(0.299*rr+0.587*gg+0.114*bb)<128
        except Exception: dark=False
        col=RGBColor(0xE8,0xF0,0xF7) if dark else GREY
        tb=_box(s,12.35,0.30,0.9,0.35); p=tb.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
        r=p.add_run(); r.text=f"{i} / {total}"; r.font.size=Pt(10); r.font.color.rgb=col

# ============================================================
#                        BUILD
# ============================================================
add_title("Unit 3 — Relational Algebra & Relational Calculus",
          "IT 220: Database Management System  ·  BIM 4th Semester  ·  Sessions S13–S17",
          "Running schema: Student · Enrollment · Course.  Operators shown as symbol + plain English + a worked example.  "
          "Self-contained slides — every concept complete on the slide face; exports to PDF with no information lost.")

add_outcomes("Unit 3 — Learning Outcomes","overview  ·  s13–s17",
    "By the end of this unit, you will be able to:",
    ["Explain relational algebra as the procedural foundation of SQL (S13)",
     "Apply the unary operations SELECT (σ) and PROJECT (π) (S13)",
     "Apply the set operations UNION (∪), INTERSECTION (∩), DIFFERENCE (−), CARTESIAN PRODUCT (×) (S14)",
     "Apply JOIN (⋈: theta / equi / natural) and DIVISION (÷) (S15)",
     "Use aggregates (ℱ) & outer joins, and read/write Tuple (TRC) and Domain (DRC) calculus (S16–S17)"],
    "Unit 2 gave you the tables (ER → relational). Unit 3 is how we QUERY them formally. Unit 5 (SQL) is the practical syntax for all of it.")

add_roadmap("Unit 3 — Roadmap","Where each session fits (S13–S17)",
    ["S13   What RA is · SELECT (σ) & PROJECT (π)",
     "S14   Set ops: ∪, ∩, −, and CARTESIAN PRODUCT (×)",
     "S15   JOIN (⋈: theta/equi/natural) & DIVISION (÷)",
     "S16   Aggregates (ℱ), outer joins, Tuple calculus (TRC)",
     "S17   Domain calculus (DRC) · safety · synthesis"],
    ["Unit 4   Normalization",
     "Unit 5   SQL — the practical syntax for all of this",
     "Unit 6   Transactions & concurrency",
     "Unit 7   Advanced topics"])

# ============================ S13 ============================
add_divider("Session 13 · Unit 3","Introduction to Relational Algebra · SELECT (σ) & PROJECT (π)",
    "When you search 'students from Pokhara with grade A' on the college portal, the system silently does two tiny operations — pick the right ROWS, then keep only the COLUMNS you asked for. Today we name them.",
    "OPENING HOOK [~5 min]. Agenda: what relational algebra is → SELECT (σ) → PROJECT (π) → composing them.")

concept_understand("S13 · Concept 1 · [THEORY]","What Relational Algebra Is",
    "Relational algebra is a formal, PROCEDURAL query language: a set of operations that each take relation(s) as input and produce a relation as output — so you specify HOW to compute the answer, step by step.",
    ["Closure property: output is always a relation → operations CHAIN into bigger queries.",
     "Procedural (says HOW) — contrast with relational calculus (says WHAT), coming in S16–S17.",
     "It is the theoretical basis of SQL: the optimizer turns your SQL into a relational-algebra expression.",
     "By definition it removes duplicates and has no display formatting."],
    "s13_closure.png","Relation in → relation out → chain.",
    "~12 min. Stress closure — it's what lets small operations snap together like LEGO.")
concept_apply("S13 · Concept 1 · [THEORY]","What Relational Algebra Is",
    "The eSewa transaction-history screen is a chain of relational-algebra steps over a Transactions relation: filter to your account, keep the columns shown, sort. Every SQL SELECT you write is translated into relational algebra under the hood to plan execution.",
    "\"Relational algebra is just SQL.\" No — SQL is a practical language; relational algebra is the underlying math: closed-form, removes duplicates by definition, no display formatting.",
    "Relational algebra is a formal, procedural query language: a set of operations that take relation(s) as input and yield a relation as output. Its closure property (output is always a relation) lets operations be composed. It is procedural (specifies HOW) and is the theoretical foundation of SQL.",
    "relational algebra · procedural · relation (in/out) · closure property · basis of SQL")

concept_understand("S13 · Concept 2 · [THEORY]","SELECT (σ) & PROJECT (π) — the unary operations",
    "The two one-input operations: SELECT (σ) keeps ROWS that satisfy a condition (a horizontal cut); PROJECT (π) keeps chosen COLUMNS (a vertical cut) and removes duplicate rows.",
    ["SELECT σ_condition(R): condition uses =, ≠, <, >, ≤, ≥ and AND/OR/NOT; same schema, rows ≤ original.",
     "PROJECT π_attrs(R): result schema = the listed attributes only; duplicates auto-removed (relations are sets).",
     "Example: σ_city='Pokhara'(Student) → Pokhara students; π_name,program(Student) → those two columns.",
     "NAMING TRAP: SQL's keyword SELECT does the COLUMN job — that is π, not σ. (σ = rows, π = columns.)"],
    "s13_select_project.png","σ picks rows, π picks columns — and SQL SELECT is really π.",
    "~13 min. Hammer the σ-vs-SQL-SELECT naming trap; show the horizontal vs vertical cut.")
concept_apply("S13 · Concept 2 · [THEORY]","SELECT (σ) & PROJECT (π) — the unary operations",
    "Every 'search + show selected fields' feature — a college result lookup, a Khalti transaction filter — is σ then π: filter the rows you want, then keep the columns to display. These are the two most-used operations in all of querying.",
    "\"PROJECT (π) keeps duplicate rows like SQL does.\" Pure relational algebra removes duplicates automatically; SQL needs DISTINCT. And remember σ = rows, not columns.",
    "SELECT σ_condition(R) filters ROWS satisfying a condition (horizontal; same schema; cardinality ≤ R). PROJECT π_attrs(R) keeps chosen COLUMNS (vertical; schema = listed attrs; duplicates removed). σ = rows, π = columns; SQL's SELECT keyword actually performs π.",
    "SELECT (σ, rows) · PROJECT (π, columns) · condition · same-schema · duplicate removal · SQL-SELECT-is-π trap")

add_solved_problem("S13 · Solved Problem","Compose σ then π — 'names of students from Pokhara'",
    "Using Student(sid, name, program, city), write a relational-algebra expression that returns just the "
    "NAMES of students whose city is Pokhara.",
    ["Identify what to filter (rows) and what to keep (columns): rows where city = 'Pokhara'; keep only name.",
     "Filter the rows first with SELECT: σ_city='Pokhara'(Student) → the Pokhara students (all columns).",
     "Then trim to the wanted column with PROJECT: π_name( … ).",
     "Compose them (inner operation first): π_name( σ_city='Pokhara'(Student) ).",
     "Why σ before π? Filtering rows first is the common, efficient pattern — and π drops the duplicate names automatically."],
    "π_name( σ_city='Pokhara'(Student) ).  Read it inside-out: SELECT the Pokhara rows, then PROJECT the name column. "
    "This σ-then-π pattern is the backbone of almost every simple query.",
    "Let students attempt before revealing; then flip it (π then σ) and discuss why row-filter-first is preferred.")

add_activity("Spot the σ and the π",
    ["In pairs (2 min): open any app with a filter — Daraz price filter, eSewa date filter, a result lookup.",
     "Identify which part is doing SELECT (σ) — choosing rows — and which is doing PROJECT (π) — choosing columns.",
     "Write one line: 'σ = … , π = …' for your chosen screen.",
     "Share (3 min)."],
    "e.g. Daraz: σ = the price/brand filter that keeps only matching products (rows); π = the product card showing just image, name, price (columns). Every 'search + show fields' screen is σ then π.",
    "ACTIVITY [~5 min].")
add_quiz("S13 — Quick Check",
    [("Q1.  Which operation reduces the number of COLUMNS?","q"),
     ("a) σ (SELECT)   b) ✅ π (PROJECT)   c) ∪ (UNION)   d) ⋈ (JOIN)","a"),
     ("Q2.  σ_credit>3(Course) changes which of these?","q"),
     ("a) the schema   b) ✅ the number of rows   c) the column names   d) nothing","a"),
     ("Discussion: on a phone-app filter, which part is σ (rows) and which is π (columns)?","o")],
    "QUIZ [~5 min].")
add_summary("S13 · Summary  ·  [~2 min]",
    ["Relational algebra = procedural, closed operations over relations (output is always a relation, so they chain).",
     "SELECT (σ) filters ROWS by a condition; PROJECT (π) keeps chosen COLUMNS and drops duplicates.",
     "Naming trap: SQL's SELECT keyword does π's job (columns), not σ's."],
    "Every 'search then show selected fields' screen — result lookups, transaction filters — is σ then π: the two most-used operations in all of querying, and the first thing an interviewer expects you to reason about.",
    "S14 — treating whole relations as sets: UNION, INTERSECTION, DIFFERENCE, and CARTESIAN PRODUCT.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S14 ============================
add_divider("Session 14 · Unit 3","Set-Theory Operations: ∪, ∩, −, ×",
    "A college merges its morning and day-shift student lists for one notice. How do you combine two tables without listing anyone twice — and what if their columns don't match?",
    "OPENING HOOK [~5 min]. Agenda: union compatibility + UNION/INTERSECTION/DIFFERENCE → CARTESIAN PRODUCT.")

concept_understand("S14 · Concept 1 · [THEORY]","UNION (∪), INTERSECTION (∩), DIFFERENCE (−)",
    "These treat relations as SETS of tuples. All three need the relations to be UNION-COMPATIBLE: the same number of attributes with matching domains, column by column (names may differ, domains must align).",
    ["UNION (R ∪ S): tuples in R OR S, duplicates removed.",
     "INTERSECTION (R ∩ S): tuples in BOTH.",
     "DIFFERENCE (R − S): tuples in R but NOT in S — and it is NOT commutative (R − S ≠ S − R).",
     "e.g. AllStudents − GraduatedStudents = currently active students."],
    "s14_set_ops.png","Same shape first; ∪ = either, ∩ = both, − = only-in-R.",
    "~17 min. Draw the three Venn diagrams; stress union compatibility and the order-sensitivity of −.")
concept_apply("S14 · Concept 1 · [THEORY]","UNION (∪), INTERSECTION (∩), DIFFERENCE (−)",
    "UNION merges two mailing lists into one deduplicated list; INTERSECTION finds customers with BOTH an eSewa AND a Khalti account; DIFFERENCE finds lapsed users = AllUsers − ActiveUsers. Everyday analytics tasks.",
    "\"Any two tables can be UNIONed\" / \"R − S = S − R.\" Both wrong — the relations must be union-compatible first, and DIFFERENCE is order-sensitive.",
    "UNION (∪, in either), INTERSECTION (∩, in both), and DIFFERENCE (−, in R not S) all require UNION-COMPATIBLE relations (same number of attributes, matching domains). DIFFERENCE is not commutative. All use set semantics (duplicates removed).",
    "union-compatible · UNION (∪) · INTERSECTION (∩) · DIFFERENCE (−, not commutative) · set semantics")

concept_understand("S14 · Concept 2 · [THEORY]","CARTESIAN PRODUCT (×)",
    "R × S pairs EVERY tuple of R with EVERY tuple of S. The result's degree = sum of the two degrees; its cardinality = product of the two cardinalities. It does NOT need union compatibility.",
    ["On its own it is usually meaningless and huge — 4 students × 3 courses = 12 pairings, most nonsensical.",
     "It becomes useful only when followed by a SELECT (σ) that keeps the meaningful pairs.",
     "That combination — × then σ — is exactly a JOIN (S15).",
     "If R has 5 rows and S has 4, R × S has 20 rows."],
    "s14_cartesian.png","× pairs everything; a condition makes it a JOIN.",
    "~13 min. Show the pairing grid ballooning; motivate that you almost always want it filtered.")
concept_apply("S14 · Concept 2 · [THEORY]","CARTESIAN PRODUCT (×)",
    "You rarely want the raw product — you want it filtered. That is why × is the raw material for JOIN: pair everything, then keep only the rows whose keys match (S15). Alone, × just explodes the row count.",
    "\"Cartesian product combines related rows.\" No — it blindly pairs everything; relating them meaningfully needs a following condition (a σ), which turns it into a JOIN.",
    "R × S pairs every tuple of R with every tuple of S; degree = sum, cardinality = product; no union compatibility needed. Alone it is meaningless/large; followed by a selection it becomes a JOIN.",
    "CARTESIAN PRODUCT (×) · degree = sum · cardinality = product · no compatibility needed · raw material for JOIN")

add_activity("Overlap, missing, and the explosion",
    ["In pairs (3 min): give one Nepali example where you'd want INTERSECTION (overlap) and one for DIFFERENCE (what's missing).",
     "e.g. INTERSECTION = users on both eSewa and Khalti; DIFFERENCE = all users − active users = lapsed users.",
     "Then: if R has 5 rows and S has 4 rows, how many rows in R × S?",
     "Share (2 min)."],
    "INTERSECTION finds who is in both sets; DIFFERENCE finds who is in one but not the other (order matters). R × S = 5 × 4 = 20 rows — most meaningless until filtered into a JOIN.",
    "ACTIVITY [~5 min].")
add_quiz("S14 — Quick Check",
    [("Q1.  Which operation does NOT require union compatibility?","q"),
     ("a) UNION   b) INTERSECTION   c) ✅ CARTESIAN PRODUCT   d) DIFFERENCE","a"),
     ("Q2.  If R has 5 rows and S has 4 rows, R × S has how many rows?","q"),
     ("✅ 20  (cardinality = product)","a"),
     ("Discussion: a real Nepali example of INTERSECTION (overlap) vs DIFFERENCE (missing).","o")],
    "QUIZ [~5 min].")
add_summary("S14 · Summary  ·  [~2 min]",
    ["UNION / INTERSECTION / DIFFERENCE need union-compatible relations; DIFFERENCE is order-sensitive.",
     "CARTESIAN PRODUCT (×) pairs everything: degree = sum, cardinality = product; no compatibility needed.",
     "× alone is meaningless — it is the raw material a condition turns into a JOIN."],
    "Merging mailing lists (UNION), finding common customers across eSewa/Khalti (INTERSECTION), and finding lapsed users (DIFFERENCE) are everyday analytics — and × is where joins come from.",
    "S15 — add a condition to × to get meaningful matches: the JOIN, plus DIVISION.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S15 ============================
add_divider("Session 15 · Unit 3","Binary Operations: JOIN (⋈) & DIVISION (÷)",
    "Your marksheet shows the course name and your grade together — but those two facts live in two different tables. What stitches them into one row?",
    "OPENING HOOK [~5 min]. Agenda: × → JOIN (theta) → equijoin vs natural join → DIVISION.")

concept_understand("S15 · Concept 1 · [THEORY]","JOIN: from × to the Theta Join (⋈)",
    "A JOIN is a CARTESIAN PRODUCT followed by a SELECT on a matching condition: R ⋈_θ S = σ_θ(R × S). It keeps only the pairs that satisfy the join condition, discarding the meaningless pairings of ×.",
    ["A theta join (⋈_θ) joins on ANY comparison θ: =, <, >, ≤, ≥, ≠.",
     "Recap: S14's Student × Course exploded to 12 rows; a condition keeps only the sensible ones.",
     "Example: pair students with scholarship slabs where marks ≥ slab_threshold.",
     "Joins build every report that combines tables."],
    "s15_join.png","JOIN = × + condition.",
    "~11 min. Morph Student × Enrollment → JOIN: the meaningless pairs fall away.")
concept_apply("S15 · Concept 1 · [THEORY]","JOIN: from × to the Theta Join (⋈)",
    "A bank statement is Account ⋈ Transactions; a Daraz invoice is Order ⋈ Product ⋈ Customer. Every report that stitches two or more tables together is built on joins.",
    "\"A join is a separate, magic operation.\" It is literally × then σ — a Cartesian product filtered by a condition. Understanding that demystifies every join you'll ever write.",
    "A join is a Cartesian product followed by a selection on a condition: R ⋈_θ S = σ_θ(R × S). A theta join uses any comparison operator in θ (=, <, >, ≤, ≥, ≠); it keeps only the matching pairs.",
    "JOIN (⋈) · theta join (⋈_θ) · join condition · × + σ")

concept_understand("S15 · Concept 2 · [THEORY]","Equijoin vs Natural Join (⋈)",
    "An equijoin is a theta join whose condition uses only EQUALITY (=); the join column then appears twice. A natural join (⋈) is an equijoin on ALL same-named attributes that then automatically removes the duplicate column.",
    ["Equijoin: Student ⋈_(Student.sid = Enrollment.sid) Enrollment → keeps BOTH sid columns.",
     "Natural join: Student ⋈ Enrollment → joins on sid automatically, yields ONE clean sid.",
     "Then π_name,course,grade(Student ⋈ Enrollment) produces a marksheet row.",
     "Natural join matches by identical NAME, not position — a rename can silently break it."],
    "s15_natural_join.png","Equijoin keeps both; natural join zips to one.",
    "~12 min. Use the zipper image; contrast the two result schemas side by side.")
concept_apply("S15 · Concept 2 · [THEORY]","Equijoin vs Natural Join (⋈)",
    "Your marksheet is π_name,course,grade(Student ⋈ Enrollment): the natural join on sid stitches the Student and Enrollment tables into one clean row per enrollment, then project keeps the three columns you see.",
    "\"Natural join matches by column position\" / \"equijoin removes the duplicate.\" Natural join matches by identical NAME (a rename breaks it into a product); only the NATURAL join drops the duplicate column — the equijoin keeps both.",
    "An equijoin joins on an equality condition and keeps BOTH matching columns. A natural join (⋈) joins on ALL like-named attributes and REMOVES the duplicate column. Natural join = equijoin on same-named attributes + drop the duplicate.",
    "equijoin (= only, both columns) · natural join (⋈, same-named attrs, duplicate removed) · join on name not position")

concept_understand("S15 · Concept 3 · [THEORY]","DIVISION (÷)",
    "R ÷ S returns the tuples in R that are associated with ALL tuples in S — the 'for all' (universal) operation. It answers 'find X related to EVERY Y' queries.",
    ["Example: 'students enrolled in ALL core courses' = Enrollment(sid,course) ÷ CoreCourses(course).",
     "On the sample data: sid 1 and 3 took both DBMS and Java → they survive; sid 2 took only DBMS → excluded.",
     "It is the hardest operation — anchor it in one strong example.",
     "It is ALL (universal), not ANY (existential)."],
    "s15_division.png","Division = related to ALL of a set (for-all), not any.",
    "~12 min. Walk the division grid slowly: which sid matched every required course?")
concept_apply("S15 · Concept 3 · [THEORY]","DIVISION (÷)",
    "'Which delivery riders have covered EVERY district in the valley?' = Coverage(rider,district) ÷ ValleyDistricts(district). Division powers 'qualified for ALL requirements' checks: graduation eligibility, completed all KYC steps.",
    "\"Division finds X in ANY of the set.\" No — it is ALL (universal), not any (existential). 'Any' is a simple σ/JOIN; 'all' is division.",
    "R ÷ S returns the tuples of R associated with ALL tuples of S — the universal ('for all') operation. e.g. students enrolled in EVERY core course = Enrollment ÷ CoreCourses; a rider covering EVERY district = Coverage ÷ Districts.",
    "DIVISION (÷) · universal / 'for all' · R ÷ S · all-not-any")

add_solved_problem("S15 · Solved Problem","Build a marksheet from two tables",
    "Using Student(sid, name, program, city) and Enrollment(sid, course, grade), write an expression that "
    "produces each student's NAME with their COURSE and GRADE on one row.",
    ["The needed facts live in two tables: name is in Student, course & grade in Enrollment — they share sid.",
     "Stitch them on the shared attribute with a natural join: Student ⋈ Enrollment (joins on sid, one clean sid).",
     "The joined relation now has sid, name, program, city, course, grade together.",
     "Keep only the wanted columns with PROJECT: π_name,course,grade( … ).",
     "Compose: π_name,course,grade( Student ⋈ Enrollment )."],
    "π_name,course,grade( Student ⋈ Enrollment ).  The natural join merges on sid (no duplicate column), then project "
    "trims to the marksheet columns. This join-then-project shape is behind almost every real report.",
    "Let students try; check they use the NATURAL join (one sid) and remember join-then-project order.")

add_quiz("S15 — Quick Check",
    [("Q1.  Which join automatically drops the duplicate join column?","q"),
     ("a) theta   b) equijoin   c) ✅ natural join   d) Cartesian product","a"),
     ("Q2.  'Students enrolled in every core course' is best expressed with:","q"),
     ("✅ DIVISION (÷)","a"),
     ("Discussion: describe a real 'must match ALL of them' situation — that's division.","o")],
    "QUIZ [~5 min].")
add_summary("S15 · Summary  ·  [~2 min]",
    ["JOIN = CARTESIAN PRODUCT + a matching condition (⋈_θ = σ_θ(R × S)).",
     "theta → equi → natural narrows down; the natural join removes the duplicate join column.",
     "DIVISION (÷) answers 'related to ALL of a set' — universal, not existential."],
    "Joins build every report that combines tables — bank statements, Daraz invoices; division powers 'qualified for all requirements' checks like graduation eligibility or completed KYC.",
    "S16 — operations beyond the basics (aggregates, outer joins) and the declarative world of Tuple Relational Calculus.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S16 ============================
add_divider("Session 16 · Unit 3","Aggregates (ℱ), Outer Joins & Tuple Relational Calculus",
    "Basic relational algebra can't answer 'what's the AVERAGE grade?' or keep a student who enrolled in NOTHING. We need a few extra tools — and a totally different way of asking questions.",
    "OPENING HOOK [~5 min]. Agenda: aggregate/grouping (ℱ) → outer joins → Tuple Relational Calculus (declarative).")

concept_understand("S16 · Concept 1 · [THEORY]","Aggregate Functions & Grouping (ℱ)",
    "The aggregate operator ℱ (script F) applies COUNT, SUM, AVG, MAX, MIN — optionally GROUPED by attributes: <grouping attrs> ℱ <functions>(R). Grouping partitions rows into groups, then summarizes each group.",
    ["It is NOT one of the basic operations — it is an EXTENSION (basic algebra can't count).",
     "program ℱ COUNT(sid)(Student) → number of students per program.",
     "course ℱ AVG(grade)(Enrollment) → average grade per course.",
     "Grouping first partitions, then one summary value is produced per group."],
    "s16_aggregate.png","ℱ = group, then summarize each group.",
    "~11 min. Show raw rows → grouped summary; note aggregates are an added/extended operation.")
concept_apply("S16 · Concept 1 · [THEORY]","Aggregate Functions & Grouping (ℱ)",
    "Dashboards live on aggregates: merchant ℱ SUM(amount)(Transactions) = total collected per merchant on Khalti; program ℱ COUNT(sid)(Student) = students per program on a college portal.",
    "\"Aggregate is one of the basic relational operations.\" It is an additional/extended operation — the basic six have no built-in counting or summarizing.",
    "The aggregate/grouping operation <grouping> ℱ <functions>(R) partitions R by the grouping attributes and applies aggregate functions (COUNT/SUM/AVG/MAX/MIN) per group. It is an EXTENDED (not basic) operation. e.g. course ℱ AVG(grade)(Enrollment).",
    "aggregate (ℱ) · COUNT / SUM / AVG / MAX / MIN · grouping · extended operation")

concept_understand("S16 · Concept 2 · [THEORY]","Outer Joins (⟕ ⟖ ⟗)",
    "An inner join (natural/theta) DROPS rows with no match. Outer joins KEEP the unmatched rows too, padding the missing side with NULL.",
    ["LEFT outer (⟕) keeps all LEFT rows; RIGHT outer (⟖) all RIGHT rows; FULL outer (⟗) keeps both.",
     "Example: Student ⟕ Enrollment keeps a newly-admitted student who has enrolled in nothing yet.",
     "That student's course & grade columns become NULL.",
     "Contrast: an inner join would silently hide them."],
    "s16_outer_join.png","Inner join hides the unmatched; outer join keeps them with NULLs.",
    "~12 min. Show Ram (no enrolment) dropped by inner, kept with NULLs by left outer.")
concept_apply("S16 · Concept 2 · [THEORY]","Outer Joins (⟕ ⟖ ⟗)",
    "A bank report listing EVERY account, including those with zero transactions this month, is Account ⟕ Transactions (left outer). An inner join would silently omit the inactive accounts — exactly the ones an auditor asks about.",
    "\"Inner join shows everyone.\" It silently drops rows with no match. To keep unmatched rows (a new student, a dormant account), you need an OUTER join, which fills the gaps with NULL.",
    "An outer join keeps unmatched tuples (unlike an inner join, which drops them), padding missing values with NULL: LEFT (⟕) keeps all left rows, RIGHT (⟖) all right rows, FULL (⟗) both. e.g. Student ⟕ Enrollment keeps students with no enrolment.",
    "inner vs outer join · LEFT (⟕) / RIGHT (⟖) / FULL (⟗) · NULL padding")

concept_understand("S16 · Concept 3 · [THEORY]","Tuple Relational Calculus (TRC)",
    "TRC is a DECLARATIVE (non-procedural) language: you describe WHAT tuples you want, not HOW to get them. Form: { t | COND(t) }, where t is a tuple variable ranging over a relation.",
    ["COND uses predicates and the quantifiers ∃ (there exists) and ∀ (for all).",
     "Same expressive power as relational algebra (relational completeness).",
     "{ t.name | t ∈ Student ∧ t.program='BIM' } → names of BIM students (declarative π∘σ).",
     "'Enrolled in at least one course': { t | t ∈ Student ∧ ∃ e (e ∈ Enrollment ∧ e.sid=t.sid) }."],
    "s16_how_vs_what.png","Algebra = HOW; calculus = WHAT.",
    "~12 min. Pair one algebra expression with its TRC twin; introduce ∃ gently.")
concept_apply("S16 · Concept 3 · [THEORY]","Tuple Relational Calculus (TRC)",
    "Declarative thinking is exactly how you reason in SQL day to day: you state the result you want ('give me BIM students' names') and let the DBMS figure out the steps — that's the calculus mindset, not the algebra one.",
    "\"Calculus and algebra give different answers.\" They are equivalent in power (relational completeness) — same result set, different style: algebra says HOW, calculus says WHAT.",
    "Tuple relational calculus is a declarative (non-procedural) language of the form { t | COND(t) }, where t is a tuple variable and COND uses predicates and quantifiers (∃, ∀). It describes WHAT to retrieve, not HOW, and is equal in power to relational algebra.",
    "TRC · declarative / non-procedural · tuple variable · ∃ / ∀ · relational completeness")

add_activity("HOW or WHAT?",
    ["In pairs (3 min): phrase one query about your college in plain English.",
     "Now express it two ways: as algebra steps (σ, π, ⋈ …) and as a calculus description ({ t | … }).",
     "Say which felt like giving directions (HOW = algebra) and which like naming a destination (WHAT = calculus).",
     "Share (2 min)."],
    "e.g. 'names of BIM students' — algebra: π_name(σ_program='BIM'(Student)) [HOW]; TRC: { t.name | t ∈ Student ∧ t.program='BIM' } [WHAT]. Same answer, different mindset.",
    "ACTIVITY [~5 min].")
add_quiz("S16 — Quick Check",
    [("Q1.  Which keeps students who enrolled in NO course?","q"),
     ("a) natural join   b) equijoin   c) ✅ LEFT outer join   d) division","a"),
     ("Q2.  Tuple relational calculus is best described as:","q"),
     ("✅ declarative / non-procedural","a"),
     ("Discussion: phrase a college query in English — did you describe HOW (algebra) or WHAT (calculus)?","o")],
    "QUIZ [~5 min].")
add_summary("S16 · Summary  ·  [~2 min]",
    ["Aggregates/grouping (ℱ) add counting & summarizing — an extended operation basic algebra lacks.",
     "Outer joins (⟕ ⟖ ⟗) keep unmatched rows, padding missing values with NULL.",
     "Tuple calculus (TRC) describes WHAT you want with a tuple variable and quantifiers (∃/∀); equal in power to algebra."],
    "Dashboards live on aggregates (sales per merchant, students per program); outer joins build 'include everyone, even the inactive' reports; declarative thinking is exactly how you reason in SQL every day.",
    "S17 — the second calculus, Domain Relational Calculus, plus expression safety and the unit synthesis.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============================ S17 ============================
add_divider("Session 17 · Unit 3 (closes the unit)","Domain Relational Calculus (DRC) · Safety · Synthesis",
    "Yesterday we asked for whole tuples. Today we go one level finer — we name a variable for each column VALUE and describe the answer field by field. Same power, different feel.",
    "OPENING HOOK [~5 min]. Agenda: DRC → TRC vs DRC & safety → Unit 3 synthesis (procedural ⇄ declarative ⇄ SQL).")

concept_understand("S17 · Concept 1 · [THEORY]","Domain Relational Calculus (DRC)",
    "DRC is a declarative calculus whose variables range over DOMAINS (individual attribute VALUES) rather than whole tuples. Form: { ⟨x1, …, xn⟩ | COND(x1, …, xn) }, each xi a domain variable bound to a column.",
    ["Uses the quantifiers ∃ and ∀, like TRC.",
     "It is relationally complete (= algebra = TRC).",
     "It is the basis of the QBE (Query-By-Example) 'fill in the columns' interface.",
     "Names of Pokhara students: { ⟨n⟩ | ∃ s,p ( ⟨s,n,p,'Pokhara'⟩ ∈ Student ) }."],
    "s17_three_ways.png","DRC = one variable per column.",
    "~12 min. Show the same query as algebra / TRC / DRC stacked; DRC names each column.")
concept_apply("S17 · Concept 1 · [THEORY]","Domain Relational Calculus (DRC)",
    "QBE interfaces — the visual 'fill the columns' query builders in tools like MS Access and report designers used in many Nepali offices and banks — are DRC made visual: each blank you fill is a domain variable.",
    "\"DRC and TRC are the same thing with different symbols.\" They differ in variable SCOPE — a domain variable is a single column value, a tuple variable is a whole row — though they have equal power.",
    "Domain relational calculus is a declarative calculus of the form { ⟨x1,…,xn⟩ | COND } where each variable ranges over a DOMAIN (a single attribute value), using ∃/∀. It is equal in power to algebra and TRC and underlies QBE.",
    "DRC · domain variable · ⟨…⟩ list · ∃ / ∀ · QBE · relational completeness")

concept_understand("S17 · Concept 2 · [THEORY]","TRC vs DRC · Safety of Expressions",
    "TRC and DRC differ in variable SCOPE, not power; and both must be SAFE — producing a finite result drawn from values actually in the database.",
    ["TRC: a tuple variable ranges over whole ROWS (t.name).",
     "DRC: a domain variable ranges over single column VALUES (one variable per attribute).",
     "DRC names every attribute in the ⟨…⟩ list — more verbose, but maps directly onto QBE 'fill the blanks'.",
     "Safe = finite result. { ⟨x⟩ | NOT(⟨x⟩ ∈ Student) } is UNSAFE — infinitely many values aren't students."],
    None,"Tuple var = a row; domain var = a value. Safe = a finite answer.",
    "~11 min. Show the same query both ways; give one clearly unsafe expression and why it's banned.")
concept_apply("S17 · Concept 2 · [THEORY]","TRC vs DRC · Safety of Expressions",
    "A report tool that let you ask for 'everything NOT in this table' would try to list infinitely many values — which is why real query languages only allow safe, bounded expressions drawn from the stored data.",
    "\"Any condition you can write is a valid query.\" Unsafe expressions (like NOT ∈) would return infinite sets and are disallowed; a valid query must be safe (finite, drawn from database values).",
    "In TRC a variable ranges over whole tuples; in DRC over individual domain values (one per column) — same power, different variable scope. A SAFE expression yields a finite result drawn from database values; unsafe expressions (e.g. NOT ∈) return infinite sets and are disallowed.",
    "tuple vs domain variable · scope · safe / unsafe expression · finite result")

concept_understand("S17 · Concept 3 · [THEORY]","Unit 3 Synthesis: Procedural ⇄ Declarative ⇄ SQL",
    "Relational algebra (procedural) and relational calculus (declarative: TRC over tuples, DRC over domains) are all RELATIONALLY EQUIVALENT — they express the same queries; algebra says HOW, calculus says WHAT.",
    ["Algebra: σ, π, ∪, ∩, −, ×, ⋈, ÷, plus aggregates (ℱ) and outer joins.",
     "Calculus: TRC { t | … } over tuples; DRC { ⟨x⟩ | … } over domains.",
     "All three have the same expressive power (relational completeness).",
     "All three underlie SQL (Unit 5) — the practical, friendly syntax for everything here."],
    "s17_synthesis.png","Algebra (HOW) ⇄ TRC/DRC (WHAT) → all equal → all become SQL.",
    "~7 min. This is the slide to photograph — the whole unit as one map into SQL.")
concept_apply("S17 · Concept 3 · [THEORY]","Unit 3 Synthesis: Procedural ⇄ Declarative ⇄ SQL",
    "Every SQL query you write in Unit 5 is, underneath, one of these formal expressions: SELECT…WHERE is σ, the column list is π, JOIN is ⋈, GROUP BY is ℱ, and the whole thing is planned as a relational-algebra tree by the optimizer.",
    "\"Algebra, TRC, and DRC are three unrelated topics.\" They are three equivalent faces of the same query power — and SQL is the practical language sitting on top of all of them.",
    "Relational algebra (procedural) and relational calculus (TRC over tuples, DRC over domains — declarative) are all relationally equivalent (relational completeness): each expresses the same queries. All three underlie SQL.",
    "relational equivalence · procedural vs declarative · foundation of SQL")

add_activity("Why is 'fill in the blanks' easier?",
    ["In pairs (3 min): why might a non-programmer find DRC-style 'fill the columns' (QBE) easier than writing algebra?",
     "Where have you seen fill-in-the-blank query forms (MS Access, report builders, search filters)?",
     "Map one such form onto DRC: each blank you fill is a domain variable.",
     "Share (2 min)."],
    "QBE lets you state WHAT you want by filling column blanks — no procedure to write — which is exactly DRC's model (a variable per column). Many office/bank report tools and advanced search forms are QBE-style.",
    "ACTIVITY [~5 min].")
add_quiz("S17 — Quick Check",
    [("Q1.  In domain relational calculus, variables range over:","q"),
     ("✅ individual attribute values / domains","a"),
     ("Q2.  An expression that would return an infinite set of values is called:","q"),
     ("✅ unsafe","a"),
     ("Discussion: where have you seen QBE-style 'fill in the blank' query forms?","o")],
    "QUIZ [~5 min].")
add_summary("S17 · Summary  ·  [~2 min]",
    ["DRC uses one variable per COLUMN and describes the answer field-by-field; it underlies QBE.",
     "TRC vs DRC differ in variable scope (row vs value), not power; only SAFE (finite) expressions are allowed.",
     "Algebra, TRC, and DRC are all relationally equivalent — and all feed SQL."],
    "QBE 'fill the columns' interfaces (MS Access, report designers used in Nepali offices/banks) are DRC made visual; understanding it demystifies those drag-and-fill report tools.",
    "Unit 4 — putting good STRUCTURE under all this: Database Normalization.",
    "REAL-LIFE APPLICATION [~3 min] + SUMMARY [~2 min].")

# ============= END-OF-UNIT REVISION AIDS =============
add_cheatsheet("Unit 3 · Cheat Sheet","One-page revision reference",
    [("SELECT (σ) & PROJECT (π)","σ = rows by condition (same schema). π = chosen columns (drops duplicates). SQL's SELECT keyword = π."),
     ("Set ops","∪ = either · ∩ = both · − = only-in-R (order matters) — all need union-compatible relations. × pairs all (degree=sum, card=product)."),
     ("JOIN (⋈)","⋈ = × + condition. theta (any comparison) → equijoin (=, keeps both cols) → natural (same-name, drops duplicate)."),
     ("DIVISION (÷)","R ÷ S = tuples related to ALL of S (universal 'for all', not 'any'). e.g. enrolled in every core course."),
     ("Extended ops","ℱ = group + aggregate (COUNT/SUM/AVG/MAX/MIN). Outer joins ⟕ ⟖ ⟗ keep unmatched rows, padding NULL."),
     ("Calculus","TRC { t | COND } over tuples; DRC { ⟨x⟩ | COND } over domains; ∃/∀; safe = finite. algebra ≡ TRC ≡ DRC ≡ SQL.")])

add_glossary("Unit 3 · Glossary","Key terms — quick reference",
    [("Relational algebra","procedural query language; operations on relations."),
     ("Closure","output is always a relation, so operations chain."),
     ("SELECT (σ)","keeps ROWS satisfying a condition (horizontal)."),
     ("PROJECT (π)","keeps chosen COLUMNS, removes duplicates (vertical)."),
     ("Union-compatible","same #attributes and matching domains."),
     ("UNION / INTERSECTION","in either (∪) / in both (∩)."),
     ("DIFFERENCE (−)","in R not S; not commutative."),
     ("CARTESIAN PRODUCT (×)","pairs every row with every row."),
     ("JOIN (⋈)","× followed by a matching condition."),
     ("Theta / equijoin","join on any comparison / on equality (keeps both cols)."),
     ("Natural join","equijoin on same-named attrs; drops the duplicate."),
     ("DIVISION (÷)","tuples related to ALL of a set ('for all')."),
     ("Aggregate (ℱ)","group + COUNT/SUM/AVG/MAX/MIN (extended op)."),
     ("Outer join (⟕⟖⟗)","keeps unmatched rows, padding NULL."),
     ("TRC / DRC","calculus over tuples / over domains (declarative)."),
     ("Safe expression","yields a finite result from database values.")])

# ============= CONSOLIDATED END-OF-UNIT QUIZ =============
add_divider("Unit 3 · Revision","Consolidated end-of-unit quiz (S13–S17)",
    "Twelve MCQs across the whole unit (answers shown), then write-the-expression and discussion questions to work from the concept slides and Unit3_material.md.",
    "Use for revision/assessment after S17.")
add_quiz("Part A — Multiple choice (answers shown)",
    [("1.  Relational algebra is a ___ query language   →  ✅ procedural","a"),
     ("2.  Selects ROWS by a condition   →  ✅ σ (SELECT)","a"),
     ("3.  π automatically removes   →  ✅ duplicate rows","a"),
     ("4.  Does NOT require union compatibility   →  ✅ × (Cartesian product)","a"),
     ("5.  R has 6 rows, S has 5 → R × S has   →  ✅ 30 rows","a"),
     ("6.  SET DIFFERENCE (R − S) is   →  ✅ not commutative","a"),
     ("7.  A JOIN = Cartesian product followed by a   →  ✅ SELECT (σ)","a"),
     ("8.  Removes the duplicate matching column   →  ✅ natural join","a"),
     ("9.  'Enrolled in ALL core courses'   →  ✅ DIVISION (÷)","a"),
     ("10.  Keep students who enrolled in NO course   →  ✅ LEFT OUTER JOIN","a"),
     ("11.  In TRC a variable ranges over   →  ✅ whole tuples (rows)","a"),
     ("12.  In DRC a variable ranges over   →  ✅ single attribute values (domains)","a")],
    "Consolidated quiz Part A.",compact=True)
add_quiz("Parts B–C — write the expression & discuss",
    [("Part B — Write the relational-algebra / calculus expression (sample relations)","q"),
     ("13. Names of students from Kathmandu → π_name(σ_city='Kathmandu'(Student))","o"),
     ("14. Each student's name with course and grade → π_name,course,grade(Student ⋈ Enrollment)","o"),
     ("15. Average grade per course → course ℱ AVG(grade)(Enrollment)","o"),
     ("16. TRC: names of BIM students → { t.name | t ∈ Student ∧ t.program='BIM' }","o"),
     ("Part C — Discussion","q"),
     ("17. Decompose one screen of a Nepali app into σ / π / JOIN / aggregate operations.","o"),
     ("18. Why do algebra (HOW) and calculus (WHAT) always give the same answer, and why did SQL need both ideas?","o")],
    "Consolidated quiz Parts B–C. Model answers on the concept slides and in Unit3_material.md.",compact=True)

# ---- close ----
add_title("End of Unit 3  ·  IT 220",
          "S13–S17 complete: relational algebra (σ, π, ∪, ∩, −, ×, ⋈, ÷, ℱ, outer joins) and relational calculus (TRC, DRC) — the procedural and declarative foundations of SQL",
          "Built to COURSE_MATERIAL_STANDARD.md · self-contained slides · exports cleanly to PDF · "
          "Next unit (Unit 4): Database Normalization.")

_add_page_numbers()
out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"IT220_Unit3.pptx")
prs.save(out)
print("Saved",out,"with",len(prs.slides._sldIdLst),"slides")
