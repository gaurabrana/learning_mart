#!/usr/bin/env python3
"""deckkit — shared slide-builder toolkit for the IT 246 (sixth-sem) decks.

Single source of the visual system + slide helpers so every `build_unitN_pptx.py`
imports ONE copy instead of pasting ~150 lines each (see PROJECT_STATUS.md gap #3).
Built to COURSE_MATERIAL_STANDARD.md — palette/fonts/slide-types in §5, the §7A depth
tables via `add_table_slide`. Import with `from deckkit import *`, build, then call
`_add_page_numbers()` and `save("IT246_UnitN.pptx")`.

Palette (hex): NAVY 0C2B4A · BLUE 185FA5 · TEAL 0F6E56 · AMBER 854F0B · CORAL 993C1D;
card tints BLUE_T E6F1FB / CORAL_T FAECE7 / TEAL_T E1F5EE / AMBER_T FAEEDA. 13.333×7.5in.

NOTE: `prs` is a module-level singleton created at import — one deck per process, which
is how the build scripts run (one script = one process = one deck). Images resolve from
`images/` next to THIS file, so all sixth-sem decks share the same image folder.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "images")
NAVY=RGBColor(0x0C,0x2B,0x4A); BLUE=RGBColor(0x18,0x5F,0xA5); TEAL=RGBColor(0x0F,0x6E,0x56)
AMBER=RGBColor(0x85,0x4F,0x0B); CORAL=RGBColor(0x99,0x3C,0x1D); LIGHT=RGBColor(0xF2,0xF5,0xF8)
PALET=RGBColor(0xE1,0xF5,0xEE); WHITE=RGBColor(0xFF,0xFF,0xFF); DARK=RGBColor(0x1A,0x1A,0x1A)
GREY=RGBColor(0x55,0x60,0x6B); ROWALT=RGBColor(0xEE,0xF3,0xF8)
BLUE_T=RGBColor(0xE6,0xF1,0xFB); CORAL_T=RGBColor(0xFA,0xEC,0xE7)
TEAL_T=RGBColor(0xE1,0xF5,0xEE); AMBER_T=RGBColor(0xFA,0xEE,0xDA)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

def _bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def _box(s,l,t,w,h): return s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)).text_frame
def _bar(s,c,l,t,w,h):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=c; sp.line.fill.background(); sp.shadow.inherit=False
def _notes(s,t):
    if t: s.notes_slide.notes_text_frame.text=t
def _header(s,kicker,header,kc=BLUE):
    _bar(s,kc,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=kc
    th=_box(s,0.7,0.66,12,1.0); th.word_wrap=True; r=th.paragraphs[0].add_run(); r.text=header
    r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
def _card(s,l,t,w,h,accent,fill,label,body,body_sz=14):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    sp.line.color.rgb=accent; sp.line.width=Pt(1); sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP
    tf.margin_left=Inches(0.22); tf.margin_right=Inches(0.20); tf.margin_top=Inches(0.10); tf.margin_bottom=Inches(0.08)
    p=tf.paragraphs[0]; r=p.add_run(); r.text=label; r.font.size=Pt(12.5); r.font.bold=True; r.font.color.rgb=accent
    p2=tf.add_paragraph(); p2.space_before=Pt(3); r2=p2.add_run(); r2.text=body; r2.font.size=Pt(body_sz); r2.font.color.rgb=DARK
    return sp

def add_title(title,subtitle,footer):
    s=prs.slides.add_slide(BLANK); _bg(s,NAVY); _bar(s,AMBER,0.9,3.05,1.6,0.10)
    tf=_box(s,0.9,2.0,11.5,1.1); tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(33); r.font.bold=True; r.font.color.rgb=WHITE
    t2=_box(s,0.9,3.3,11.5,1.2); t2.word_wrap=True
    r=t2.paragraphs[0].add_run(); r.text=subtitle; r.font.size=Pt(18); r.font.color.rgb=RGBColor(0xCF,0xDD,0xEA)
    t3=_box(s,0.9,6.4,11.5,0.8); t3.word_wrap=True; r=t3.paragraphs[0].add_run(); r.text=footer
    r.font.size=Pt(13); r.font.color.rgb=RGBColor(0x9A,0xB2,0xC6); return s

def add_outcomes(header,kicker,intro,items,nxt):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,header)
    tf=_box(s,0.7,1.75,12.0,4.6); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=intro; r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=NAVY; p.space_after=Pt(8)
    for it in items:
        p=tf.add_paragraph(); r=p.add_run(); r.text="•  "+it; r.font.size=Pt(14.5); r.font.color.rgb=DARK; p.space_after=Pt(5)
    _card(s,0.7,6.35,12.0,0.85,BLUE,BLUE_T,"➡️ Course context",nxt,body_sz=13)
    return s

def add_roadmap(kicker,title,done,todo):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    _card(s,0.7,1.8,5.85,4.55,TEAL,TEAL_T,"✅  THIS UNIT  ·  S1–S5","\n".join(done),body_sz=14)
    _card(s,6.78,1.8,5.85,4.55,GREY,LIGHT,"⏭️  LATER UNITS","\n".join(todo),body_sz=14)
    return s

def add_divider(kicker,title,hook,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,BLUE)
    tk=_box(s,0.9,1.9,11.5,0.6); r=tk.paragraphs[0].add_run(); r.text=kicker.upper()
    r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=RGBColor(0xBF,0xDC,0xF2)
    tt=_box(s,0.9,2.5,11.5,1.4); tt.word_wrap=True; r=tt.paragraphs[0].add_run(); r.text=title
    r.font.size=Pt(30); r.font.bold=True; r.font.color.rgb=WHITE
    _bar(s,AMBER,0.9,4.05,1.4,0.09)
    th=_box(s,0.9,4.25,11.5,2.4); th.word_wrap=True; p=th.paragraphs[0]; r=p.add_run(); r.text="🎣  "+hook
    r.font.size=Pt(18); r.font.color.rgb=RGBColor(0xE8,0xF1,0xF9); r.font.italic=True
    _notes(s,timing); return s

def concept_understand(kicker,heading,definition,mechanism,image,hook,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker+"  ·  understand",heading)
    tw=6.2 if image else 11.9
    tf=_box(s,0.7,1.72,tw,5.05); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=definition; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=NAVY; p.space_after=Pt(10)
    p=tf.add_paragraph(); r=p.add_run(); r.text="HOW IT WORKS"; r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=BLUE; p.space_after=Pt(5)
    for m in mechanism:
        p=tf.add_paragraph(); r=p.add_run(); r.text="•  "+m; r.font.size=Pt(14); r.font.color.rgb=DARK; p.space_after=Pt(6)
    if image:
        path=os.path.join(IMG,image)
        if os.path.exists(path):
            pic=s.shapes.add_picture(path,Inches(6.95),Inches(1.9),width=Inches(6.05))
            band_top,band_h=1.9,4.55; h_in=pic.height/914400.0
            pic.top=Inches(band_top+max(0.0,(band_h-h_in)/2.0))
    _bar(s,TEAL,0,6.9,SW.inches,0.6)
    ft=_box(s,0.7,6.96,11.9,0.5); p=ft.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text="🧠  Memory hook:  "+hook; r.font.size=Pt(13.5); r.font.bold=True; r.font.color.rgb=WHITE
    _notes(s,timing); return s

def concept_apply(kicker,heading,example,trap,exam,keyterms):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker+"  ·  apply",heading)
    _card(s,0.7,1.66,12.0,1.48,BLUE ,BLUE_T ,"🌍  Real example (Nepal / IT)",example)
    _card(s,0.7,3.26,12.0,1.30,CORAL,CORAL_T,"⚠️  Common trap — and the fix",trap)
    _card(s,0.7,4.68,12.0,1.92,TEAL ,TEAL_T ,"🎯  Exam-ready answer",exam)
    _card(s,0.7,6.72,12.0,0.60,AMBER,AMBER_T,"🔑  Key terms",keyterms,body_sz=11.5)
    return s

def add_table_slide(kicker,title,headers,rows,per_page=8,widths=None,note="",fs=11):
    """Native PPTX table on its OWN slide; auto-paginates so nothing is squeezed."""
    pages=[rows[i:i+per_page] for i in range(0,len(rows),per_page)] or [[]]
    ncols=len(headers)
    if widths is None: widths=[1]*ncols
    tot=sum(widths); ws=[12.5*w/tot for w in widths]
    for pi,chunk in enumerate(pages):
        s=prs.slides.add_slide(BLANK); _bg(s,WHITE)
        _header(s,kicker,title+("  (cont.)" if pi else ""))
        nrows=len(chunk)+1
        top=1.72; avail=(6.55 if note else 6.95)-top
        gt=s.shapes.add_table(nrows,ncols,Inches(0.42),Inches(top),Inches(12.5),Inches(avail))
        table=gt.table; table.first_row=False; table.horz_banding=False
        for j,w in enumerate(ws): table.columns[j].width=Inches(w)
        for j,h in enumerate(headers):
            c=table.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=NAVY
            c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Inches(0.08); c.margin_right=Inches(0.06); c.margin_top=Inches(0.03); c.margin_bottom=Inches(0.03)
            tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; r=p.add_run(); r.text=h
            r.font.size=Pt(fs+0.5); r.font.bold=True; r.font.color.rgb=WHITE
        for i,row in enumerate(chunk,1):
            for j,val in enumerate(row):
                c=table.cell(i,j); c.fill.solid(); c.fill.fore_color.rgb=(WHITE if i%2 else ROWALT)
                c.vertical_anchor=MSO_ANCHOR.MIDDLE
                c.margin_left=Inches(0.08); c.margin_right=Inches(0.06); c.margin_top=Inches(0.03); c.margin_bottom=Inches(0.03)
                tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; r=p.add_run(); r.text=str(val)
                r.font.size=Pt(fs); r.font.color.rgb=DARK
                if j==0: r.font.bold=True; r.font.color.rgb=NAVY
        if note:
            nb=_box(s,0.42,6.62,12.5,0.5); p=nb.paragraphs[0]; r=p.add_run(); r.text="ℹ️  "+note
            r.font.size=Pt(11.5); r.font.italic=True; r.font.color.rgb=GREY
    return

# Standard §14 aliases — semantic names for the two mandated §7A table kinds.
# Both delegate to add_table_slide (one generic engine); use whichever reads clearest
# at the call site. See COURSE_MATERIAL_STANDARD.md §14.
def add_comparison_table(kicker,title,headers,rows,per_page=8,widths=None,note="",fs=11):
    """A confusable-set comparison table (X vs Y vs Z). Alias of add_table_slide."""
    return add_table_slide(kicker,title,headers,rows,per_page,widths,note,fs)
def add_examples_table(kicker,title,headers,rows,per_page=8,widths=None,note="",fs=11):
    """A concrete-example table (X vs not-X, ≥6 localised rows). Alias of add_table_slide."""
    return add_table_slide(kicker,title,headers,rows,per_page,widths,note,fs)

def add_activity(title,steps,expected,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,PALET); _bar(s,TEAL,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text="🛠  IN-CLASS ACTIVITY  ·  THINK–PAIR–SHARE"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=TEAL
    th=_box(s,0.7,0.66,12,0.9); r=th.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
    tf=_box(s,0.7,1.75,11.9,3.6); tf.word_wrap=True; first=True
    for text in steps:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text="•  "+text; r.font.size=Pt(15.5); r.font.color.rgb=DARK; p.space_after=Pt(7)
    _card(s,0.7,5.5,11.9,1.7,TEAL,WHITE,"✅  Expected answer / what to draw out",expected)
    _notes(s,timing); return s

def add_quiz(title,items,timing="",compact=False):
    s=prs.slides.add_slide(BLANK); _bg(s,LIGHT); _bar(s,AMBER,0,0,SW.inches,0.18)
    tk=_box(s,0.7,0.30,12,0.42); r=tk.paragraphs[0].add_run(); r.text="CHECK FOR UNDERSTANDING"
    r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=AMBER
    th=_box(s,0.7,0.66,12,0.9); r=th.paragraphs[0].add_run(); r.text=title; r.font.size=Pt(25); r.font.bold=True; r.font.color.rgb=NAVY
    qz,az,sb=(14,12.5,5) if compact else (17,15,9)
    tf=_box(s,0.7,1.75,11.9,5.3); tf.word_wrap=True; first=True
    for text,kind in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=text
        if kind=="q": r.font.size=Pt(qz); r.font.bold=True; r.font.color.rgb=NAVY; p.space_before=Pt(sb)
        elif kind=="a": r.font.size=Pt(az); r.font.color.rgb=TEAL; r.font.bold=True; p.level=1
        else: r.font.size=Pt(az); r.font.color.rgb=DARK; r.font.italic=True; p.level=1; p.space_before=Pt(sb-3)
    _notes(s,timing); return s

def add_summary(kicker,takeaways,why,nxt,timing=""):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,"Summary & Key Takeaways",kc=NAVY)
    tf=_box(s,0.7,1.7,12.0,2.7); tf.word_wrap=True; first=True
    for i,tk in enumerate(takeaways,1):
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        r=p.add_run(); r.text=f"{i}.  "+tk; r.font.size=Pt(15); r.font.color.rgb=DARK; p.space_after=Pt(8)
    _card(s,0.7,4.55,12.0,1.35,TEAL,TEAL_T,"💡  Why this matters",why,body_sz=13.5)
    _card(s,0.7,6.02,12.0,1.15,BLUE,BLUE_T,"➡️  Next",nxt,body_sz=13.5)
    _notes(s,timing); return s

def add_cheatsheet(kicker,title,blocks):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    xs=[0.7,6.78]; ys=[1.72,3.36,5.0]; hh=1.55
    for idx,(label,body) in enumerate(blocks[:6]):
        _card(s,xs[idx%2],ys[idx//2],5.85,hh,BLUE,BLUE_T,label,body,body_sz=12)
    return s

def add_glossary(kicker,title,terms):
    s=prs.slides.add_slide(BLANK); _bg(s,WHITE); _header(s,kicker,title)
    half=(len(terms)+1)//2
    for x,group in [(0.7,terms[:half]),(6.78,terms[half:])]:
        tf=_box(s,x,1.75,5.85,5.3); tf.word_wrap=True; first=True
        for term,defn in group:
            p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            r=p.add_run(); r.text=term+" — "; r.font.size=Pt(11.5); r.font.bold=True; r.font.color.rgb=NAVY
            r2=p.add_run(); r2.text=defn; r2.font.size=Pt(11.5); r2.font.color.rgb=DARK; p.space_after=Pt(6)
    return s

def _add_page_numbers():
    total=len(prs.slides._sldIdLst)
    for i,s in enumerate(prs.slides,1):
        try:
            h=str(s.background.fill.fore_color.rgb); rr,gg,bb=int(h[0:2],16),int(h[2:4],16),int(h[4:6],16)
            dark=(0.299*rr+0.587*gg+0.114*bb)<128
        except Exception: dark=False
        col=RGBColor(0xE8,0xF0,0xF7) if dark else GREY
        tb=_box(s,12.35,0.30,0.9,0.35); p=tb.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
        r=p.add_run(); r.text=f"{i} / {total}"; r.font.size=Pt(10); r.font.color.rgb=col

def save(filename):
    """Save the deck next to the calling build script's folder (sixth/)."""
    out=os.path.join(os.path.dirname(os.path.abspath(__file__)),filename)
    prs.save(out); print("Saved",out,"with",len(prs.slides._sldIdLst),"slides"); return out

# Everything public + the underscore primitives (so unit scripts can build custom
# layouts) + the shared globals. `from deckkit import *` pulls all of these.
__all__=["prs","BLANK","SW","SH","IMG",
    "NAVY","BLUE","TEAL","AMBER","CORAL","LIGHT","PALET","WHITE","DARK","GREY","ROWALT",
    "BLUE_T","CORAL_T","TEAL_T","AMBER_T","Inches","Pt","RGBColor","PP_ALIGN","MSO_ANCHOR","MSO_SHAPE",
    "_bg","_box","_bar","_notes","_header","_card",
    "add_title","add_outcomes","add_roadmap","add_divider","concept_understand","concept_apply",
    "add_table_slide","add_comparison_table","add_examples_table","add_activity","add_quiz",
    "add_summary","add_cheatsheet","add_glossary","_add_page_numbers","save"]
